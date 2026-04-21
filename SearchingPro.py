import os
import re
import sys
import time
import json
import sqlite3
import datetime
import shutil
from concurrent.futures import ThreadPoolExecutor

try:
    import winreg
except ImportError:
    winreg = None
try:
    import psutil
except ImportError:
    psutil = None
try:
    import ctypes
    from ctypes import wintypes
except ImportError:
    ctypes = None
    wintypes = None
try:
    import matplotlib
    matplotlib.use('Qt5Agg')
    from matplotlib.backends.backend_qt5agg import FigureCanvasQTAgg as FigureCanvas
    from matplotlib.figure import Figure
    import matplotlib.pyplot as plt
    import matplotlib.font_manager as fm
    _cjk_font_found = False
    for font_name in ['PingFang TC', 'PingFang SC', 'Heiti TC', 'Microsoft JhengHei', 'Microsoft YaHei', 'Noto Sans CJK TC', 'Arial Unicode MS']:
        if font_name in [f.name for f in fm.fontManager.ttflist]:
            plt.rcParams['font.sans-serif'] = [font_name] + plt.rcParams['font.sans-serif']
            _cjk_font_found = True
            break
    if not _cjk_font_found:
        plt.rcParams['font.sans-serif'] = ['Arial', 'Helvetica'] + plt.rcParams['font.sans-serif']
    plt.rcParams['axes.unicode_minus'] = False
except ImportError:
    matplotlib = None
    FigureCanvas = None
    Figure = None
    plt = None

from PyQt5.QtWidgets import (
    QApplication,
    QMainWindow,
    QWidget,
    QVBoxLayout,
    QHBoxLayout,
    QLabel,
    QLineEdit,
    QPushButton,
    QComboBox,
    QTableWidget,
    QTableWidgetItem,
    QFileDialog,
    QCheckBox,
    QGroupBox,
    QStatusBar,
    QTabWidget,
    QTextEdit,
    QSpinBox,
    QHeaderView,
    QMessageBox,
    QProgressBar,
    QDateEdit,
    QGridLayout,
    QMenu,
    QAction,
    QSplitter,
    QInputDialog,
    QListWidget,
    QListWidgetItem,
    QSystemTrayIcon,
    QStyle,
    QTreeWidget,
    QTreeWidgetItem,
    QDialog,
    QDialogButtonBox,
    QRadioButton,
    QFrame,
    QToolButton,
    QSizePolicy,
    QStackedWidget,
    QTimeEdit,
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, QDate, QEvent, QTimer, QPropertyAnimation, QEasingCurve, QParallelAnimationGroup, QRect
from PyQt5.QtGui import QIcon, QColor, QBrush, QFont, QTextDocument, QTextCharFormat, QTextCursor
from PyQt5.QtWidgets import QFrame, QToolButton, QSizePolicy, QStyledItemDelegate


SETTINGS_FILENAME = "settings.json"
APP_NAME = "AdvancedFileSearcher"


def _platform_shortcut(shortcut):
    """將跨平台快捷捷徑字串轉換為當前平台格式"""
    import platform
    system = platform.system()
    if system == "Windows":
        return shortcut.replace("Cmd+", "Ctrl+")
    return shortcut  # macOS and others keep Cmd+


def _shortcut_label(shortcut):
    """回傳適合顯示給用戶的快捷捷徑標籤"""
    import platform
    system = platform.system()
    if system == "Windows":
        return shortcut.replace("Cmd+", "Ctrl+")
    return shortcut


def _monospace_font():
    """回傳跨平台的等寬字型"""
    import platform
    system = platform.system()
    if system == "Windows":
        return QFont("Consolas", 11)
    elif system == "Darwin":
        return QFont("Menlo", 12)
    else:
        return QFont("DejaVu Sans Mono", 11)


class ContentExtractor:
    """多格式文件內容提取器"""

    @staticmethod
    def extract_text(file_path, ext):
        """根據檔案類型提取文字內容"""
        ext = ext.lower()
        try:
            if ext == ".pdf":
                return ContentExtractor._extract_pdf(file_path)
            elif ext == ".docx":
                return ContentExtractor._extract_docx(file_path)
            elif ext == ".xlsx":
                return ContentExtractor._extract_xlsx(file_path)
            elif ext == ".pptx":
                return ContentExtractor._extract_pptx(file_path)
            elif ext in {".txt", ".py", ".java", ".c", ".cpp", ".html", ".xml", ".json", ".csv", ".md", ".log", ".ini", ".conf", ".js", ".css", ".yml", ".yaml"}:
                return ContentExtractor._extract_plain(file_path)
            else:
                return ""
        except Exception as e:
            print(f"提取 {ext} 內容失敗: {e}")
            return ""

    @staticmethod
    def _extract_pdf(file_path):
        from PyPDF2 import PdfReader
        text = []
        reader = PdfReader(file_path)
        for page in reader.pages[:20]:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        return "\n".join(text)[:1000000]

    @staticmethod
    def _extract_docx(file_path):
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)[:1000000]

    @staticmethod
    def _extract_xlsx(file_path):
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text = []
        for sheet in wb.worksheets[:5]:
            for row in sheet.iter_rows(max_row=1000):
                for cell in row:
                    if cell.value:
                        text.append(str(cell.value))
        wb.close()
        return "\n".join(text)[:1000000]

    @staticmethod
    def _extract_pptx(file_path):
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides[:20]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)[:1000000]

    @staticmethod
    def _extract_plain(file_path):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(1000000)

    @staticmethod
    def get_supported_extensions():
        return {
            ".pdf", ".docx", ".xlsx", ".pptx",
            ".txt", ".py", ".java", ".c", ".cpp", ".html", ".xml", ".json", ".csv", ".md",
            ".log", ".ini", ".conf", ".js", ".css", ".yml", ".yaml"
        }


class HighlightDelegate(QStyledItemDelegate):
    """表格關鍵字高亮代理"""

    def __init__(self, keywords=None, parent=None):
        super().__init__(parent)
        self.keywords = keywords or []
        self.highlight_color = QColor("#007aff")
        self.highlight_bg = QColor("#e5f2ff")

    def set_keywords(self, keywords):
        self.keywords = keywords if isinstance(keywords, list) else [keywords] if keywords else []

    def paint(self, painter, option, index):
        text = index.data()
        if not text or not self.keywords:
            super().paint(painter, option, index)
            return

        super().paint(painter, option, index)

        text_lower = text.lower()
        painter.save()
        painter.setRenderHint(painter.TextAntialiasing)

        font_metrics = painter.fontMetrics()
        rect = option.rect

        for keyword in self.keywords:
            if not keyword:
                continue
            kw_lower = keyword.lower()
            start = 0
            while True:
                pos = text_lower.find(kw_lower, start)
                if pos == -1:
                    break
                highlight_rect = QRect(rect.x() + font_metrics.horizontalAdvance(text[:pos]),
                                       rect.y(),
                                       font_metrics.horizontalAdvance(text[pos:pos + len(keyword)]),
                                       rect.height())
                painter.fillRect(highlight_rect, self.highlight_bg)
                painter.setPen(self.highlight_color)
                painter.drawText(highlight_rect, Qt.AlignLeft | Qt.AlignVCenter, text[pos:pos + len(keyword)])
                start = pos + len(keyword)

        painter.restore()


class SearchTemplateManager:
    """搜尋模板管理器"""

    def __init__(self, db_path):
        self.db_path = db_path
        self._create_table()

    def _create_table(self):
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute("""
                CREATE TABLE IF NOT EXISTS search_templates (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    name TEXT UNIQUE,
                    filename TEXT,
                    content TEXT,
                    file_type TEXT,
                    boolean_expr TEXT,
                    path_filter TEXT,
                    use_regex_filename INTEGER,
                    use_regex_content INTEGER,
                    use_fuzzy_filename INTEGER,
                    use_fuzzy_content INTEGER,
                    use_boolean INTEGER,
                    created_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    last_used TIMESTAMP,
                    use_count INTEGER DEFAULT 0
                )
            """)
            conn.commit()
            conn.close()
        except Exception as e:
            print(f"建立模板表失敗: {e}")

    def save_template(self, name, config):
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute("""
                INSERT OR REPLACE INTO search_templates 
                (name, filename, content, file_type, boolean_expr, path_filter,
                 use_regex_filename, use_regex_content, use_fuzzy_filename, use_fuzzy_content, use_boolean)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, (
                name,
                config.get("filename", ""),
                config.get("content", ""),
                config.get("file_type", ""),
                config.get("boolean_expr", ""),
                config.get("path_filter", ""),
                int(config.get("use_regex_filename", False)),
                int(config.get("use_regex_content", False)),
                int(config.get("use_fuzzy_filename", False)),
                int(config.get("use_fuzzy_content", False)),
                int(config.get("boolean_mode", False)),
            ))
            conn.commit()
            conn.close()
            return True
        except Exception as e:
            print(f"儲存模板失敗: {e}")
            return False

    def load_template(self, name):
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute("""
                UPDATE search_templates SET last_used = CURRENT_TIMESTAMP, use_count = use_count + 1
                WHERE name = ?
            """, (name,))
            cursor.execute("SELECT * FROM search_templates WHERE name = ?", (name,))
            record = cursor.fetchone()
            conn.commit()
            conn.close()
            if record:
                return {
                    "filename": record[2],
                    "content": record[3],
                    "file_type": record[4],
                    "boolean_expr": record[5],
                    "path_filter": record[6],
                    "use_regex_filename": bool(record[7]),
                    "use_regex_content": bool(record[8]),
                    "use_fuzzy_filename": bool(record[9]),
                    "use_fuzzy_content": bool(record[10]),
                    "boolean_mode": bool(record[11]),
                }
            return None
        except Exception as e:
            print(f"載入模板失敗: {e}")
            return None

    def get_all_templates(self):
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute("""
                SELECT name, last_used, use_count FROM search_templates
                ORDER BY last_used DESC
            """)
            records = cursor.fetchall()
            conn.close()
            return records
        except Exception as e:
            print(f"取得模板列表失敗: {e}")
            return []

    def delete_template(self, name):
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute("DELETE FROM search_templates WHERE name = ?", (name,))
            conn.commit()
            conn.close()
            return True
        except Exception as e:
            print(f"刪除模板失敗: {e}")
            return False

    def get_template_suggestions(self, prefix, limit=10):
        """根據前綴取得模板名稱建議"""
        suggestions = []
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            prefix_pattern = f"%{prefix}%"
            cursor.execute("""
                SELECT name, filename, content, boolean_expr FROM search_templates
                WHERE (name LIKE ? OR filename LIKE ? OR content LIKE ? OR boolean_expr LIKE ?)
                ORDER BY use_count DESC, last_used DESC LIMIT ?
            """, (prefix_pattern, prefix_pattern, prefix_pattern, prefix_pattern, limit))
            for row in cursor.fetchall():
                name = row[0]
                if name and name not in suggestions:
                    suggestions.append(name)
            conn.close()
        except Exception as e:
            print(f"取得模板建議失敗: {e}")
        return suggestions[:limit]


from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler


class FileIndexWatcher(FileSystemEventHandler):
    """檔案系統監控事件處理器"""

    def __init__(self, callback=None):
        super().__init__()
        self.callback = callback
        self.pending_changes = set()
        self.process_timer = QTimer()
        self.process_timer.setSingleShot(True)
        self.process_timer.timeout.connect(self._process_pending_changes)

    def on_created(self, event):
        if not event.is_directory:
            self.pending_changes.add(("create", event.src_path))
            self.process_timer.start(2000)

    def on_modified(self, event):
        if not event.is_directory:
            self.pending_changes.add(("modify", event.src_path))
            self.process_timer.start(2000)

    def on_deleted(self, event):
        if not event.is_directory:
            self.pending_changes.add(("delete", event.src_path))
            self.process_timer.start(2000)

    def on_moved(self, event):
        if not event.is_directory:
            self.pending_changes.add(("move", event.src_path, event.dest_path))
            self.process_timer.start(2000)

    def _process_pending_changes(self):
        if self.callback and self.pending_changes:
            self.callback(list(self.pending_changes))
            self.pending_changes.clear()


class IndexWatchManager:
    """索引即時監控管理器"""

    def __init__(self, db_path, content_types, exclude_rules):
        self.db_path = db_path
        self.content_types = content_types
        self.exclude_rules = exclude_rules
        self.observer = None
        self.handler = None
        self.watched_dirs = set()
        self.is_running = False

    def start_watching(self, directories):
        if self.is_running:
            self.stop_watching()

        self.handler = FileIndexWatcher(callback=self._handle_file_changes)
        self.observer = Observer()

        for dir_path in directories:
            if os.path.isdir(dir_path):
                self.observer.schedule(self.handler, dir_path, recursive=True)
                self.watched_dirs.add(dir_path)

        self.observer.start()
        self.is_running = True
        print(f"開始監控 {len(self.watched_dirs)} 個目錄")

    def stop_watching(self):
        if self.observer:
            self.observer.stop()
            self.observer.join()
            self.observer = None
        self.handler = None
        self.watched_dirs.clear()
        self.is_running = False
        print("停止檔案監控")

    def _handle_file_changes(self, changes):
        conn = None
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            for change in changes:
                action = change[0]

                if action == "create":
                    path = change[1]
                    self._add_file(cursor, path)

                elif action == "modify":
                    path = change[1]
                    self._update_file(cursor, path)

                elif action == "delete":
                    path = change[1]
                    cursor.execute("DELETE FROM files WHERE path = ?", (path,))

                elif action == "move":
                    old_path, new_path = change[1], change[2]
                    cursor.execute("DELETE FROM files WHERE path = ?", (old_path,))
                    self._add_file(cursor, new_path)

            conn.commit()
        except Exception as e:
            print(f"處理檔案變更時出錯: {e}")
        finally:
            if conn:
                conn.close()

    def _add_file(self, cursor, path):
        if not os.path.exists(path):
            return
        try:
            filename = os.path.basename(path)
            ext = os.path.splitext(filename)[1].lower()
            stats = os.stat(path)
            modified = datetime.datetime.fromtimestamp(stats.st_mtime).strftime("%Y-%m-%d %H:%M:%S")

            content = ""
            if ext in ContentExtractor.get_supported_extensions() or ext in self.content_types:
                content = ContentExtractor.extract_text(path, ext)

            cursor.execute(
                "INSERT OR REPLACE INTO files (path, filename, ext, size, modified, content) VALUES (?, ?, ?, ?, ?, ?)",
                (path, filename, ext, stats.st_size, modified, content)
            )
        except Exception as e:
            print(f"新增檔案 {path} 時出錯: {e}")

    def _update_file(self, cursor, path):
        self._add_file(cursor, path)


class CollapsibleGroupBox(QFrame):
    """可摺疊區域元件，帶動畫效果"""

    def __init__(self, title="", parent=None):
        super().__init__(parent)
        self._title = title
        self._collapsed = False
        self._content_widget = None
        self._main_layout = None
        self._header_layout = None
        self._toggle_button = None
        self._animation = None
        self._max_height = 0
        self._init_ui()

    def _init_ui(self):
        self.setFrameShape(QFrame.StyledPanel)
        self.setFrameShadow(QFrame.Raised)

        self._main_layout = QVBoxLayout(self)
        self._main_layout.setContentsMargins(0, 0, 0, 0)
        self._main_layout.setSpacing(0)

        self._header_layout = QHBoxLayout()
        self._header_layout.setContentsMargins(8, 4, 8, 4)

        self._toggle_button = QToolButton()
        self._toggle_button.setCheckable(True)
        self._toggle_button.setChecked(True)
        self._toggle_button.setArrowType(Qt.DownArrow)
        self._toggle_button.setFixedSize(20, 20)
        self._toggle_button.setStyleSheet("""
            QToolButton {
                border: none;
                background: transparent;
            }
            QToolButton:hover {
                background: rgba(0, 122, 255, 0.1);
                border-radius: 4px;
            }
        """)
        self._toggle_button.clicked.connect(self._on_toggle)

        title_label = QLabel(self._title)
        title_label.setStyleSheet("font-weight: 600; font-size: 14px;")

        self._header_layout.addWidget(self._toggle_button)
        self._header_layout.addWidget(title_label)
        self._header_layout.addStretch()

        self._main_layout.addLayout(self._header_layout)

        self._content_widget = QWidget()
        self._main_layout.addWidget(self._content_widget)

    def _on_toggle(self):
        self._collapsed = not self._collapsed
        self._toggle_button.setArrowType(Qt.UpArrow if self._collapsed else Qt.DownArrow)

        if self._collapsed:
            self._max_height = self._content_widget.sizeHint().height()
            self._animation = QPropertyAnimation(self._content_widget, b"maximumHeight")
            self._animation.setDuration(200)
            self._animation.setStartValue(self._max_height)
            self._animation.setEndValue(0)
            self._animation.setEasingCurve(QEasingCurve.InOutCubic)
            self._animation.start()
            self._content_widget.setMaximumHeight(0)
        else:
            self._content_widget.setMaximumHeight(16777215)

    def contentLayout(self):
        if self._content_widget.layout() is None:
            self._content_widget.setLayout(QVBoxLayout())
            self._content_widget.layout().setContentsMargins(8, 4, 8, 8)
            self._content_widget.layout().setSpacing(8)
        return self._content_widget.layout()

    def setContentLayout(self, layout):
        old_layout = self._content_widget.layout()
        if old_layout is not None:
            while old_layout.count():
                item = old_layout.takeAt(0)
                if item.widget():
                    item.widget().setParent(None)
            del old_layout
        self._content_widget.setLayout(layout)

    def isCollapsed(self):
        return self._collapsed

    def setCollapsed(self, collapsed):
        if self._collapsed != collapsed:
            self._on_toggle()


LIGHT_STYLE = """
    QMainWindow {
        background-color: #f5f5f7;
        color: #1d1d1f;
    }

    QTabWidget::pane {
        border: 1px solid #d1d1d6;
        border-top: 2px solid #007aff;
        background-color: #f5f5f7;
        border-radius: 0 0 8px 8px;
        margin-top: -1px;
    }

    QTabWidget::tab-bar {
        alignment: center;
    }

    QTabBar::tab {
        background-color: #e5e5ea;
        color: #1d1d1f;
        border: 1px solid #d1d1d6;
        border-bottom: none;
        padding: 10px 28px;
        margin-right: 0;
        border-top-left-radius: 8px;
        border-top-right-radius: 8px;
        min-width: 100px;
        font-size: 13px;
    }

    QTabBar::tab:selected {
        background-color: #f5f5f7;
        border-bottom: 2px solid #f5f5f7;
        color: #007aff;
        font-weight: 600;
    }

    QTabBar::tab:hover:!selected {
        background-color: #d1d1d6;
    }

    QGroupBox {
        font-weight: 600;
        color: #1d1d1f;
        border: 1px solid #d1d1d6;
        border-radius: 8px;
        margin-top: 12px;
        padding-top: 16px;
        background-color: #ffffff;
    }

    QGroupBox::title {
        subcontrol-origin: margin;
        left: 12px;
        padding: 0 8px 0 8px;
        background-color: #ffffff;
    }

    QPushButton {
        background-color: #007aff;
        color: white;
        border: none;
        border-radius: 6px;
        padding: 8px 16px;
        font-weight: 500;
        min-width: 80px;
        font-size: 13px;
    }

    QPushButton:hover {
        background-color: #0056cc;
    }

    QPushButton:pressed {
        background-color: #004499;
    }

    QPushButton:disabled {
        background-color: #d1d1d6;
        color: #8e8e93;
    }

    QPushButton#deleteButton {
        background-color: #ff3b30;
    }

    QPushButton#deleteButton:hover {
        background-color: #d70015;
    }

    QPushButton#warningButton {
        background-color: #ff9500;
    }

    QPushButton#warningButton:hover {
        background-color: #e68600;
    }

    QPushButton#successButton {
        background-color: #34c759;
    }

    QPushButton#successButton:hover {
        background-color: #28a745;
    }

    QPushButton#secondaryButton {
        background-color: #e5e5ea;
        color: #1d1d1f;
    }

    QPushButton#secondaryButton:hover {
        background-color: #d1d1d6;
    }

    QLineEdit {
        border: 1px solid #d1d1d6;
        border-radius: 6px;
        padding: 8px 12px;
        background-color: #ffffff;
        color: #1d1d1f;
        font-size: 13px;
    }

    QLineEdit:focus {
        border: 2px solid #007aff;
        outline: none;
    }

    QComboBox {
        border: 1px solid #d1d1d6;
        border-radius: 6px;
        padding: 8px 12px;
        background-color: #ffffff;
        color: #1d1d1f;
        min-width: 120px;
        font-size: 13px;
    }

    QComboBox:focus {
        border: 2px solid #007aff;
    }

    QComboBox::drop-down {
        border: none;
        width: 20px;
    }

    QComboBox::down-arrow {
        image: none;
        border-left: 4px solid transparent;
        border-right: 4px solid transparent;
        border-top: 6px solid #8e8e93;
        margin-right: 8px;
    }

    QSpinBox {
        border: 1px solid #d1d1d6;
        border-radius: 6px;
        padding: 8px 12px;
        background-color: #ffffff;
        color: #1d1d1f;
        font-size: 13px;
    }

    QSpinBox:focus {
        border: 2px solid #007aff;
    }

    QDateEdit {
        border: 1px solid #d1d1d6;
        border-radius: 6px;
        padding: 8px 12px;
        background-color: #ffffff;
        color: #1d1d1f;
        font-size: 13px;
    }

    QDateEdit:focus {
        border: 2px solid #007aff;
    }

    QCheckBox {
        color: #1d1d1f;
        spacing: 8px;
        font-size: 13px;
    }

    QCheckBox::indicator {
        width: 18px;
        height: 18px;
        border: 2px solid #d1d1d6;
        border-radius: 4px;
        background-color: #ffffff;
    }

    QCheckBox::indicator:checked {
        background-color: #007aff;
        border-color: #007aff;
    }

    QCheckBox::indicator:hover {
        border-color: #007aff;
    }

    QTableWidget {
        border: 1px solid #d1d1d6;
        border-radius: 8px;
        background-color: #ffffff;
        alternate-background-color: #f9f9f9;
        gridline-color: #e5e5ea;
        selection-background-color: #007aff;
        selection-color: white;
        font-size: 13px;
    }

    QTableWidget::item {
        padding: 8px;
        border: none;
        color: #1d1d1f;
    }

    QTableWidget::item:selected {
        background-color: #007aff;
        color: white;
    }

    QHeaderView::section {
        background-color: #f5f5f7;
        color: #1d1d1f;
        padding: 10px;
        border: none;
        border-bottom: 2px solid #d1d1d6;
        font-weight: 600;
        font-size: 13px;
    }

    QHeaderView::section:vertical {
        background-color: #f5f5f7;
        color: #8e8e93;
        border: none;
        border-right: 1px solid #e5e5ea;
        font-weight: 500;
        font-size: 12px;
        padding: 4px 8px;
    }

    QHeaderView::section:horizontal {
        background-color: #f5f5f7;
        color: #1d1d1f;
        border: none;
        border-bottom: 2px solid #d1d1d6;
        font-weight: 600;
        font-size: 13px;
        padding: 10px;
    }

    QTextEdit {
        border: 1px solid #d1d1d6;
        border-radius: 8px;
        background-color: #ffffff;
        color: #1d1d1f;
        padding: 8px;
        font-size: 13px;
    }

    QTextEdit:focus {
        border: 2px solid #007aff;
    }

    QStatusBar {
        background-color: #f5f5f7;
        color: #1d1d1f;
        border-top: 1px solid #d1d1d6;
        font-size: 12px;
    }

    QProgressBar {
        border: 1px solid #d1d1d6;
        border-radius: 6px;
        background-color: #e5e5ea;
        text-align: center;
        color: #1d1d1f;
        font-size: 12px;
    }

    QProgressBar::chunk {
        background-color: #007aff;
        border-radius: 5px;
    }

    QSplitter::handle {
        background-color: #d1d1d6;
    }

    QSplitter::handle:horizontal {
        width: 3px;
    }

    QSplitter::handle:vertical {
        height: 3px;
    }

    QListWidget {
        border: 1px solid #d1d1d6;
        border-radius: 8px;
        background-color: #ffffff;
        alternate-background-color: #f9f9f9;
        selection-background-color: #007aff;
        selection-color: white;
        font-size: 13px;
    }

    QListWidget::item {
        padding: 8px;
        border: none;
        color: #1d1d1f;
    }

    QListWidget::item:selected {
        background-color: #007aff;
        color: white;
    }

    QListWidget::item:hover {
        background-color: #e5e5ea;
        color: #1d1d1f;
    }

    QTreeWidget {
        border: 1px solid #d1d1d6;
        border-radius: 8px;
        background-color: #ffffff;
        alternate-background-color: #f9f9f9;
        selection-background-color: #007aff;
        selection-color: white;
        font-size: 13px;
    }

    QTreeWidget::item {
        padding: 6px;
        color: #1d1d1f;
    }

    QTreeWidget::item:selected {
        background-color: #007aff;
        color: white;
    }

    QLabel {
        color: #1d1d1f;
        font-size: 13px;
    }

    QMenuBar {
        background-color: #ffffff;
        color: #1d1d1f;
        border-bottom: 1px solid #d1d1d6;
        font-size: 13px;
    }

    QMenuBar::item {
        padding: 8px 12px;
        background-color: transparent;
        border-radius: 4px;
    }

    QMenuBar::item:selected {
        background-color: #e5e5ea;
    }

    QMenu {
        background-color: #ffffff;
        border: 1px solid #d1d1d6;
        border-radius: 8px;
        padding: 4px;
    }

    QMenu::item {
        padding: 8px 24px 8px 12px;
        border-radius: 4px;
    }

    QMenu::item:selected {
        background-color: #007aff;
        color: white;
    }

    QMenu::separator {
        height: 1px;
        background-color: #e5e5ea;
        margin: 4px 8px;
    }

    QFrame {
        background-color: transparent;
    }

    QScrollBar:vertical {
        border: none;
        background: #f5f5f7;
        width: 10px;
        border-radius: 5px;
    }

    QScrollBar::handle:vertical {
        background: #c1c1c6;
        border-radius: 5px;
        min-height: 30px;
    }

    QScrollBar::handle:vertical:hover {
        background: #a1a1a6;
    }

    QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {
        height: 0px;
    }

    QScrollBar:horizontal {
        border: none;
        background: #f5f5f7;
        height: 10px;
        border-radius: 5px;
    }

    QScrollBar::handle:horizontal {
        background: #c1c1c6;
        border-radius: 5px;
        min-width: 30px;
    }

    QScrollBar::handle:horizontal:hover {
        background: #a1a1a6;
    }

    QScrollBar::add-line:horizontal, QScrollBar::sub-line:horizontal {
        width: 0px;
    }
"""

DARK_STYLE = """
    QMainWindow {
        background-color: #1c1c1e;
        color: #ffffff;
    }

    QTabWidget::pane {
        border: 1px solid #38383a;
        border-top: 2px solid #0a84ff;
        background-color: #1c1c1e;
        border-radius: 0 0 8px 8px;
        margin-top: -1px;
    }

    QTabWidget::tab-bar {
        alignment: center;
    }

    QTabBar::tab {
        background-color: #2c2c2e;
        color: #ffffff;
        border: 1px solid #38383a;
        border-bottom: none;
        padding: 10px 28px;
        margin-right: 0;
        border-top-left-radius: 8px;
        border-top-right-radius: 8px;
        min-width: 100px;
        font-size: 13px;
    }

    QTabBar::tab:selected {
        background-color: #1c1c1e;
        border-bottom: 2px solid #1c1c1e;
        color: #0a84ff;
        font-weight: 600;
    }

    QTabBar::tab:hover:!selected {
        background-color: #3a3a3c;
    }

    QGroupBox {
        font-weight: 600;
        color: #ffffff;
        border: 1px solid #38383a;
        border-radius: 8px;
        margin-top: 12px;
        padding-top: 16px;
        background-color: #2c2c2e;
    }

    QGroupBox::title {
        subcontrol-origin: margin;
        left: 12px;
        padding: 0 8px 0 8px;
        background-color: #2c2c2e;
    }

    QPushButton {
        background-color: #0a84ff;
        color: white;
        border: none;
        border-radius: 6px;
        padding: 8px 16px;
        font-weight: 500;
        min-width: 80px;
        font-size: 13px;
    }

    QPushButton:hover {
        background-color: #409cff;
    }

    QPushButton:pressed {
        background-color: #0066cc;
    }

    QPushButton:disabled {
        background-color: #38383a;
        color: #636366;
    }

    QPushButton#deleteButton {
        background-color: #ff453a;
    }

    QPushButton#deleteButton:hover {
        background-color: #ff6961;
    }

    QPushButton#warningButton {
        background-color: #ff9f0a;
    }

    QPushButton#warningButton:hover {
        background-color: #ffb340;
    }

    QPushButton#successButton {
        background-color: #30d158;
    }

    QPushButton#successButton:hover {
        background-color: #4ae16a;
    }

    QPushButton#secondaryButton {
        background-color: #3a3a3c;
        color: #ffffff;
    }

    QPushButton#secondaryButton:hover {
        background-color: #48484a;
    }

    QLineEdit {
        border: 1px solid #38383a;
        border-radius: 6px;
        padding: 8px 12px;
        background-color: #1c1c1e;
        color: #ffffff;
        font-size: 13px;
    }

    QLineEdit:focus {
        border: 2px solid #0a84ff;
        outline: none;
    }

    QComboBox {
        border: 1px solid #38383a;
        border-radius: 6px;
        padding: 8px 12px;
        background-color: #1c1c1e;
        color: #ffffff;
        min-width: 120px;
        font-size: 13px;
    }

    QComboBox:focus {
        border: 2px solid #0a84ff;
    }

    QComboBox::drop-down {
        border: none;
        width: 20px;
    }

    QComboBox::down-arrow {
        image: none;
        border-left: 4px solid transparent;
        border-right: 4px solid transparent;
        border-top: 6px solid #8e8e93;
        margin-right: 8px;
    }

    QSpinBox {
        border: 1px solid #38383a;
        border-radius: 6px;
        padding: 8px 12px;
        background-color: #1c1c1e;
        color: #ffffff;
        font-size: 13px;
    }

    QSpinBox:focus {
        border: 2px solid #0a84ff;
    }

    QDateEdit {
        border: 1px solid #38383a;
        border-radius: 6px;
        padding: 8px 12px;
        background-color: #1c1c1e;
        color: #ffffff;
        font-size: 13px;
    }

    QDateEdit:focus {
        border: 2px solid #0a84ff;
    }

    QCheckBox {
        color: #ffffff;
        spacing: 8px;
        font-size: 13px;
    }

    QCheckBox::indicator {
        width: 18px;
        height: 18px;
        border: 2px solid #38383a;
        border-radius: 4px;
        background-color: #1c1c1e;
    }

    QCheckBox::indicator:checked {
        background-color: #0a84ff;
        border-color: #0a84ff;
    }

    QCheckBox::indicator:hover {
        border-color: #0a84ff;
    }

    QTableWidget {
        border: 1px solid #38383a;
        border-radius: 8px;
        background-color: #2c2c2e;
        alternate-background-color: #323234;
        gridline-color: #38383a;
        selection-background-color: #0a84ff;
        selection-color: white;
        font-size: 13px;
    }

    QTableWidget::item {
        padding: 8px;
        border: none;
        color: #ffffff;
    }

    QTableWidget::item:selected {
        background-color: #0a84ff;
        color: white;
    }

    QHeaderView::section {
        background-color: #1c1c1e;
        color: #ffffff;
        padding: 10px;
        border: none;
        border-bottom: 2px solid #38383a;
        font-weight: 600;
        font-size: 13px;
    }

    QHeaderView::section:vertical {
        background-color: #1c1c1e;
        color: #636366;
        border: none;
        border-right: 1px solid #38383a;
        font-weight: 500;
        font-size: 12px;
        padding: 4px 8px;
    }

    QHeaderView::section:horizontal {
        background-color: #1c1c1e;
        color: #ffffff;
        border: none;
        border-bottom: 2px solid #38383a;
        font-weight: 600;
        font-size: 13px;
        padding: 10px;
    }

    QTextEdit {
        border: 1px solid #38383a;
        border-radius: 8px;
        background-color: #1c1c1e;
        color: #ffffff;
        padding: 8px;
        font-size: 13px;
    }

    QTextEdit:focus {
        border: 2px solid #0a84ff;
    }

    QStatusBar {
        background-color: #1c1c1e;
        color: #ffffff;
        border-top: 1px solid #38383a;
        font-size: 12px;
    }

    QProgressBar {
        border: 1px solid #38383a;
        border-radius: 6px;
        background-color: #38383a;
        text-align: center;
        color: #ffffff;
        font-size: 12px;
    }

    QProgressBar::chunk {
        background-color: #0a84ff;
        border-radius: 5px;
    }

    QSplitter::handle {
        background-color: #38383a;
    }

    QSplitter::handle:horizontal {
        width: 3px;
    }

    QSplitter::handle:vertical {
        height: 3px;
    }

    QListWidget {
        border: 1px solid #38383a;
        border-radius: 8px;
        background-color: #2c2c2e;
        alternate-background-color: #323234;
        selection-background-color: #0a84ff;
        selection-color: white;
        font-size: 13px;
    }

    QListWidget::item {
        padding: 8px;
        border: none;
        color: #ffffff;
    }

    QListWidget::item:selected {
        background-color: #0a84ff;
        color: white;
    }

    QListWidget::item:hover {
        background-color: #3a3a3c;
        color: #ffffff;
    }

    QTreeWidget {
        border: 1px solid #38383a;
        border-radius: 8px;
        background-color: #2c2c2e;
        alternate-background-color: #323234;
        selection-background-color: #0a84ff;
        selection-color: white;
        font-size: 13px;
    }

    QTreeWidget::item {
        padding: 6px;
        color: #ffffff;
    }

    QTreeWidget::item:selected {
        background-color: #0a84ff;
        color: white;
    }

    QLabel {
        color: #ffffff;
        font-size: 13px;
    }

    QMenuBar {
        background-color: #2c2c2e;
        color: #ffffff;
        border-bottom: 1px solid #38383a;
        font-size: 13px;
    }

    QMenuBar::item {
        padding: 8px 12px;
        background-color: transparent;
        border-radius: 4px;
    }

    QMenuBar::item:selected {
        background-color: #3a3a3c;
    }

    QMenu {
        background-color: #2c2c2e;
        border: 1px solid #38383a;
        border-radius: 8px;
        padding: 4px;
    }

    QMenu::item {
        padding: 8px 24px 8px 12px;
        border-radius: 4px;
    }

    QMenu::item:selected {
        background-color: #0a84ff;
        color: white;
    }

    QMenu::separator {
        height: 1px;
        background-color: #38383a;
        margin: 4px 8px;
    }

    QFrame {
        background-color: transparent;
    }

    QScrollBar:vertical {
        border: none;
        background: #1c1c1e;
        width: 10px;
        border-radius: 5px;
    }

    QScrollBar::handle:vertical {
        background: #48484a;
        border-radius: 5px;
        min-height: 30px;
    }

    QScrollBar::handle:vertical:hover {
        background: #636366;
    }

    QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {
        height: 0px;
    }

    QScrollBar:horizontal {
        border: none;
        background: #1c1c1e;
        height: 10px;
        border-radius: 5px;
    }

    QScrollBar::handle:horizontal {
        background: #48484a;
        border-radius: 5px;
        min-width: 30px;
    }

    QScrollBar::handle:horizontal:hover {
        background: #636366;
    }

    QScrollBar::add-line:horizontal, QScrollBar::sub-line:horizontal {
        width: 0px;
    }
"""


SETTINGS_FILENAME = "settings.json"
APP_NAME = "AdvancedFileSearcher"


def levenshtein_distance(s1, s2):
    """計算兩個字串之間的 Levenshtein 距離"""
    if len(s1) < len(s2):
        return levenshtein_distance(s2, s1)
    if len(s2) == 0:
        return len(s1)
    previous_row = range(len(s2) + 1)
    for i, c1 in enumerate(s1):
        current_row = [i + 1]
        for j, c2 in enumerate(s2):
            insertions = previous_row[j + 1] + 1
            deletions = current_row[j] + 1
            substitutions = previous_row[j] + (c1 != c2)
            current_row.append(min(insertions, deletions, substitutions))
        previous_row = current_row
    return previous_row[-1]


def fuzzy_match(text, pattern, threshold=0.7):
    """
    模糊匹配：檢查 text 是否與 pattern 模糊匹配
    threshold: 相似度閾值 (0-1)，預設 0.7
    """
    if not text or not pattern:
        return False
    text_lower = text.lower()
    pattern_lower = pattern.lower()
    if pattern_lower in text_lower:
        return True
    if len(pattern_lower) <= 3:
        return False
    max_distance = int(len(pattern_lower) * (1 - threshold))
    distance = levenshtein_distance(text_lower, pattern_lower)
    similarity = 1 - (distance / max(len(text_lower), len(pattern_lower)))
    return similarity >= threshold


class BooleanExpressionParser:
    """布林表達式解析器，支援 AND, OR, NOT, 括號分組"""
    
    def __init__(self, expression):
        self.expression = expression
        self.tokens = []
        self.pos = 0
        self._tokenize()
    
    def _tokenize(self):
        """將表達式分割成 token"""
        pattern = r'[()]|\bAND\b|\bOR\b|\bNOT\b|\bNEAR\b|\"[^\"]+\"|[^\s()]+'
        self.tokens = re.findall(pattern, self.expression, re.IGNORECASE)
    
    def parse(self):
        """解析布林表達式，返回查詢條件和參數"""
        if not self.tokens:
            return None, []
        self.pos = 0
        return self._parse_or()
    
    def _parse_or(self):
        """解析 OR 表達式"""
        left_conditions, left_params = self._parse_and()
        conditions = [left_conditions] if left_conditions else []
        params = list(left_params)
        
        while self.pos < len(self.tokens) and self.tokens[self.pos].upper() == 'OR':
            self.pos += 1
            right_conditions, right_params = self._parse_and()
            if right_conditions:
                conditions.append(right_conditions)
                params.extend(right_params)
        
        if len(conditions) > 1:
            return f"({' OR '.join(conditions)})", params
        return conditions[0] if conditions else None, params
    
    def _parse_and(self):
        """解析 AND 表達式"""
        left_conditions, left_params = self._parse_not()
        conditions = [left_conditions] if left_conditions else []
        params = list(left_params)
        
        while self.pos < len(self.tokens) and self.tokens[self.pos].upper() == 'AND':
            self.pos += 1
            right_conditions, right_params = self._parse_not()
            if right_conditions:
                conditions.append(right_conditions)
                params.extend(right_params)
        
        if len(conditions) > 1:
            return f"({' AND '.join(conditions)})", params
        return conditions[0] if conditions else None, params
    
    def _parse_not(self):
        """解析 NOT 表達式"""
        if self.pos < len(self.tokens) and self.tokens[self.pos].upper() == 'NOT':
            self.pos += 1
            term = self._get_term()
            if term:
                return "(filename NOT LIKE ? AND content NOT LIKE ?)", [f"%{term}%", f"%{term}%"]
            return None, []
        return self._parse_term()
    
    def _parse_term(self):
        """解析項 (詞或括號表達式)"""
        if self.pos >= len(self.tokens):
            return None, []
        
        token = self.tokens[self.pos]
        
        if token == '(':
            self.pos += 1
            condition, params = self._parse_or()
            if self.pos < len(self.tokens) and self.tokens[self.pos] == ')':
                self.pos += 1
            return condition, params
        elif token == ')':
            return None, []
        else:
            return self._get_term_condition()
    
    def _get_term(self):
        """獲取當前詞"""
        if self.pos >= len(self.tokens):
            return None
        token = self.tokens[self.pos]
        if token.upper() in ('AND', 'OR', 'NOT', 'NEAR') or token in ('(', ')'):
            return None
        self.pos += 1
        if token.startswith('"') and token.endswith('"'):
            return token[1:-1]
        return token
    
    def _get_term_condition(self):
        """獲取詞的查詢條件"""
        if self.pos >= len(self.tokens):
            return None, []
        token = self.tokens[self.pos]
        if token.upper() in ('AND', 'OR', 'NOT', 'NEAR') or token in ('(', ')'):
            return None, []
        
        self.pos += 1
        
        if token.startswith('"') and token.endswith('"'):
            term = token[1:-1]
        else:
            term = token
        
        return "(filename LIKE ? OR content LIKE ?)", [f"%{term}%", f"%{term}%"]


class UpdateHistoryManager:
    """更新歷史記錄管理器"""

    def __init__(self, db_path):
        self.db_path = db_path
        self._create_table_if_not_exists()

    def _create_table_if_not_exists(self):
        """建立更新歷史記錄資料表"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute("""
                CREATE TABLE IF NOT EXISTS update_history (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    update_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    mode TEXT,
                    files_added INTEGER DEFAULT 0,
                    files_updated INTEGER DEFAULT 0,
                    files_deleted INTEGER DEFAULT 0,
                    total_files INTEGER DEFAULT 0,
                    duration_seconds REAL DEFAULT 0,
                    status TEXT,
                    target_paths TEXT,
                    error_count INTEGER DEFAULT 0,
                    error_messages TEXT
                )
            """)

            cursor.execute(
                "CREATE INDEX IF NOT EXISTS idx_update_time ON update_history(update_time)"
            )
            cursor.execute(
                "CREATE INDEX IF NOT EXISTS idx_mode ON update_history(mode)"
            )

            conn.commit()
            conn.close()
        except Exception as e:
            print(f"建立更新歷史記錄資料表失敗: {e}")

    def add_record(
        self,
        mode,
        files_added=0,
        files_updated=0,
        files_deleted=0,
        total_files=0,
        duration_seconds=0,
        status="success",
        target_paths=None,
        error_count=0,
        error_messages=None,
    ):
        """新增更新記錄"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            target_paths_str = ""
            if target_paths:
                target_paths_str = (
                    "\n".join(target_paths)
                    if isinstance(target_paths, list)
                    else str(target_paths)
                )

            error_messages_str = ""
            if error_messages:
                error_messages_str = (
                    "\n".join(error_messages)
                    if isinstance(error_messages, list)
                    else str(error_messages)
                )

            cursor.execute(
                """
                INSERT INTO update_history 
                (mode, files_added, files_updated, files_deleted, total_files, 
                 duration_seconds, status, target_paths, error_count, error_messages)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
                (
                    mode,
                    files_added,
                    files_updated,
                    files_deleted,
                    total_files,
                    duration_seconds,
                    status,
                    target_paths_str,
                    error_count,
                    error_messages_str,
                ),
            )

            conn.commit()
            conn.close()
            return True
        except Exception as e:
            print(f"新增更新記錄失敗: {e}")
            return False

    def get_history(
        self, limit=50, mode_filter=None, status_filter=None, days_limit=30
    ):
        """取得更新歷史記錄"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            query = """
                SELECT id, update_time, mode, files_added, files_updated, files_deleted,
                       total_files, duration_seconds, status, target_paths, error_count
                FROM update_history
                WHERE update_time >= datetime('now', ?)
            """
            params = [f"-{days_limit} days"]

            if mode_filter:
                query += " AND mode = ?"
                params.append(mode_filter)

            if status_filter:
                query += " AND status = ?"
                params.append(status_filter)

            query += " ORDER BY update_time DESC LIMIT ?"
            params.append(limit)

            cursor.execute(query, params)
            records = cursor.fetchall()
            conn.close()

            return records
        except Exception as e:
            print(f"取得更新歷史記錄失敗: {e}")
            return []

    def get_statistics(self, days_limit=30):
        """取得更新統計資訊"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute(
                """
                SELECT 
                    COUNT(*) as total_updates,
                    SUM(files_added) as total_added,
                    SUM(files_updated) as total_updated,
                    SUM(files_deleted) as total_deleted,
                    AVG(duration_seconds) as avg_duration,
                    SUM(error_count) as total_errors,
                    COUNT(CASE WHEN status = 'success' THEN 1 END) as success_count,
                    COUNT(CASE WHEN status = 'partial' THEN 1 END) as partial_count,
                    COUNT(CASE WHEN status = 'failed' THEN 1 END) as failed_count
                FROM update_history
                WHERE update_time >= datetime('now', ?)
            """,
                (f"-{days_limit} days",),
            )

            stats = cursor.fetchone()
            conn.close()

            return {
                "total_updates": stats[0] or 0,
                "total_added": stats[1] or 0,
                "total_updated": stats[2] or 0,
                "total_deleted": stats[3] or 0,
                "avg_duration": stats[4] or 0,
                "total_errors": stats[5] or 0,
                "success_count": stats[6] or 0,
                "partial_count": stats[7] or 0,
                "failed_count": stats[8] or 0,
            }
        except Exception as e:
            print(f"取得更新統計失敗: {e}")
            return {}

    def clear_old_records(self, days=30):
        """清除舊的歷史記錄"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute(
                """
                DELETE FROM update_history
                WHERE update_time < datetime('now', ?)
            """,
                (f"-{days} days",),
            )

            deleted_count = cursor.rowcount
            conn.commit()
            conn.close()

            return deleted_count
        except Exception as e:
            print(f"清除舊歷史記錄失敗: {e}")
            return 0


class SearchHistoryManager:
    """搜尋歷史記錄管理器"""

    def __init__(self, db_path):
        self.db_path = db_path
        self._create_table_if_not_exists()

    def _create_table_if_not_exists(self):
        """建立搜尋歷史記錄資料表"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute("""
                CREATE TABLE IF NOT EXISTS search_history (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    search_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    filename TEXT,
                    content TEXT,
                    file_type TEXT,
                    boolean_expr TEXT,
                    path_filter TEXT,
                    use_regex_filename INTEGER DEFAULT 0,
                    use_regex_content INTEGER DEFAULT 0,
                    use_boolean INTEGER DEFAULT 0,
                    use_fuzzy INTEGER DEFAULT 0,
                    result_count INTEGER DEFAULT 0,
                    search_duration REAL DEFAULT 0
                )
            """)

            cursor.execute(
                "CREATE INDEX IF NOT EXISTS idx_search_time ON search_history(search_time)"
            )

            conn.commit()
            conn.close()
        except Exception as e:
            print(f"建立搜尋歷史記錄資料表失敗: {e}")

    def add_record(self, search_config, result_count=0, search_duration=0):
        """新增搜尋記錄"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute(
                """
                INSERT INTO search_history 
                (filename, content, file_type, boolean_expr, path_filter,
                 use_regex_filename, use_regex_content, use_boolean, use_fuzzy,
                 result_count, search_duration)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
                (
                    search_config.get("filename", ""),
                    search_config.get("content", ""),
                    search_config.get("file_type", ""),
                    search_config.get("boolean_expr", ""),
                    search_config.get("path_filter", ""),
                    int(search_config.get("use_regex_filename", False)),
                    int(search_config.get("use_regex_content", False)),
                    int(search_config.get("boolean_mode", False)),
                    int(search_config.get("use_fuzzy", False)),
                    result_count,
                    search_duration,
                ),
            )

            conn.commit()
            conn.close()
            return True
        except Exception as e:
            print(f"新增搜尋記錄失敗: {e}")
            return False

    def get_history(self, limit=50):
        """取得搜尋歷史記錄"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute(
                """
                SELECT id, search_time, filename, content, file_type, boolean_expr,
                       path_filter, use_regex_filename, use_regex_content, use_boolean,
                       use_fuzzy, result_count, search_duration
                FROM search_history
                ORDER BY search_time DESC LIMIT ?
            """,
                (limit,),
            )

            records = cursor.fetchall()
            conn.close()

            return records
        except Exception as e:
            print(f"取得搜尋歷史記錄失敗: {e}")
            return []

    def get_unique_searches(self, limit=20):
        """取得唯一的搜尋條件（用於下拉選單）"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute(
                """
                SELECT DISTINCT filename, content, file_type, boolean_expr, path_filter,
                       use_regex_filename, use_regex_content, use_boolean, use_fuzzy,
                       MAX(search_time) as last_used, COUNT(*) as use_count
                FROM search_history
                WHERE filename != '' OR content != '' OR boolean_expr != ''
                GROUP BY filename, content, file_type, boolean_expr, path_filter,
                         use_regex_filename, use_regex_content, use_boolean, use_fuzzy
                ORDER BY last_used DESC LIMIT ?
            """,
                (limit,),
            )

            records = cursor.fetchall()
            conn.close()

            return records
        except Exception as e:
            print(f"取得唯一搜尋條件失敗: {e}")
            return []

    def clear_all_records(self):
        """清除所有搜尋歷史記錄"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute("DELETE FROM search_history")

            deleted_count = cursor.rowcount
            conn.commit()
            conn.close()

            return deleted_count
        except Exception as e:
            print(f"清除搜尋歷史記錄失敗: {e}")
            return 0

    def get_suggestions(self, prefix, limit=10):
        """根據前綴取得搜尋建議"""
        suggestions = []
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            prefix_pattern = f"%{prefix}%"
            cursor.execute("""
                SELECT DISTINCT filename FROM search_history
                WHERE filename IS NOT NULL AND filename != '' AND filename LIKE ?
                ORDER BY search_time DESC LIMIT ?
            """, (prefix_pattern, limit))
            for row in cursor.fetchall():
                if row[0] and row[0] not in suggestions:
                    suggestions.append(row[0])
            cursor.execute("""
                SELECT DISTINCT content FROM search_history
                WHERE content IS NOT NULL AND content != '' AND content LIKE ?
                ORDER BY search_time DESC LIMIT ?
            """, (prefix_pattern, limit))
            for row in cursor.fetchall():
                if row[0] and row[0] not in suggestions:
                    suggestions.append(row[0])
            cursor.execute("""
                SELECT DISTINCT boolean_expr FROM search_history
                WHERE boolean_expr IS NOT NULL AND boolean_expr != '' AND boolean_expr LIKE ?
                ORDER BY search_time DESC LIMIT ?
            """, (prefix_pattern, limit))
            for row in cursor.fetchall():
                if row[0] and row[0] not in suggestions:
                    suggestions.append(row[0])
            conn.close()
        except Exception as e:
            print(f"取得搜尋建議失敗: {e}")
        return suggestions[:limit]


class UpdateHistoryDialog(QDialog):
    """更新歷史記錄查看對話框"""

    def __init__(self, history_manager, parent=None):
        super().__init__(parent)
        self.history_manager = history_manager
        self.setWindowTitle("更新歷史記錄")
        self.setMinimumWidth(800)
        self.setMinimumHeight(600)

        self._init_ui()
        self._load_history()

    def _init_ui(self):
        layout = QVBoxLayout()

        filter_group = QGroupBox("篩選選項")
        filter_layout = QHBoxLayout()

        self.mode_filter_combo = QComboBox()
        self.mode_filter_combo.addItems(["全部", "全量索引", "增量更新"])
        filter_layout.addWidget(QLabel("更新模式:"))
        filter_layout.addWidget(self.mode_filter_combo)

        self.status_filter_combo = QComboBox()
        self.status_filter_combo.addItems(["全部", "成功", "部分成功", "失敗"])
        filter_layout.addWidget(QLabel("狀態:"))
        filter_layout.addWidget(self.status_filter_combo)

        self.days_filter_spin = QSpinBox()
        self.days_filter_spin.setRange(1, 365)
        self.days_filter_spin.setValue(30)
        self.days_filter_spin.setSuffix(" 天")
        filter_layout.addWidget(QLabel("時間範圍:"))
        filter_layout.addWidget(self.days_filter_spin)

        refresh_button = QPushButton("重新載入")
        refresh_button.clicked.connect(self._load_history)
        filter_layout.addWidget(refresh_button)

        filter_layout.addStretch()
        filter_group.setLayout(filter_layout)
        layout.addWidget(filter_group)

        stats_group = QGroupBox("統計摘要")
        stats_layout = QGridLayout()

        self.stats_labels = {}
        stats_items = [
            ("total_updates", "總更新次數"),
            ("success_count", "成功次數"),
            ("total_added", "新增檔案總數"),
            ("total_updated", "更新檔案總數"),
            ("total_deleted", "刪除檔案總數"),
            ("avg_duration", "平均耗時"),
            ("total_errors", "錯誤總數"),
        ]

        for i, (key, label) in enumerate(stats_items):
            row = i // 4
            col = i % 4
            stats_layout.addWidget(QLabel(f"{label}:"), row, col * 2)
            self.stats_labels[key] = QLabel("0")
            self.stats_labels[key].setStyleSheet("font-weight: bold; color: #007aff;")
            stats_layout.addWidget(self.stats_labels[key], row, col * 2 + 1)

        stats_group.setLayout(stats_layout)
        layout.addWidget(stats_group)

        table_group = QGroupBox("歷史記錄")
        table_layout = QVBoxLayout()

        self.history_table = QTableWidget(0, 10)
        self.history_table.setHorizontalHeaderLabels(
            [
                "時間",
                "模式",
                "新增",
                "更新",
                "刪除",
                "總檔案數",
                "耗時",
                "狀態",
                "目標路徑",
                "錯誤數",
            ]
        )
        self.history_table.horizontalHeader().setSectionResizeMode(
            0, QHeaderView.ResizeToContents
        )
        self.history_table.horizontalHeader().setSectionResizeMode(
            8, QHeaderView.Stretch
        )
        self.history_table.setSelectionBehavior(QTableWidget.SelectRows)
        self.history_table.setEditTriggers(QTableWidget.NoEditTriggers)
        self.history_table.setAlternatingRowColors(True)

        self.history_table.setContextMenuPolicy(Qt.CustomContextMenu)
        self.history_table.customContextMenuRequested.connect(self._show_context_menu)

        table_layout.addWidget(self.history_table)
        table_group.setLayout(table_layout)
        layout.addWidget(table_group)

        buttons = QDialogButtonBox(QDialogButtonBox.Close)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)

        self.setLayout(layout)

        self.mode_filter_combo.currentIndexChanged.connect(self._load_history)
        self.status_filter_combo.currentIndexChanged.connect(self._load_history)
        self.days_filter_spin.valueChanged.connect(self._load_history)

    def _load_history(self):
        mode_filter = None
        if self.mode_filter_combo.currentIndex() == 1:
            mode_filter = "full"
        elif self.mode_filter_combo.currentIndex() == 2:
            mode_filter = "update"

        status_filter = None
        if self.status_filter_combo.currentIndex() == 1:
            status_filter = "success"
        elif self.status_filter_combo.currentIndex() == 2:
            status_filter = "partial"
        elif self.status_filter_combo.currentIndex() == 3:
            status_filter = "failed"

        days_limit = self.days_filter_spin.value()

        records = self.history_manager.get_history(
            limit=100,
            mode_filter=mode_filter,
            status_filter=status_filter,
            days_limit=days_limit,
        )

        self.history_table.setRowCount(len(records))

        for i, record in enumerate(records):
            update_time = record[1]
            if isinstance(update_time, str):
                time_str = update_time
            else:
                try:
                    time_str = datetime.datetime.fromtimestamp(update_time).strftime(
                        "%Y-%m-%d %H:%M:%S"
                    )
                except (ValueError, OSError):
                    time_str = str(update_time)

            mode = record[2]
            mode_display = "全量索引" if mode == "full" else "增量更新"

            self.history_table.setItem(i, 0, QTableWidgetItem(time_str))
            self.history_table.setItem(i, 1, QTableWidgetItem(mode_display))
            self.history_table.setItem(i, 2, QTableWidgetItem(str(record[3])))
            self.history_table.setItem(i, 3, QTableWidgetItem(str(record[4])))
            self.history_table.setItem(i, 4, QTableWidgetItem(str(record[5])))
            self.history_table.setItem(i, 5, QTableWidgetItem(str(record[6])))

            duration = record[7]
            duration_str = (
                f"{duration:.2f} 秒" if duration < 60 else f"{duration / 60:.2f} 分"
            )
            self.history_table.setItem(i, 6, QTableWidgetItem(duration_str))

            status = record[8]
            status_display = {
                "success": "成功",
                "partial": "部分成功",
                "failed": "失敗",
            }.get(status, status)
            status_item = QTableWidgetItem(status_display)
            if status == "success":
                status_item.setForeground(QBrush(QColor(0, 128, 0)))
            elif status == "failed":
                status_item.setForeground(QBrush(QColor(255, 0, 0)))
            else:
                status_item.setForeground(QBrush(QColor(255, 165, 0)))
            self.history_table.setItem(i, 7, status_item)

            target_paths = record[9] or ""
            target_paths_display = (
                target_paths[:100] + "..." if len(target_paths) > 100 else target_paths
            )
            self.history_table.setItem(i, 8, QTableWidgetItem(target_paths_display))

            self.history_table.setItem(i, 9, QTableWidgetItem(str(record[10])))

        self.history_table.resizeColumnsToContents()

        stats = self.history_manager.get_statistics(days_limit=days_limit)
        for key, label in self.stats_labels.items():
            value = stats.get(key, 0)
            if key == "avg_duration":
                value_str = f"{value:.2f} 秒" if value < 60 else f"{value / 60:.2f} 分"
            else:
                value_str = str(int(value))
            label.setText(value_str)

    def _show_context_menu(self, position):
        row = self.history_table.indexAt(position).row()
        if row < 0:
            return

        menu = QMenu()
        details_action = QAction("查看詳細資訊", self)
        details_action.triggered.connect(lambda: self._show_record_details(row))
        menu.addAction(details_action)

        menu.exec_(self.history_table.viewport().mapToGlobal(position))

    def _show_record_details(self, row):
        mode_filter = None
        if self.mode_filter_combo.currentIndex() == 1:
            mode_filter = "full"
        elif self.mode_filter_combo.currentIndex() == 2:
            mode_filter = "update"

        status_filter = None
        if self.status_filter_combo.currentIndex() == 1:
            status_filter = "success"
        elif self.status_filter_combo.currentIndex() == 2:
            status_filter = "partial"
        elif self.status_filter_combo.currentIndex() == 3:
            status_filter = "failed"

        days_limit = self.days_filter_spin.value()
        records = self.history_manager.get_history(
            limit=100,
            mode_filter=mode_filter,
            status_filter=status_filter,
            days_limit=days_limit,
        )

        if row >= len(records):
            return

        record = records[row]
        record_id = record[0]

        try:
            conn = sqlite3.connect(self.history_manager.db_path)
            cursor = conn.cursor()
            cursor.execute("SELECT * FROM update_history WHERE id = ?", (record_id,))
            full_record = cursor.fetchone()
            conn.close()

            details_text = f"時間: {full_record[1]}\n"
            details_text += (
                f"模式: {'全量索引' if full_record[2] == 'full' else '增量更新'}\n"
            )
            details_text += f"新增檔案: {full_record[3]}\n"
            details_text += f"更新檔案: {full_record[4]}\n"
            details_text += f"刪除檔案: {full_record[5]}\n"
            details_text += f"總檔案數: {full_record[6]}\n"
            details_text += f"耗時: {full_record[7]:.2f} 秒\n"
            details_text += f"狀態: {full_record[8]}\n"
            details_text += f"\n目標路徑:\n{full_record[9] or '無'}\n"
            details_text += f"\n錯誤數: {full_record[11]}\n"
            if full_record[12]:
                details_text += f"\n錯誤訊息:\n{full_record[12]}"

            QMessageBox.information(self, "詳細資訊", details_text)
        except Exception as e:
            QMessageBox.warning(self, "錯誤", f"無法讀取詳細資訊: {str(e)}")


class UpdateSummaryDialog(QDialog):
    """更新摘要對話框"""

    def __init__(self, stats, parent=None):
        super().__init__(parent)
        self.stats = stats
        self.setWindowTitle("更新摘要")
        self.setMinimumWidth(400)
        self.setModal(True)

        self._init_ui()

    def _init_ui(self):
        layout = QVBoxLayout()

        summary_group = QGroupBox("更新結果摘要")
        summary_layout = QGridLayout()

        mode = self.stats.get("mode", "unknown")
        mode_display = "全量索引" if mode == "full" else "增量更新"
        summary_layout.addWidget(QLabel("更新模式:"), 0, 0)
        mode_label = QLabel(mode_display)
        mode_label.setStyleSheet("font-weight: bold; color: #007aff;")
        summary_layout.addWidget(mode_label, 0, 1)

        duration = self.stats.get("duration", 0)
        duration_str = (
            f"{duration:.2f} 秒" if duration < 60 else f"{duration / 60:.2f} 分"
        )
        summary_layout.addWidget(QLabel("執行耗時:"), 1, 0)
        duration_label = QLabel(duration_str)
        duration_label.setStyleSheet("font-weight: bold;")
        summary_layout.addWidget(duration_label, 1, 1)

        files_added = self.stats.get("files_added", 0)
        summary_layout.addWidget(QLabel("新增檔案:"), 2, 0)
        added_label = QLabel(str(files_added))
        added_label.setStyleSheet("font-weight: bold; color: #28a745;")
        summary_layout.addWidget(added_label, 2, 1)

        files_updated = self.stats.get("files_updated", 0)
        summary_layout.addWidget(QLabel("更新檔案:"), 3, 0)
        updated_label = QLabel(str(files_updated))
        updated_label.setStyleSheet("font-weight: bold; color: #17a2b8;")
        summary_layout.addWidget(updated_label, 3, 1)

        files_deleted = self.stats.get("files_deleted", 0)
        summary_layout.addWidget(QLabel("刪除檔案:"), 4, 0)
        deleted_label = QLabel(str(files_deleted))
        deleted_label.setStyleSheet("font-weight: bold; color: #dc3545;")
        summary_layout.addWidget(deleted_label, 4, 1)

        total_files = self.stats.get("total_files", 0)
        summary_layout.addWidget(QLabel("總處理檔案:"), 5, 0)
        total_label = QLabel(str(total_files))
        total_label.setStyleSheet("font-weight: bold;")
        summary_layout.addWidget(total_label, 5, 1)

        error_count = self.stats.get("error_count", 0)
        summary_layout.addWidget(QLabel("錯誤數量:"), 6, 0)
        error_label = QLabel(str(error_count))
        if error_count > 0:
            error_label.setStyleSheet("font-weight: bold; color: #dc3545;")
        else:
            error_label.setStyleSheet("font-weight: bold; color: #28a745;")
        summary_layout.addWidget(error_label, 6, 1)

        summary_group.setLayout(summary_layout)
        layout.addWidget(summary_group)

        if error_count > 0 and self.stats.get("error_messages"):
            errors_group = QGroupBox("錯誤訊息")
            errors_layout = QVBoxLayout()

            error_text = QTextEdit()
            error_text.setReadOnly(True)
            error_text.setPlainText("\n".join(self.stats["error_messages"][:10]))
            if len(self.stats["error_messages"]) > 10:
                error_text.append(
                    f"\n... (還有 {len(self.stats['error_messages']) - 10} 個錯誤)"
                )
            errors_layout.addWidget(error_text)

            errors_group.setLayout(errors_layout)
            layout.addWidget(errors_group)

        buttons = QDialogButtonBox(QDialogButtonBox.Ok)
        buttons.accepted.connect(self.accept)
        layout.addWidget(buttons)

        self.setLayout(layout)


class AdvancedExcludeRules:
    """進階排除規則管理"""

    def __init__(self):
        self.exclude_extensions = []
        self.min_size_mb = None
        self.max_size_mb = None
        self.exclude_older_than_days = None
        self.exclude_newer_than_days = None
        self.exclude_patterns = []
        self.exclude_dirs = []

    def should_exclude(self, file_path, file_stat=None):
        """檢查檔案是否應被排除"""
        if not file_stat:
            try:
                file_stat = os.stat(file_path)
            except Exception:
                return False

        file_name = os.path.basename(file_path)
        _, ext = os.path.splitext(file_name)
        ext = ext.lower()

        if self.exclude_extensions and ext in self.exclude_extensions:
            return True

        if self.min_size_mb is not None:
            file_size_mb = file_stat.st_size / (1024 * 1024)
            if file_size_mb < self.min_size_mb:
                return True

        if self.max_size_mb is not None:
            file_size_mb = file_stat.st_size / (1024 * 1024)
            if file_size_mb > self.max_size_mb:
                return True

        if self.exclude_older_than_days is not None:
            modified_time = datetime.datetime.fromtimestamp(file_stat.st_mtime)
            age_days = (datetime.datetime.now() - modified_time).days
            if age_days > self.exclude_older_than_days:
                return True

        if self.exclude_newer_than_days is not None:
            modified_time = datetime.datetime.fromtimestamp(file_stat.st_mtime)
            age_days = (datetime.datetime.now() - modified_time).days
            if age_days < self.exclude_newer_than_days:
                return True

        for pattern in self.exclude_patterns:
            try:
                if re.search(pattern, file_path):
                    return True
            except re.error:
                continue

        dir_path = os.path.dirname(file_path)
        for exclude_dir in self.exclude_dirs:
            if dir_path.startswith(exclude_dir) or file_path.startswith(exclude_dir):
                return True

        return False

    def to_dict(self):
        """轉換為字典格式"""
        return {
            "exclude_extensions": self.exclude_extensions,
            "min_size_mb": self.min_size_mb,
            "max_size_mb": self.max_size_mb,
            "exclude_older_than_days": self.exclude_older_than_days,
            "exclude_newer_than_days": self.exclude_newer_than_days,
            "exclude_patterns": self.exclude_patterns,
            "exclude_dirs": self.exclude_dirs,
        }

    def from_dict(self, data):
        """從字典載入"""
        self.exclude_extensions = data.get("exclude_extensions", [])
        self.min_size_mb = data.get("min_size_mb")
        self.max_size_mb = data.get("max_size_mb")
        self.exclude_older_than_days = data.get("exclude_older_than_days")
        self.exclude_newer_than_days = data.get("exclude_newer_than_days")
        self.exclude_patterns = data.get("exclude_patterns", [])
        self.exclude_dirs = data.get("exclude_dirs", [])

    def clear(self):
        """清除所有規則"""
        self.exclude_extensions = []
        self.min_size_mb = None
        self.max_size_mb = None
        self.exclude_older_than_days = None
        self.exclude_newer_than_days = None
        self.exclude_patterns = []
        self.exclude_dirs = []


class AdvancedExcludeDialog(QDialog):
    """進階排除規則設定對話框"""

    def __init__(self, current_rules=None, parent=None):
        super().__init__(parent)
        self.setWindowTitle("進階排除規則設定")
        self.setMinimumWidth(500)
        self.setModal(True)

        self.rules = current_rules or AdvancedExcludeRules()

        self._init_ui()

    def _init_ui(self):
        layout = QVBoxLayout()

        extensions_group = QGroupBox("副檔名排除")
        extensions_layout = QVBoxLayout()

        extensions_layout.addWidget(
            QLabel("排除特定副檔名的檔案 (逗號分隔，如: .tmp, .bak):")
        )
        self.extensions_input = QLineEdit()
        self.extensions_input.setText(", ".join(self.rules.exclude_extensions))
        extensions_layout.addWidget(self.extensions_input)

        extensions_group.setLayout(extensions_layout)
        layout.addWidget(extensions_group)

        size_group = QGroupBox("檔案大小排除")
        size_layout = QGridLayout()

        self.exclude_min_size = QCheckBox("排除小於")
        size_layout.addWidget(self.exclude_min_size, 0, 0)
        self.min_size_spin = QSpinBox()
        self.min_size_spin.setRange(1, 1000000)
        self.min_size_spin.setSuffix(" MB")
        self.min_size_spin.setValue(self.rules.min_size_mb or 1)
        if self.rules.min_size_mb is not None:
            self.exclude_min_size.setChecked(True)
        self.min_size_spin.setEnabled(self.rules.min_size_mb is not None)
        self.exclude_min_size.toggled.connect(self.min_size_spin.setEnabled)
        size_layout.addWidget(self.min_size_spin, 0, 1)

        self.exclude_max_size = QCheckBox("排除大於")
        size_layout.addWidget(self.exclude_max_size, 1, 0)
        self.max_size_spin = QSpinBox()
        self.max_size_spin.setRange(1, 1000000)
        self.max_size_spin.setSuffix(" MB")
        self.max_size_spin.setValue(self.rules.max_size_mb or 100)
        if self.rules.max_size_mb is not None:
            self.exclude_max_size.setChecked(True)
        self.max_size_spin.setEnabled(self.rules.max_size_mb is not None)
        self.exclude_max_size.toggled.connect(self.max_size_spin.setEnabled)
        size_layout.addWidget(self.max_size_spin, 1, 1)

        size_group.setLayout(size_layout)
        layout.addWidget(size_group)

        date_group = QGroupBox("檔案日期排除")
        date_layout = QGridLayout()

        self.exclude_older = QCheckBox("排除超過")
        date_layout.addWidget(self.exclude_older, 0, 0)
        self.older_days_spin = QSpinBox()
        self.older_days_spin.setRange(1, 3650)
        self.older_days_spin.setSuffix(" 天未修改")
        self.older_days_spin.setValue(self.rules.exclude_older_than_days or 365)
        if self.rules.exclude_older_than_days is not None:
            self.exclude_older.setChecked(True)
        self.older_days_spin.setEnabled(self.rules.exclude_older_than_days is not None)
        self.exclude_older.toggled.connect(self.older_days_spin.setEnabled)
        date_layout.addWidget(self.older_days_spin, 0, 1)

        self.exclude_newer = QCheckBox("排除最近")
        date_layout.addWidget(self.exclude_newer, 1, 0)
        self.newer_days_spin = QSpinBox()
        self.newer_days_spin.setRange(1, 3650)
        self.newer_days_spin.setSuffix(" 天內修改")
        self.newer_days_spin.setValue(self.rules.exclude_newer_than_days or 7)
        if self.rules.exclude_newer_than_days is not None:
            self.exclude_newer.setChecked(True)
        self.newer_days_spin.setEnabled(self.rules.exclude_newer_than_days is not None)
        self.exclude_newer.toggled.connect(self.newer_days_spin.setEnabled)
        date_layout.addWidget(self.newer_days_spin, 1, 1)

        date_group.setLayout(date_layout)
        layout.addWidget(date_group)

        dirs_group = QGroupBox("特定目錄排除")
        dirs_layout = QVBoxLayout()

        dirs_layout.addWidget(QLabel("排除特定目錄下的所有檔案:"))
        self.dirs_list = QTextEdit()
        self.dirs_list.setPlaceholderText("每行一個目錄路徑")
        self.dirs_list.setText("\n".join(self.rules.exclude_dirs))
        self.dirs_list.setMaximumHeight(100)
        dirs_layout.addWidget(self.dirs_list)

        add_dir_button = QPushButton("新增目錄")
        add_dir_button.clicked.connect(self._add_exclude_dir)
        dirs_layout.addWidget(add_dir_button)

        dirs_group.setLayout(dirs_layout)
        layout.addWidget(dirs_group)

        patterns_group = QGroupBox("正則表達式排除")
        patterns_layout = QVBoxLayout()

        patterns_layout.addWidget(QLabel("使用正則表達式排除檔案 (逗號分隔):"))
        self.patterns_input = QLineEdit()
        self.patterns_input.setText(", ".join(self.rules.exclude_patterns))
        self.patterns_input.setPlaceholderText("例如: .*\\.git.*, .*\\.svn.*")
        patterns_layout.addWidget(self.patterns_input)

        patterns_group.setLayout(patterns_layout)
        layout.addWidget(patterns_group)

        buttons = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
        buttons.accepted.connect(self._accept_rules)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)

        self.setLayout(layout)

    def _add_exclude_dir(self):
        directory = QFileDialog.getExistingDirectory(self, "選擇要排除的目錄")
        if directory:
            current_text = self.dirs_list.toPlainText()
            if current_text:
                self.dirs_list.setPlainText(f"{current_text}\n{directory}")
            else:
                self.dirs_list.setPlainText(directory)

    def _accept_rules(self):
        self.rules.exclude_extensions = [
            ext.strip().lower()
            if ext.strip().startswith(".")
            else "." + ext.strip().lower()
            for ext in self.extensions_input.text().split(",")
            if ext.strip()
        ]

        self.rules.min_size_mb = (
            self.min_size_spin.value() if self.exclude_min_size.isChecked() else None
        )
        self.rules.max_size_mb = (
            self.max_size_spin.value() if self.exclude_max_size.isChecked() else None
        )

        self.rules.exclude_older_than_days = (
            self.older_days_spin.value() if self.exclude_older.isChecked() else None
        )
        self.rules.exclude_newer_than_days = (
            self.newer_days_spin.value() if self.exclude_newer.isChecked() else None
        )

        self.rules.exclude_dirs = [
            d.strip() for d in self.dirs_list.toPlainText().split("\n") if d.strip()
        ]

        self.rules.exclude_patterns = [
            p.strip() for p in self.patterns_input.text().split(",") if p.strip()
        ]

        self.accept()

    def get_rules(self):
        return self.rules


class IndexCleanupDialog(QDialog):
    """索引清理對話框"""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("索引清理選項")
        self.setModal(True)
        self.setMinimumWidth(400)

        layout = QVBoxLayout()

        # 清理選項
        self.remove_deleted = QCheckBox("移除不存在的檔案")
        self.remove_deleted.setChecked(True)
        layout.addWidget(self.remove_deleted)

        self.remove_by_path = QCheckBox("移除指定路徑下的索引")
        layout.addWidget(self.remove_by_path)

        self.path_input = QLineEdit()
        self.path_input.setPlaceholderText("輸入要移除的路徑前綴")
        self.path_input.setEnabled(False)
        layout.addWidget(self.path_input)

        self.remove_by_type = QCheckBox("移除指定類型的檔案")
        layout.addWidget(self.remove_by_type)

        self.type_input = QLineEdit()
        self.type_input.setPlaceholderText("輸入副檔名，如: .tmp, .bak")
        self.type_input.setEnabled(False)
        layout.addWidget(self.type_input)

        self.remove_by_size = QCheckBox("移除超過指定大小的檔案")
        layout.addWidget(self.remove_by_size)

        size_layout = QHBoxLayout()
        self.size_input = QSpinBox()
        self.size_input.setRange(1, 1000000)
        self.size_input.setValue(100)
        self.size_input.setSuffix(" MB")
        self.size_input.setEnabled(False)
        size_layout.addWidget(self.size_input)
        size_layout.addStretch()
        layout.addLayout(size_layout)

        self.remove_duplicates = QCheckBox("移除重複的索引項目")
        self.remove_duplicates.setChecked(True)
        layout.addWidget(self.remove_duplicates)

        # 連接信號
        self.remove_by_path.toggled.connect(self.path_input.setEnabled)
        self.remove_by_type.toggled.connect(self.type_input.setEnabled)
        self.remove_by_size.toggled.connect(self.size_input.setEnabled)

        # 按鈕
        buttons = QDialogButtonBox(
            QDialogButtonBox.Ok | QDialogButtonBox.Cancel, Qt.Horizontal, self
        )
        buttons.accepted.connect(self.accept)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)

        self.setLayout(layout)

    def get_options(self):
        """取得清理選項"""
        return {
            "remove_deleted": self.remove_deleted.isChecked(),
            "remove_by_path": self.remove_by_path.isChecked(),
            "path_prefix": self.path_input.text()
            if self.remove_by_path.isChecked()
            else "",
            "remove_by_type": self.remove_by_type.isChecked(),
            "file_types": [
                t.strip() for t in self.type_input.text().split(",") if t.strip()
            ]
            if self.remove_by_type.isChecked()
            else [],
            "remove_by_size": self.remove_by_size.isChecked(),
            "max_size_mb": self.size_input.value()
            if self.remove_by_size.isChecked()
            else 0,
            "remove_duplicates": self.remove_duplicates.isChecked(),
        }


class DuplicateFilesDialog(QDialog):
    """重複檔案偵測結果對話框"""

    def __init__(self, dup_groups, parent=None):
        super().__init__(parent)
        self.dup_groups = dup_groups
        self.setWindowTitle("重複檔案偵測結果")
        self.setModal(True)
        self.setMinimumSize(700, 500)
        self._init_ui()

    def _init_ui(self):
        layout = QVBoxLayout()

        # 統計標籤
        total_dup_files = sum(len(files) for _, files in self.dup_groups)
        total_dup_size = sum(f[2] for _, group in self.dup_groups for f in group)
        saveable = total_dup_size - sum(
            min(f[2] for f in group) for _, group in self.dup_groups
        )
        info_label = QLabel(
            f"找到 <b>{len(self.dup_groups)}</b> 個重複群組，"
            f"共 <b>{total_dup_files}</b> 個重複檔案，"
            f"刪除後可節省約 <b>{self._format_size(saveable)}</b>"
        )
        info_label.setWordWrap(True)
        layout.addWidget(info_label)

        # 樹形列表
        self.tree = QTreeWidget()
        self.tree.setHeaderLabels(["檔名", "路徑", "大小", "Hash"])
        self.tree.setSelectionMode(QTreeWidget.MultiSelection)

        for group_hash, files in self.dup_groups:
            group_item = QTreeWidgetItem()
            group_item.setText(0, f"[{len(files)} 個重複] — Hash: {group_hash[:8]}...")
            for path, filename, size in files:
                child = QTreeWidgetItem()
                child.setText(0, filename)
                child.setText(1, path)
                child.setText(2, self._format_size(size))
                child.setText(3, group_hash)
                child.setFlags(child.flags() | Qt.ItemIsUserCheckable)
                child.setCheckState(0, Qt.Unchecked)
                group_item.addChild(child)
            self.tree.addTopLevelItem(group_item)
            group_item.setExpanded(True)

        layout.addWidget(self.tree)

        # 按鈕
        btn_layout = QHBoxLayout()
        select_all_btn = QPushButton("全選")
        select_all_btn.setObjectName("secondaryButton")
        select_all_btn.clicked.connect(self._select_all)
        btn_layout.addWidget(select_all_btn)

        deselect_all_btn = QPushButton("取消全選")
        deselect_all_btn.setObjectName("secondaryButton")
        deselect_all_btn.clicked.connect(self._deselect_all)
        btn_layout.addWidget(deselect_all_btn)

        delete_btn = QPushButton("刪除選中")
        delete_btn.setObjectName("deleteButton")
        delete_btn.clicked.connect(self._delete_selected)
        btn_layout.addWidget(delete_btn)

        close_btn = QPushButton("關閉")
        close_btn.setObjectName("secondaryButton")
        close_btn.clicked.connect(self.reject)
        btn_layout.addWidget(close_btn)

        btn_layout.addStretch()
        layout.addLayout(btn_layout)

        self.setLayout(layout)

    def _select_all(self):
        for i in range(self.tree.topLevelItemCount()):
            group = self.tree.topLevelItem(i)
            for j in range(group.childCount()):
                group.child(j).setCheckState(0, Qt.Checked)

    def _deselect_all(self):
        for i in range(self.tree.topLevelItemCount()):
            group = self.tree.topLevelItem(i)
            for j in range(group.childCount()):
                group.child(j).setCheckState(0, Qt.Unchecked)

    def _delete_selected(self):
        to_delete = []
        for i in range(self.tree.topLevelItemCount()):
            group = self.tree.topLevelItem(i)
            for j in range(group.childCount()):
                child = group.child(j)
                if child.checkState(0) == Qt.Checked:
                    to_delete.append(child.text(1))

        if not to_delete:
            QMessageBox.information(self, "提示", "請先勾選要刪除的檔案。")
            return

        reply = QMessageBox.warning(
            self,
            "確認刪除",
            f"確定要刪除 {len(to_delete)} 個檔案嗎？\n此操作無法復原。",
            QMessageBox.Yes | QMessageBox.No,
        )
        if reply != QMessageBox.Yes:
            return

        deleted = 0
        for path in to_delete:
            try:
                if os.path.exists(path):
                    os.remove(path)
                    deleted += 1
            except Exception as e:
                print(f"刪除 {path} 失敗: {e}")

        # Also remove from database
        try:
            parent = self.parent()
            if parent and hasattr(parent, "db_path"):
                import sqlite3
                conn = sqlite3.connect(parent.db_path)
                cursor = conn.cursor()
                for path in to_delete:
                    cursor.execute("DELETE FROM files WHERE path = ?", (path,))
                conn.commit()
                conn.close()
        except Exception as e:
            print(f"更新資料庫時出錯: {e}")

        QMessageBox.information(self, "刪除完成", f"已刪除 {deleted}/{len(to_delete)} 個檔案。")
        self.accept()

    def _format_size(self, size_bytes):
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.1f} KB"
        elif size_bytes < 1024 * 1024 * 1024:
            return f"{size_bytes / (1024 * 1024):.1f} MB"
        else:
            return f"{size_bytes / (1024 * 1024 * 1024):.2f} GB"


class EmptyFilesDialog(QDialog):
    """空檔案偵測結果對話框"""

    def __init__(self, empty_files, parent=None):
        super().__init__(parent)
        self.empty_files = empty_files
        self.setWindowTitle("空檔案偵測結果")
        self.setModal(True)
        self.setMinimumSize(600, 400)
        self._init_ui()

    def _init_ui(self):
        layout = QVBoxLayout()

        info_label = QLabel(f"找到 <b>{len(self.empty_files)}</b> 個空檔案（大小為 0）")
        layout.addWidget(info_label)

        self.table = QTableWidget(len(self.empty_files), 2)
        self.table.setHorizontalHeaderLabels(["檔名", "路徑"])
        self.table.setSelectionBehavior(QTableWidget.SelectRows)
        self.table.setEditTriggers(QTableWidget.NoEditTriggers)

        for i, (path, filename) in enumerate(self.empty_files):
            self.table.setItem(i, 0, QTableWidgetItem(filename))
            self.table.setItem(i, 1, QTableWidgetItem(path))

        self.table.horizontalHeader().setSectionResizeMode(0, QHeaderView.Stretch)
        self.table.horizontalHeader().setSectionResizeMode(1, QHeaderView.Stretch)
        layout.addWidget(self.table)

        btn_layout = QHBoxLayout()

        open_btn = QPushButton("開啟所選檔案")
        open_btn.setObjectName("secondaryButton")
        open_btn.clicked.connect(self._open_selected)
        btn_layout.addWidget(open_btn)

        delete_btn = QPushButton("刪除所選檔案")
        delete_btn.setObjectName("deleteButton")
        delete_btn.clicked.connect(self._delete_selected)
        btn_layout.addWidget(delete_btn)

        btn_layout.addStretch()

        close_btn = QPushButton("關閉")
        close_btn.setObjectName("secondaryButton")
        close_btn.clicked.connect(self.reject)
        btn_layout.addWidget(close_btn)

        layout.addLayout(btn_layout)
        self.setLayout(layout)

    def _open_selected(self):
        selected = self.table.selectedItems()
        if not selected:
            return
        row = selected[0].row()
        path = self.table.item(row, 1).text()
        if os.path.exists(path):
            import sys
            import subprocess
            if sys.platform == "win32":
                os.startfile(path)
            elif sys.platform == "darwin":
                subprocess.call(["open", path])
            else:
                subprocess.call(["xdg-open", path])

    def _delete_selected(self):
        selected_rows = set(item.row() for item in self.table.selectedItems())
        if not selected_rows:
            QMessageBox.information(self, "提示", "請先選擇要刪除的檔案。")
            return

        reply = QMessageBox.warning(
            self,
            "確認刪除",
            f"確定要刪除 {len(selected_rows)} 個空檔案嗎？",
            QMessageBox.Yes | QMessageBox.No,
        )
        if reply != QMessageBox.Yes:
            return

        deleted = 0
        for row in selected_rows:
            path = self.table.item(row, 1).text()
            try:
                if os.path.exists(path):
                    os.remove(path)
                    deleted += 1
            except Exception as e:
                print(f"刪除 {path} 失敗: {e}")

        # Also remove from database
        try:
            parent = self.parent()
            if parent and hasattr(parent, "db_path"):
                import sqlite3
                conn = sqlite3.connect(parent.db_path)
                cursor = conn.cursor()
                for row in selected_rows:
                    path = self.table.item(row, 1).text()
                    cursor.execute("DELETE FROM files WHERE path = ?", (path,))
                conn.commit()
                conn.close()
        except Exception as e:
            print(f"更新資料庫時出錯: {e}")

        QMessageBox.information(self, "刪除完成", f"已刪除 {deleted}/{len(selected_rows)} 個檔案。")
        self.accept()


class IndexMaintenanceThread(QThread):
    """索引維護執行緒"""

    progress_update = pyqtSignal(str, int)
    maintenance_complete = pyqtSignal(dict)

    def __init__(self, db_path, action, options=None):
        super().__init__()
        self.db_path = db_path
        self.action = action  # 'cleanup', 'verify', 'optimize', 'rebuild'
        self.options = options or {}
        self.is_running = True

    def run(self):
        if self.action == "cleanup":
            self.cleanup_index()
        elif self.action == "verify":
            self.verify_index()
        elif self.action == "optimize":
            self.optimize_database()
        elif self.action == "rebuild":
            self.rebuild_index()

    def cleanup_index(self):
        """清理索引"""
        stats = {
            "deleted_files": 0,
            "removed_by_path": 0,
            "removed_by_type": 0,
            "removed_by_size": 0,
            "removed_duplicates": 0,
            "total_removed": 0,
            "errors": [],
        }

        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            # 1. 移除不存在的檔案
            if self.options.get("remove_deleted"):
                self.progress_update.emit("檢查不存在的檔案...", 10)
                cursor.execute("SELECT id, path FROM files")
                all_files = cursor.fetchall()
                to_delete = []

                for i, (file_id, path) in enumerate(all_files):
                    if not self.is_running:
                        break
                    if i % 100 == 0:
                        progress = 10 + int((i / len(all_files)) * 20)
                        self.progress_update.emit(
                            f"檢查檔案 {i}/{len(all_files)}...", progress
                        )

                    if not os.path.exists(path):
                        to_delete.append(file_id)

                if to_delete:
                    cursor.executemany(
                        "DELETE FROM files WHERE id = ?", [(id,) for id in to_delete]
                    )
                    stats["deleted_files"] = len(to_delete)
                    conn.commit()

            # 2. 移除指定路徑下的索引
            if self.options.get("remove_by_path") and self.options.get("path_prefix"):
                self.progress_update.emit("移除指定路徑索引...", 35)
                path_prefix = self.options["path_prefix"]
                cursor.execute(
                    "SELECT COUNT(*) FROM files WHERE path LIKE ?", (f"{path_prefix}%",)
                )
                count = cursor.fetchone()[0]
                cursor.execute(
                    "DELETE FROM files WHERE path LIKE ?", (f"{path_prefix}%",)
                )
                stats["removed_by_path"] = count
                conn.commit()

            # 3. 移除指定類型的檔案
            if self.options.get("remove_by_type") and self.options.get("file_types"):
                self.progress_update.emit("移除指定類型檔案...", 50)
                for file_type in self.options["file_types"]:
                    if not file_type.startswith("."):
                        file_type = "." + file_type
                    cursor.execute(
                        "SELECT COUNT(*) FROM files WHERE ext = ?", (file_type.lower(),)
                    )
                    count = cursor.fetchone()[0]
                    cursor.execute(
                        "DELETE FROM files WHERE ext = ?", (file_type.lower(),)
                    )
                    stats["removed_by_type"] += count
                conn.commit()

            # 4. 移除超過指定大小的檔案
            if self.options.get("remove_by_size") and self.options.get("max_size_mb"):
                self.progress_update.emit("移除大型檔案索引...", 65)
                max_size_bytes = self.options["max_size_mb"] * 1024 * 1024
                cursor.execute(
                    "SELECT COUNT(*) FROM files WHERE size > ?", (max_size_bytes,)
                )
                count = cursor.fetchone()[0]
                cursor.execute("DELETE FROM files WHERE size > ?", (max_size_bytes,))
                stats["removed_by_size"] = count
                conn.commit()

            # 5. 移除重複的索引項目
            if self.options.get("remove_duplicates"):
                self.progress_update.emit("移除重複項目...", 80)
                cursor.execute("""
                    DELETE FROM files
                    WHERE id NOT IN (
                        SELECT MIN(id)
                        FROM files
                        GROUP BY path
                    )
                """)
                stats["removed_duplicates"] = cursor.rowcount
                conn.commit()

            # 計算總計
            stats["total_removed"] = (
                stats["deleted_files"]
                + stats["removed_by_path"]
                + stats["removed_by_type"]
                + stats["removed_by_size"]
                + stats["removed_duplicates"]
            )

            self.progress_update.emit("清理完成", 100)

        except Exception as e:
            stats["errors"].append(str(e))
            self.progress_update.emit(f"錯誤: {str(e)}", 0)
        finally:
            conn.close()

        self.maintenance_complete.emit(stats)

    def verify_index(self):
        """驗證索引完整性"""
        stats = {
            "total_files": 0,
            "missing_files": 0,
            "corrupted_entries": 0,
            "duplicate_entries": 0,
            "orphaned_content": 0,
            "errors": [],
        }

        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            self.progress_update.emit("檢查索引完整性...", 10)

            # 檢查總檔案數
            cursor.execute("SELECT COUNT(*) FROM files")
            stats["total_files"] = cursor.fetchone()[0]

            # 檢查不存在的檔案
            self.progress_update.emit("檢查檔案存在性...", 30)
            cursor.execute("SELECT path FROM files")
            all_paths = cursor.fetchall()

            for i, (path,) in enumerate(all_paths):
                if not self.is_running:
                    break
                if i % 100 == 0:
                    progress = 30 + int((i / len(all_paths)) * 40)
                    self.progress_update.emit(
                        f"驗證檔案 {i}/{len(all_paths)}...", progress
                    )

                if not os.path.exists(path):
                    stats["missing_files"] += 1

            # 檢查重複項目
            self.progress_update.emit("檢查重複項目...", 75)
            cursor.execute("""
                SELECT path, COUNT(*) as cnt
                FROM files
                GROUP BY path
                HAVING cnt > 1
            """)
            duplicates = cursor.fetchall()
            stats["duplicate_entries"] = sum(cnt - 1 for _, cnt in duplicates)

            # 檢查資料庫完整性
            self.progress_update.emit("檢查資料庫完整性...", 90)
            cursor.execute("PRAGMA integrity_check")
            result = cursor.fetchone()[0]
            if result != "ok":
                stats["errors"].append(f"資料庫完整性問題: {result}")

            self.progress_update.emit("驗證完成", 100)

        except Exception as e:
            stats["errors"].append(str(e))
            self.progress_update.emit(f"錯誤: {str(e)}", 0)
        finally:
            conn.close()

        self.maintenance_complete.emit(stats)

    def optimize_database(self):
        """優化資料庫"""
        stats = {
            "original_size": 0,
            "optimized_size": 0,
            "space_saved": 0,
            "errors": [],
        }

        try:
            # 取得原始大小
            stats["original_size"] = os.path.getsize(self.db_path)

            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            self.progress_update.emit("重建索引...", 20)
            cursor.execute("REINDEX")

            self.progress_update.emit("清理空間...", 50)
            cursor.execute("VACUUM")

            self.progress_update.emit("分析表格...", 80)
            cursor.execute("ANALYZE")

            conn.commit()

            # 取得優化後大小
            stats["optimized_size"] = os.path.getsize(self.db_path)
            stats["space_saved"] = stats["original_size"] - stats["optimized_size"]

            self.progress_update.emit("優化完成", 100)

        except Exception as e:
            stats["errors"].append(str(e))
            self.progress_update.emit(f"錯誤: {str(e)}", 0)
        finally:
            conn.close()

        self.maintenance_complete.emit(stats)

    def rebuild_index(self):
        """重建索引結構"""
        stats = {"tables_rebuilt": 0, "indexes_rebuilt": 0, "errors": []}

        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            self.progress_update.emit("備份資料...", 10)

            # 備份現有資料
            cursor.execute("SELECT * FROM files")
            backup_data = cursor.fetchall()

            self.progress_update.emit("重建表格結構...", 30)

            # 刪除舊表格
            cursor.execute("DROP TABLE IF EXISTS files")

            # 重新建立表格
            cursor.execute("""
            CREATE TABLE files (
                id INTEGER PRIMARY KEY,
                path TEXT UNIQUE,
                filename TEXT,
                ext TEXT,
                size INTEGER,
                modified TIMESTAMP,
                content_indexed BOOLEAN,
                content TEXT,
                created TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                last_verified TIMESTAMP
            )
            """)
            stats["tables_rebuilt"] = 1

            self.progress_update.emit("重建索引...", 50)

            # 重建索引
            cursor.execute("CREATE INDEX idx_filename ON files(filename)")
            cursor.execute("CREATE INDEX idx_ext ON files(ext)")
            cursor.execute("CREATE INDEX idx_path ON files(path)")
            cursor.execute("CREATE INDEX idx_size ON files(size)")
            cursor.execute("CREATE INDEX idx_modified ON files(modified)")
            stats["indexes_rebuilt"] = 5

            self.progress_update.emit("還原資料...", 70)

            # 還原資料（處理舊資料格式）
            for row in backup_data:
                try:
                    # 根據舊表格的欄位數量調整
                    if len(row) >= 8:
                        cursor.execute(
                            """
                            INSERT INTO files (id, path, filename, ext, size, modified, content_indexed, content)
                            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                        """,
                            row[:8],
                        )
                except Exception as e:
                    stats["errors"].append(f"還原資料錯誤: {str(e)}")

            conn.commit()

            self.progress_update.emit("優化資料庫...", 90)
            cursor.execute("VACUUM")
            cursor.execute("ANALYZE")

            conn.commit()

            self.progress_update.emit("重建完成", 100)

        except Exception as e:
            stats["errors"].append(str(e))
            self.progress_update.emit(f"錯誤: {str(e)}", 0)
        finally:
            conn.close()

        self.maintenance_complete.emit(stats)

    def stop(self):
        self.is_running = False


class FileIndexer(QThread):
    progress_update = pyqtSignal(int, int)
    indexing_complete = pyqtSignal(dict)

    def __init__(
        self,
        db_path,
        directories,
        exclude_dirs=None,
        index_content=True,
        content_types=None,
        max_workers=None,
        batch_size=10000,
        mode="full",
        update_targets=None,
        advanced_exclude_rules=None,
    ):
        super().__init__()
        self.db_path = db_path
        self.directories = directories
        self.exclude_dirs = exclude_dirs or []
        self.file_count = 0
        self.is_running = True
        self.batch_size = batch_size
        self.total_estimated = 0

        # 索引內容控制
        self.index_content = bool(index_content)
        self.content_types = (
            self._normalize_ext_list(content_types)
            if content_types
            else set(
                [
                    ".txt",
                    ".py",
                    ".java",
                    ".c",
                    ".cpp",
                    ".html",
                    ".xml",
                    ".json",
                    ".csv",
                    ".md",
                ]
            )
        )
        self.max_workers = max_workers or os.cpu_count()
        print(f"索引執行緒數設定為: {self.max_workers}")  # 除錯輸出

        # 模式：full 或 update
        self.mode = mode
        self.update_targets = set(update_targets or [])

        # 進階排除規則
        self.advanced_exclude_rules = advanced_exclude_rules or AdvancedExcludeRules()

        # 統計資料收集
        self.files_added = 0
        self.files_updated = 0
        self.files_deleted = 0
        self.error_count = 0
        self.error_messages = []
        self.start_time = None

    def _normalize_ext_list(self, types_input):
        if types_input is None:
            return set()
        if isinstance(types_input, str):
            parts = [p.strip() for p in types_input.split(",")]
        else:
            parts = [str(p).strip() for p in types_input]
        norm = set()
        for p in parts:
            if not p:
                continue
            if not p.startswith("."):
                p = "." + p
            norm.add(p.lower())
        return norm

    def run(self):
        self.start_time = time.time()

        conn = sqlite3.connect(self.db_path)
        try:
            cursor = conn.cursor()

            # 建立增強的表格結構
            cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY,
                path TEXT UNIQUE,
                filename TEXT,
                ext TEXT,
                size INTEGER,
                modified TIMESTAMP,
                content_indexed BOOLEAN,
                content TEXT,
                created TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                last_verified TIMESTAMP
            )
            """)
            cursor.execute("CREATE INDEX IF NOT EXISTS idx_filename ON files(filename)")
            cursor.execute("CREATE INDEX IF NOT EXISTS idx_ext ON files(ext)")
            cursor.execute("CREATE INDEX IF NOT EXISTS idx_path ON files(path)")
            cursor.execute("CREATE INDEX IF NOT EXISTS idx_size ON files(size)")
            cursor.execute("CREATE INDEX IF NOT EXISTS idx_modified ON files(modified)")

            if self.mode == "full":
                cursor.execute("DELETE FROM files")
                conn.commit()

            self.file_count = 0
            self.total_estimated = (
                self._estimate_total_files() if self.mode == "full" else 1000
            )
            self.progress_update.emit(0, self.total_estimated)

            all_files = []
            if self.mode == "full":
                with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                    futures = [
                        executor.submit(self._index_directory_collect, d)
                        for d in self.directories
                    ]
                    for f in futures:
                        if not self.is_running:
                            break
                        try:
                            all_files.extend(f.result())
                        except Exception as e:
                            print(f"索引任務執行出錯: {e}")
                            self.error_count += 1
                            self.error_messages.append(str(e))
                if all_files and self.is_running:
                    self._batch_insert_files(conn, cursor, all_files)
            else:
                for root_dir in list(self.update_targets):
                    if not self.is_running:
                        break
                    try:
                        self._incremental_update_for_root(conn, cursor, root_dir)
                    except Exception as e:
                        print(f"更新目標 {root_dir} 時出錯: {e}")
                        self.error_count += 1
                        self.error_messages.append(str(e))

            conn.commit()
        finally:
            conn.close()

        duration = time.time() - self.start_time

        stats = {
            "success": bool(self.is_running),
            "mode": self.mode,
            "files_added": self.files_added,
            "files_updated": self.files_updated,
            "files_deleted": self.files_deleted,
            "total_files": self.file_count,
            "duration": duration,
            "error_count": self.error_count,
            "error_messages": self.error_messages[:10],
        }

        self.indexing_complete.emit(stats)

    def _estimate_total_files(self):
        file_count = 0
        dir_count = 0
        sample_dirs = min(1000, sum(1 for _ in self._walk_directories(sample=True)))
        if sample_dirs == 0:
            return 10000
        for directory in self.directories:
            for root, dirs, files in os.walk(directory, topdown=True):
                if self._should_exclude(root):
                    dirs[:] = []
                    continue
                file_count += len(files)
                dir_count += 1
                if dir_count >= 100:
                    break
            if dir_count >= 100:
                break
        if dir_count == 0:
            return 10000
        avg_files_per_dir = max(file_count / dir_count, 1)
        estimated_total_dirs = sample_dirs * len(self.directories)
        estimated_total_files = int(estimated_total_dirs * avg_files_per_dir)
        return max(estimated_total_files, 1000)

    def _walk_directories(self, sample=False):
        sample_count = 0
        for directory in self.directories:
            try:
                for root, dirs, _ in os.walk(directory, topdown=True):
                    if self._should_exclude(root):
                        dirs[:] = []
                        continue
                    yield root
                    sample_count += 1
                    if sample and sample_count >= 1000:
                        return
            except Exception as e:
                print(f"遍歷目錄 {directory} 時出錯: {e}")

    def _index_directory_collect(self, directory):
        files_info = []
        try:
            for root, dirs, files in os.walk(directory, topdown=True):
                if not self.is_running:
                    return files_info
                if self._should_exclude(root):
                    dirs[:] = []
                    continue
                for file in files:
                    if not self.is_running:
                        return files_info
                    file_path = os.path.join(root, file)
                    rec = self._build_file_record(file_path, file)
                    if rec:
                        files_info.append(rec)
                        self.file_count += 1
                        if self.file_count % 100 == 0:
                            self.progress_update.emit(
                                self.file_count, self.total_estimated
                            )
        except Exception as e:
            print(f"處理目錄 {directory} 時出錯: {e}")
            self.error_count += 1
            self.error_messages.append(str(e))
        return files_info

    def _build_file_record(self, file_path, file_name=None):
        try:
            if file_name is None:
                file_name = os.path.basename(file_path)
            stats = os.stat(file_path)

            if self._should_exclude(file_path, stats):
                return None

            _, ext = os.path.splitext(file_name)
            ext = ext.lower()
            modified = datetime.datetime.fromtimestamp(stats.st_mtime)
            content = ""
            content_indexed = False
            if self.index_content:
                supported_exts = ContentExtractor.get_supported_extensions()
                if ext in supported_exts or ext in self.content_types:
                    try:
                        content = ContentExtractor.extract_text(file_path, ext)
                        content_indexed = True
                    except Exception:
                        content = ""
                        content_indexed = False
            return {
                "path": os.path.normpath(file_path),  # 標準化路徑分隔符
                "filename": file_name,
                "ext": ext,
                "size": stats.st_size,
                "modified": modified,
                "content_indexed": content_indexed,
                "content": content,
            }
        except Exception as e:
            print(f"處理檔案 {file_path} 時出錯: {e}")
            self.error_count += 1
            self.error_messages.append(str(e))
            return None

    def _incremental_update_for_root(self, conn, cursor, root_dir):
        disk_records = {}
        for root, dirs, files in os.walk(root_dir, topdown=True):
            if not self.is_running:
                return
            if self._should_exclude(root):
                dirs[:] = []
                continue
            for fname in files:
                if not self.is_running:
                    return
                fpath = os.path.join(root, fname)
                rec = self._build_file_record(fpath, fname)
                if rec:
                    disk_records[fpath] = rec

        cursor.execute(
            "SELECT path, size, modified FROM files WHERE path LIKE ?",
            (f"{root_dir}%",),
        )
        db_rows = cursor.fetchall()
        db_index = {p: (sz, mod) for p, sz, mod in db_rows}

        def _norm_mod(v):
            if isinstance(v, str):
                return v
            try:
                return datetime.datetime.fromtimestamp(v).strftime("%Y-%m-%d %H:%M:%S")
            except Exception:
                return str(v)

        to_insert, to_update = [], []
        processed = 0
        total = max(len(disk_records), 1)

        for path, rec in disk_records.items():
            processed += 1
            if processed % 200 == 0:
                self.progress_update.emit(processed, total)
            if path not in db_index:
                to_insert.append(rec)
            else:
                old_size, old_mod = db_index[path]
                if (old_size != rec["size"]) or (
                    _norm_mod(old_mod) != _norm_mod(rec["modified"])
                ):
                    to_update.append(rec)

        to_delete = [path for path in db_index.keys() if path not in disk_records]

        if to_insert:
            self._batch_insert_files(conn, cursor, to_insert)
        if to_update:
            self._batch_update_files(conn, cursor, to_update)
        if to_delete:
            self._batch_delete_paths(conn, cursor, to_delete)

        self.progress_update.emit(total, total)

    def _batch_insert_files(self, conn, cursor, files_info):
        total = len(files_info)
        if total == 0:
            return
        batch_size = self.batch_size
        inserted_count = 0
        for i in range(0, total, batch_size):
            if not self.is_running:
                return
            batch = files_info[i : i + batch_size]
            insert_data = []
            for file_info in batch:
                insert_data.append(
                    (
                        file_info["path"],
                        file_info["filename"],
                        file_info["ext"],
                        file_info["size"],
                        file_info["modified"],
                        file_info["content_indexed"],
                        file_info["content"],
                    )
                )
            cursor.executemany(
                "INSERT OR REPLACE INTO files (path, filename, ext, size, modified, content_indexed, content) VALUES (?, ?, ?, ?, ?, ?, ?)",
                insert_data,
            )
            inserted_count += len(batch)
            conn.commit()
        self.files_added += inserted_count

    def _batch_update_files(self, conn, cursor, files_info):
        total = len(files_info)
        if total == 0:
            return
        batch_size = self.batch_size
        updated_count = 0
        for i in range(0, total, batch_size):
            if not self.is_running:
                return
            batch = files_info[i : i + batch_size]
            update_data = []
            for file_info in batch:
                update_data.append(
                    (
                        file_info["filename"],
                        file_info["ext"],
                        file_info["size"],
                        file_info["modified"],
                        file_info["content_indexed"],
                        file_info["content"],
                        datetime.datetime.now(),  # last_verified
                        file_info["path"],
                    )
                )
            cursor.executemany(
                "UPDATE files SET filename=?, ext=?, size=?, modified=?, content_indexed=?, content=?, last_verified=? WHERE path=?",
                update_data,
            )
            updated_count += len(batch)
            conn.commit()
        self.files_updated += updated_count

    def _batch_delete_paths(self, conn, cursor, paths):
        total = len(paths)
        if total == 0:
            return
        batch_size = self.batch_size
        deleted_count = 0
        for i in range(0, total, batch_size):
            if not self.is_running:
                return
            batch = paths[i : i + batch_size]
            cursor.executemany("DELETE FROM files WHERE path=?", [(p,) for p in batch])
            deleted_count += len(batch)
            conn.commit()
        self.files_deleted += deleted_count

    def _should_exclude(self, path, file_stat=None):
        for exclude_pattern in self.exclude_dirs:
            try:
                if re.search(exclude_pattern, path):
                    return True
            except re.error:
                continue

        if self.advanced_exclude_rules:
            if self.advanced_exclude_rules.should_exclude(path, file_stat):
                return True

        return False

    def stop(self):
        self.is_running = False


class PaginationWidget(QWidget):
    """分頁控制元件"""

    page_changed = pyqtSignal(int)

    def __init__(self, parent=None):
        super().__init__(parent)
        self._current_page = 1
        self._total_pages = 1
        self._total_count = 0
        self._page_size = 100
        self._init_ui()

    def _init_ui(self):
        layout = QHBoxLayout(self)
        layout.setContentsMargins(4, 4, 4, 4)
        layout.setSpacing(8)

        self.page_size_combo = QComboBox()
        self.page_size_combo.addItems(["50", "100", "200", "500"])
        self.page_size_combo.setCurrentText("100")
        self.page_size_combo.setMaximumWidth(70)
        self.page_size_combo.currentTextChanged.connect(self._on_page_size_changed)
        layout.addWidget(QLabel("每頁:"))
        layout.addWidget(self.page_size_combo)

        layout.addSpacing(16)

        self.first_btn = QPushButton(" << ")
        self.first_btn.setMaximumWidth(40)
        self.first_btn.clicked.connect(self._go_first)
        layout.addWidget(self.first_btn)

        self.prev_btn = QPushButton(" < ")
        self.prev_btn.setMaximumWidth(40)
        self.prev_btn.clicked.connect(self._go_prev)
        layout.addWidget(self.prev_btn)

        self.page_input = QSpinBox()
        self.page_input.setMinimum(1)
        self.page_input.setMaximum(999999)
        self.page_input.setValue(1)
        self.page_input.setMaximumWidth(60)
        self.page_input.valueChanged.connect(self._on_page_input_changed)
        layout.addWidget(self.page_input)

        self.page_info = QLabel(" / 1 ")
        layout.addWidget(self.page_info)

        self.next_btn = QPushButton(" > ")
        self.next_btn.setMaximumWidth(40)
        self.next_btn.clicked.connect(self._go_next)
        layout.addWidget(self.next_btn)

        self.last_btn = QPushButton(" >> ")
        self.last_btn.setMaximumWidth(40)
        self.last_btn.clicked.connect(self._go_last)
        layout.addWidget(self.last_btn)

        layout.addSpacing(16)

        self.count_label = QLabel(" 共 0 項 ")
        layout.addWidget(self.count_label)

        layout.addStretch()

        self._update_buttons()

    def _on_page_size_changed(self, text):
        self._page_size = int(text)
        self._total_pages = max(1, (self._total_count + self._page_size - 1) // self._page_size)
        self._current_page = 1
        self.page_input.setMaximum(self._total_pages)
        self.page_input.setValue(1)
        self._update_display()
        self.page_changed.emit(self._current_page)

    def _on_page_input_changed(self, value):
        if value != self._current_page:
            self._current_page = value
            self._update_display()
            self.page_changed.emit(self._current_page)

    def _go_first(self):
        if self._current_page != 1:
            self._current_page = 1
            self.page_input.setValue(1)
            self.page_changed.emit(self._current_page)

    def _go_prev(self):
        if self._current_page > 1:
            self._current_page -= 1
            self.page_input.setValue(self._current_page)
            self.page_changed.emit(self._current_page)

    def _go_next(self):
        if self._current_page < self._total_pages:
            self._current_page += 1
            self.page_input.setValue(self._current_page)
            self.page_changed.emit(self._current_page)

    def _go_last(self):
        if self._current_page != self._total_pages:
            self._current_page = self._total_pages
            self.page_input.setValue(self._total_pages)
            self.page_changed.emit(self._current_page)

    def _update_buttons(self):
        self.first_btn.setEnabled(self._current_page > 1)
        self.prev_btn.setEnabled(self._current_page > 1)
        self.next_btn.setEnabled(self._current_page < self._total_pages)
        self.last_btn.setEnabled(self._current_page < self._total_pages)

    def _update_display(self):
        self.page_info.setText(f" / {self._total_pages} ")
        self._update_buttons()

    def set_total_count(self, total_count):
        self._total_count = total_count
        self._total_pages = max(1, (total_count + self._page_size - 1) // self._page_size)
        self.count_label.setText(f" 共 {total_count} 項 ")
        self.page_input.setMaximum(self._total_pages)
        if self._current_page > self._total_pages:
            self._current_page = self._total_pages
            self.page_input.setValue(self._total_pages)
        self._update_display()

    def current_page(self):
        return self._current_page

    def page_size(self):
        return self._page_size

    def reset(self):
        self._current_page = 1
        self._total_pages = 1
        self._total_count = 0
        self.page_input.setValue(1)
        self.page_input.setMaximum(999999)
        self._update_display()


class FileSearcher(QThread):
    search_complete = pyqtSignal(list, float)
    search_progress = pyqtSignal(int, int)
    search_count = pyqtSignal(int)

    def __init__(self, db_path, search_config, max_workers=4, max_memory_mb=1024, page=1, page_size=100):
        super().__init__()
        self.db_path = db_path
        self.search_config = search_config
        self.is_running = True
        self.max_workers = max_workers
        self.max_memory_mb = max_memory_mb
        self.page = page
        self.page_size = page_size
        self.all_results = []
        print(
            f"搜尋執行緒數設定為: {self.max_workers}, 最大記憶體: {self.max_memory_mb}MB, 分頁: {page}/{page_size}"
        )

    def run(self):
        db_dir = os.path.dirname(self.db_path)
        if not os.path.exists(db_dir):
            try:
                os.makedirs(db_dir, exist_ok=True)
            except Exception as e:
                print(f"無法創建資料庫目錄: {e}")
                self.search_complete.emit([], 0)
                return

        try:
            conn = sqlite3.connect(self.db_path)
            conn.execute("CREATE TABLE IF NOT EXISTS test_table (id INTEGER)")
            conn.execute("DROP TABLE IF EXISTS test_table")
            conn.commit()
            conn.close()
        except Exception as e:
            print(f"資料庫無法寫入: {e}")
            self.search_complete.emit([], 0)
            return

        start_time = time.time()
        conn = sqlite3.connect(self.db_path)
        try:
            conn.create_function("REGEXP", 2, self._regexp)
            cursor = conn.cursor()

            conditions = []
            params = []
            config = self.search_config

            if config.get("filename"):
                if config.get("use_regex_filename", False):
                    conditions.append("filename REGEXP ?")
                    params.append(config["filename"])
                elif config.get("use_fuzzy_filename", False):
                    conditions.append("filename LIKE ?")
                    params.append(f"%{config['filename']}%")
                else:
                    conditions.append("filename LIKE ?")
                    params.append(f"%{config['filename']}%")

            if config.get("content"):
                if config.get("use_regex_content", False):
                    conditions.append("content REGEXP ?")
                    params.append(config["content"])
                elif config.get("use_fuzzy_content", False):
                    conditions.append("content LIKE ?")
                    params.append(f"%{config['content']}%")
                else:
                    conditions.append("content LIKE ?")
                    params.append(f"%{config['content']}%")

            if config.get("file_type"):
                file_types = [
                    ft.strip() for ft in config["file_type"].split(",") if ft.strip()
                ]
                if file_types:
                    type_conditions = []
                    for ft in file_types:
                        ft_norm = ft.lower()
                        if not ft_norm.startswith("."):
                            ft_norm = "." + ft_norm
                        type_conditions.append("ext = ?")
                        params.append(ft_norm)
                    conditions.append(f"({' OR '.join(type_conditions)})")

            if config.get("min_size") is not None:
                conditions.append("size >= ?")
                params.append(config["min_size"])

            if config.get("max_size") is not None:
                conditions.append("size <= ?")
                params.append(config["max_size"])

            if config.get("date_after"):
                conditions.append("modified >= ?")
                params.append(config["date_after"])

            if config.get("date_before"):
                conditions.append("modified <= ?")
                params.append(config["date_before"])

            if config.get("path_filter"):
                conditions.append("path LIKE ?")
                params.append(f"%{config['path_filter']}%")

            if config.get("boolean_mode"):
                self._handle_boolean_search(config, conditions, params)

            if not conditions:
                self.search_count.emit(0)
                self.search_complete.emit([], 0)
                return

            where_clause = ' AND '.join(conditions)
            count_query = f"SELECT COUNT(*) FROM files WHERE {where_clause}"

            try:
                cursor.execute(count_query, params)
                total_rows = cursor.fetchone()[0]
                self.search_count.emit(total_rows)
            except Exception as e:
                print(f"計算結果總數時出錯: {e}")
                total_rows = 0
                self.search_count.emit(0)

            if total_rows == 0:
                self.search_complete.emit([], 0)
                return

            sort_field = config.get("sort_by", "filename")
            sort_order = "DESC" if config.get("sort_desc", False) else "ASC"

            offset = (self.page - 1) * self.page_size
            query = f"SELECT path, filename, ext, size, modified, content FROM files WHERE {where_clause} ORDER BY {sort_field} {sort_order} LIMIT {self.page_size} OFFSET {offset}"

            try:
                cursor.execute(query, params)
            except Exception as e:
                print(f"執行搜尋查詢時出錯: {e}")
                self.search_complete.emit([], 0)
                return

            results = []
            rows = cursor.fetchall()

            for row in rows:
                if not self.is_running:
                    break

                result = {
                    "path": row[0],
                    "filename": row[1],
                    "ext": row[2],
                    "size": row[3],
                    "modified": row[4],
                    "content": row[5],
                }

                if config.get("use_fuzzy_filename", False) and config.get("filename"):
                    if not fuzzy_match(result["filename"], config["filename"], threshold=0.6):
                        continue

                if config.get("use_fuzzy_content", False) and config.get("content"):
                    if result["content"] and not fuzzy_match(result["content"], config["content"], threshold=0.6):
                        continue

                results.append(result)

            search_time = time.time() - start_time
            print(f"搜尋完成，第 {self.page} 頁，獲得 {len(results)} 個結果，耗時 {search_time:.2f} 秒")
            self.search_complete.emit(results, search_time)
        finally:
            conn.close()

    def _regexp(self, pattern, item):
        if item is None:
            return False
        try:
            return re.search(pattern, item, re.MULTILINE) is not None
        except re.error:
            return False

    def _handle_boolean_search(self, config, conditions, params):
        if not config.get("boolean_expr"):
            return
        expr = config["boolean_expr"].strip()
        if not expr:
            return
        
        if '(' in expr or ')' in expr:
            try:
                parser = BooleanExpressionParser(expr)
                condition, extra_params = parser.parse()
                if condition:
                    conditions.append(condition)
                    params.extend(extra_params)
                    return
            except Exception as e:
                print(f"布林表達式解析錯誤: {e}")
        
        if " AND " in expr:
            terms = expr.split(" AND ")
            and_conditions = []
            for term in terms:
                term = term.strip()
                if term:
                    if term.startswith("NOT "):
                        not_term = term[4:].strip()
                        and_conditions.append(
                            "(filename NOT LIKE ? AND content NOT LIKE ?)"
                        )
                        params.append(f"%{not_term}%")
                        params.append(f"%{not_term}%")
                    else:
                        and_conditions.append("(filename LIKE ? OR content LIKE ?)")
                        params.append(f"%{term}%")
                        params.append(f"%{term}%")
            if and_conditions:
                conditions.append(f"({' AND '.join(and_conditions)})")
        elif " OR " in expr:
            terms = expr.split(" OR ")
            or_conditions = []
            for term in terms:
                term = term.strip()
                if term:
                    if term.startswith("NOT "):
                        not_term = term[4:].strip()
                        or_conditions.append(
                            "(filename NOT LIKE ? AND content NOT LIKE ?)"
                        )
                        params.append(f"%{not_term}%")
                        params.append(f"%{not_term}%")
                    else:
                        or_conditions.append("(filename LIKE ? OR content LIKE ?)")
                        params.append(f"%{term}%")
                        params.append(f"%{term}%")
            if or_conditions:
                conditions.append(f"({' OR '.join(or_conditions)})")
        elif " NEAR " in expr:
            terms = expr.split(" NEAR ")
            if len(terms) == 2:
                term1, term2 = terms[0].strip(), terms[1].strip()
                conditions.append("(content LIKE ? AND content LIKE ?)")
                params.append(f"%{term1}%")
                params.append(f"%{term2}%")
        elif expr.startswith("REGEX "):
            regex_pattern = expr[6:].strip()
            conditions.append("(filename REGEXP ? OR content REGEXP ?)")
            params.append(regex_pattern)
            params.append(regex_pattern)
        elif expr.startswith("LIKE "):
            like_pattern = expr[5:].strip()
            conditions.append("(filename LIKE ? OR content LIKE ?)")
            params.append(like_pattern)
            params.append(like_pattern)
        elif expr.startswith("NOT "):
            not_term = expr[4:].strip()
            conditions.append("(filename NOT LIKE ? AND content NOT LIKE ?)")
            params.append(f"%{not_term}%")
            params.append(f"%{not_term}%")
        else:
            conditions.append("(filename LIKE ? OR content LIKE ?)")
            params.append(f"%{expr}%")
            params.append(f"%{expr}%")

    def stop(self):
        self.is_running = False


class I18nManager:
    """多語言支援管理器"""

    DEFAULT_LANGUAGE = "zh-TW"
    SUPPORTED_LANGUAGES = ["zh-TW", "zh-CN", "en"]

    def __init__(self):
        self.current_language = self.DEFAULT_LANGUAGE
        self.translations = {}
        self._get_i18n_dir()
        self.load_language(self.DEFAULT_LANGUAGE)

    def _get_i18n_dir(self):
        """取得 i18n 資料夾路徑"""
        script_dir = os.path.dirname(os.path.abspath(__file__))
        self.i18n_dir = os.path.join(script_dir, "i18n")
        if not os.path.exists(self.i18n_dir):
            self.i18n_dir = None

    def load_language(self, language_code):
        """載入指定語言的翻譯"""
        if language_code not in self.SUPPORTED_LANGUAGES:
            language_code = self.DEFAULT_LANGUAGE
        self.current_language = language_code
        if not self.i18n_dir:
            self.translations = {}
            return False
        file_path = os.path.join(self.i18n_dir, f"{language_code}.json")
        try:
            if os.path.exists(file_path):
                with open(file_path, "r", encoding="utf-8") as f:
                    self.translations = json.load(f)
                return True
            else:
                self.translations = {}
                return False
        except Exception as e:
            print(f"載入語言檔案失敗: {e}")
            self.translations = {}
            return False

    def get(self, key, default=None):
        """取得翻譯文字"""
        keys = key.split(".")
        value = self.translations
        for k in keys:
            if isinstance(value, dict) and k in value:
                value = value[k]
            else:
                return default if default else key
        return str(value) if value else (default if default else key)

    def get_language_name(self, language_code):
        """取得語言顯示名稱"""
        names = {
            "zh-TW": self.get("settings.language_zh_tw", "繁體中文"),
            "zh-CN": self.get("settings.language_zh_cn", "簡體中文"),
            "en": self.get("settings.language_en", "英文"),
        }
        return names.get(language_code, language_code)


i18n = I18nManager()


def _tr(key, default=None):
    """取得翻譯文字的快捷函數"""
    return i18n.get(key, default)


class StatisticsWidget(QWidget):
    """統計圖表 Widget"""

    def __init__(self, db_path, parent=None):
        super().__init__(parent)
        self.db_path = db_path
        self.dark_mode = False
        self._init_ui()
        self.refresh_statistics()

    def _init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(16, 16, 16, 16)
        layout.setSpacing(16)

        header_layout = QHBoxLayout()
        self.title_label = QLabel(_tr("statistics.title"))
        self.title_label.setStyleSheet("font-size: 18px; font-weight: bold;")
        header_layout.addWidget(self.title_label)
        header_layout.addStretch()

        self.refresh_btn = QPushButton(_tr("statistics.refresh"))
        self.refresh_btn.setFixedWidth(100)
        self.refresh_btn.clicked.connect(self.refresh_statistics)
        header_layout.addWidget(self.refresh_btn)
        layout.addLayout(header_layout)

        self.summary_label = QLabel()
        self.summary_label.setStyleSheet("font-size: 12px; color: #666;")
        layout.addWidget(self.summary_label)

        charts_layout = QHBoxLayout()
        charts_layout.setSpacing(16)

        self.type_chart_container = self._create_chart_container(_tr("statistics.type_chart"))
        charts_layout.addWidget(self.type_chart_container)

        self.size_chart_container = self._create_chart_container(_tr("statistics.size_chart"))
        charts_layout.addWidget(self.size_chart_container)

        layout.addLayout(charts_layout)

        charts_layout2 = QHBoxLayout()
        charts_layout2.setSpacing(16)

        self.time_chart_container = self._create_chart_container(_tr("statistics.time_chart"))
        charts_layout2.addWidget(self.time_chart_container)

        self.dir_chart_container = self._create_chart_container(_tr("statistics.dir_chart"))
        charts_layout2.addWidget(self.dir_chart_container)

        layout.addLayout(charts_layout2)
        layout.addStretch()

    def _create_chart_container(self, title):
        container = QGroupBox(title)
        container.setStyleSheet("""
            QGroupBox {
                font-weight: bold;
                border: 1px solid #ddd;
                border-radius: 8px;
                margin-top: 12px;
                padding-top: 8px;
            }
            QGroupBox::title {
                subcontrol-origin: margin;
                left: 12px;
                padding: 0 6px;
            }
        """)
        container_layout = QVBoxLayout(container)
        container_layout.setContentsMargins(8, 16, 8, 8)

        if matplotlib and FigureCanvas:
            figure = Figure(figsize=(5, 4), dpi=80)
            canvas = FigureCanvas(figure)
            canvas.setMinimumHeight(250)
            container_layout.addWidget(canvas)
            container.figure = figure
            container.canvas = canvas
        else:
            error_label = QLabel("請安裝 matplotlib:\npip install matplotlib")
            error_label.setAlignment(Qt.AlignCenter)
            error_label.setStyleSheet("color: #999;")
            container_layout.addWidget(error_label)
            container.figure = None
            container.canvas = None

        return container

    def set_dark_mode(self, dark_mode):
        self.dark_mode = dark_mode
        self.refresh_statistics()

    def refresh_statistics(self):
        if not matplotlib or not FigureCanvas:
            self.summary_label.setText("統計功能需要安裝 matplotlib 套件")
            return

        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute("SELECT COUNT(*) FROM files")
            total_files = cursor.fetchone()[0]

            cursor.execute("SELECT COALESCE(SUM(size), 0) FROM files")
            total_size = cursor.fetchone()[0]

            if total_files == 0:
                self.summary_label.setText("目前沒有索引檔案，請先建立索引")
                self._clear_charts()
                conn.close()
                return

            size_str = self._format_size(total_size)
            self.summary_label.setText(f"總檔案數: {total_files:,} | 總大小: {size_str}")

            self._draw_type_chart(cursor)
            self._draw_size_chart(cursor)
            self._draw_time_chart(cursor)
            self._draw_dir_chart(cursor)

            conn.close()

        except Exception as e:
            self.summary_label.setText(f"載入統計資料失敗: {str(e)}")

    def _clear_charts(self):
        for container in [self.type_chart_container, self.size_chart_container,
                          self.time_chart_container, self.dir_chart_container]:
            if container.figure:
                container.figure.clear()
                container.canvas.draw()

    def _draw_type_chart(self, cursor):
        cursor.execute("""
            SELECT ext, COUNT(*) as cnt
            FROM files
            WHERE ext IS NOT NULL AND ext != ''
            GROUP BY ext
            ORDER BY cnt DESC
            LIMIT 10
        """)
        data = cursor.fetchall()

        if not data:
            return

        fig = self.type_chart_container.figure
        fig.clear()
        ax = fig.add_subplot(111)

        labels = [row[0] if row[0] else "(無副檔名)" for row in data]
        sizes = [row[1] for row in data]

        if self.dark_mode:
            ax.set_facecolor('#1e1e1e')
            fig.patch.set_facecolor('#1e1e1e')

        colors = plt.cm.Set3(range(len(labels)))
        wedges, texts, autotexts = ax.pie(
            sizes, labels=labels, autopct='%1.1f%%',
            colors=colors, pctdistance=0.75,
            wedgeprops=dict(width=0.5)
        )

        for text in texts + autotexts:
            text.set_fontsize(9)
            if self.dark_mode:
                text.set_color('white')

        ax.set_title("")
        fig.tight_layout()
        self.type_chart_container.canvas.draw()

    def _draw_size_chart(self, cursor):
        cursor.execute("""
            SELECT size FROM files
            WHERE size IS NOT NULL
            ORDER BY size
        """)
        sizes = [row[0] for row in cursor.fetchall()]

        if not sizes:
            return

        fig = self.size_chart_container.figure
        fig.clear()
        ax = fig.add_subplot(111)

        if self.dark_mode:
            ax.set_facecolor('#1e1e1e')
            fig.patch.set_facecolor('#1e1e1e')

        size_mb = [s / (1024 * 1024) for s in sizes]

        max_size = max(size_mb) if size_mb else 1
        if max_size > 1000:
            bins = [0, 1, 10, 100, 1000, max_size + 1]
            labels = ['<1MB', '1-10MB', '10-100MB', '100MB-1GB', '>1GB']
        else:
            bins = [0, 0.1, 1, 10, 100, max_size + 1]
            labels = ['<100KB', '100KB-1MB', '1-10MB', '10-100MB', '>100MB']

        counts = []
        for i in range(len(bins) - 1):
            cnt = sum(1 for s in size_mb if bins[i] <= s < bins[i + 1])
            counts.append(cnt)

        colors = ['#3498db', '#2ecc71', '#f39c12', '#e74c3c', '#9b59b6']
        bars = ax.bar(labels[:len(counts)], counts, color=colors[:len(counts)])

        for bar in bars:
            height = bar.get_height()
            if height > 0:
                ax.annotate(f'{int(height)}',
                           xy=(bar.get_x() + bar.get_width() / 2, height),
                           ha='center', va='bottom', fontsize=8)

        ax.set_ylabel('檔案數量', fontsize=9)
        if self.dark_mode:
            ax.tick_params(colors='white')
            ax.xaxis.label.set_color('white')
            ax.yaxis.label.set_color('white')
            for spine in ax.spines.values():
                spine.set_color('white')
        ax.tick_params(axis='x', labelsize=8, rotation=15)
        fig.tight_layout()
        self.size_chart_container.canvas.draw()

    def _draw_time_chart(self, cursor):
        cursor.execute("""
            SELECT date(modified) as d, COUNT(*) as cnt
            FROM files
            WHERE modified IS NOT NULL
            GROUP BY d
            ORDER BY d DESC
            LIMIT 30
        """)
        data = cursor.fetchall()

        if len(data) < 2:
            return

        fig = self.time_chart_container.figure
        fig.clear()
        ax = fig.add_subplot(111)

        if self.dark_mode:
            ax.set_facecolor('#1e1e1e')
            fig.patch.set_facecolor('#1e1e1e')

        data.reverse()
        dates = [row[0] for row in data]
        counts = [row[1] for row in data]

        display_dates = []
        for i, d in enumerate(dates):
            if i % 5 == 0:
                display_dates.append(d[5:] if len(d) > 5 else d)
            else:
                display_dates.append('')

        ax.plot(range(len(dates)), counts, marker='o', markersize=4,
               color='#3498db', linewidth=1.5)
        ax.fill_between(range(len(dates)), counts, alpha=0.3, color='#3498db')

        ax.set_xticks(range(len(dates)))
        ax.set_xticklabels(display_dates, rotation=45, ha='right', fontsize=7)
        ax.set_ylabel('檔案數量', fontsize=9)

        if self.dark_mode:
            ax.tick_params(colors='white')
            ax.xaxis.label.set_color('white')
            ax.yaxis.label.set_color('white')
            for spine in ax.spines.values():
                spine.set_color('white')

        fig.tight_layout()
        self.time_chart_container.canvas.draw()

    def _draw_dir_chart(self, cursor):
        cursor.execute("""
            SELECT
                CASE
                    WHEN instr(path, '/') > 0 THEN substr(path, 1, instr(path, '/') - 1)
                    ELSE path
                END as root_path,
                COUNT(*) as cnt
            FROM files
            WHERE path IS NOT NULL
            GROUP BY root_path
            ORDER BY cnt DESC
            LIMIT 10
        """)
        data = cursor.fetchall()

        if not data:
            return

        fig = self.dir_chart_container.figure
        fig.clear()
        ax = fig.add_subplot(111)

        if self.dark_mode:
            ax.set_facecolor('#1e1e1e')
            fig.patch.set_facecolor('#1e1e1e')

        paths = []
        for row in data:
            path = row[0]
            if len(path) > 25:
                path = path[:22] + '...'
            paths.append(path)

        counts = [row[1] for row in data]

        y_pos = range(len(paths))
        bars = ax.barh(y_pos, counts, color='#2ecc71')

        ax.set_yticks(y_pos)
        ax.set_yticklabels(paths, fontsize=8)
        ax.set_xlabel('檔案數量', fontsize=9)
        ax.invert_yaxis()

        for i, bar in enumerate(bars):
            width = bar.get_width()
            ax.annotate(f'{int(width)}',
                       xy=(width, bar.get_y() + bar.get_height() / 2),
                       ha='left', va='center', fontsize=8,
                       xytext=(3, 0), textcoords='offset points')

        if self.dark_mode:
            ax.tick_params(colors='white')
            ax.xaxis.label.set_color('white')
            ax.yaxis.label.set_color('white')
            for spine in ax.spines.values():
                spine.set_color('white')

        fig.tight_layout()
        self.dir_chart_container.canvas.draw()

    def _format_size(self, size_bytes):
        for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
            if size_bytes < 1024.0:
                return f"{size_bytes:.2f} {unit}"
            size_bytes /= 1024.0
        return f"{size_bytes:.2f} PB"


class SearchSuggestionWidget(QWidget):
    """搜尋建議下拉列表"""

    suggestion_selected = pyqtSignal(str)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowFlags(Qt.Popup | Qt.FramelessWindowHint)
        self.setAttribute(Qt.WA_TranslucentBackground, False)
        self.max_visible_items = 10
        self.item_height = 25
        self._init_ui()

    def _init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)
        self.list_widget = QListWidget()
        self.list_widget.setStyleSheet("""
            QListWidget {
                background-color: #ffffff;
                border: 1px solid #ccc;
                border-radius: 4px;
            }
            QListWidget::item {
                padding: 6px 12px;
                border-bottom: 1px solid #eee;
            }
            QListWidget::item:selected {
                background-color: #e3f2fd;
                color: #333;
            }
            QListWidget::item:hover {
                background-color: #f5f5f5;
            }
        """)
        self.list_widget.itemClicked.connect(self._on_item_clicked)
        layout.addWidget(self.list_widget)

    def show_suggestions(self, suggestions, position):
        self.list_widget.clear()
        if not suggestions:
            self.hide()
            return
        for suggestion in suggestions:
            item = QListWidgetItem(suggestion)
            self.list_widget.addItem(item)
        visible_count = min(len(suggestions), self.max_visible_items)
        self.setFixedHeight(visible_count * self.item_height + 2)
        self.setFixedWidth(300)
        if isinstance(position, tuple):
            self.move(position[0], position[1])
        else:
            self.move(position)
        self.show()
        self.list_widget.setCurrentRow(0)
        self.list_widget.setFocus()

    def _on_item_clicked(self, item):
        self.suggestion_selected.emit(item.text())
        self.hide()

    def keyPressEvent(self, event):
        if event.key() == Qt.Key_Escape:
            self.hide()
        elif event.key() == Qt.Key_Enter or event.key() == Qt.Key_Return:
            current_item = self.list_widget.currentItem()
            if current_item:
                self._on_item_clicked(current_item)
        else:
            super().keyPressEvent(event)

    def select_next(self):
        current_row = self.list_widget.currentRow()
        if current_row < self.list_widget.count() - 1:
            self.list_widget.setCurrentRow(current_row + 1)

    def select_previous(self):
        current_row = self.list_widget.currentRow()
        if current_row > 0:
            self.list_widget.setCurrentRow(current_row - 1)


class AdvancedFileSearcher(QMainWindow):
    def __init__(self):
        super().__init__()
        self.data_dir = self._get_app_data_dir()
        self.settings_path = os.path.join(self.data_dir, SETTINGS_FILENAME)

        self.db_path = self._get_safe_db_path()
        self.search_directories = []
        self.exclude_patterns = []

        self.dark_mode = False
        self._load_theme_preference()
        self._load_language_preference()

        self.update_history_manager = UpdateHistoryManager(self.db_path)
        self.search_history_manager = SearchHistoryManager(self.db_path)
        self.search_template_manager = SearchTemplateManager(self.db_path)
        self.advanced_exclude_rules = AdvancedExcludeRules()

        self.tray_icon = None
        self.close_to_tray_no_ask = False

        self.realtime_search_timer = QTimer()
        self.realtime_search_timer.setSingleShot(True)
        self.realtime_search_timer.timeout.connect(self._do_realtime_search)
        self.realtime_search_enabled = False

        # 設定應用程式圖示
        self._set_app_icon()

        self.init_ui()
        self._update_ui_language()
        self.load_settings_into_ui()
        self._init_tray_icon()
        self._init_auto_update_timer()
        self._init_search_history()
        self._init_realtime_search()
        self._init_file_watcher()
        self.setAcceptDrops(True)
        self._drag_overlay_active = False

    def _set_app_icon(self):
        """設定應用程式圖示"""
        try:
            # 取得圖示檔案路徑
            icon_path = self._get_icon_path()
            if icon_path and os.path.exists(icon_path):
                app_icon = QIcon(icon_path)
                if not app_icon.isNull():
                    self.setWindowIcon(app_icon)
                    # 同時設定應用程式層級的圖示
                    QApplication.instance().setWindowIcon(app_icon)
                else:
                    print(f"圖示檔案無效: {icon_path}")
            else:
                print(f"圖示檔案不存在: {icon_path}")
        except Exception as e:
            print(f"載入圖示時發生錯誤: {e}")

    def _get_icon_path(self):
        """取得圖示檔案路徑（跨平台支援 .ico / .icns）"""
        import platform

        # 根據平台決定優先嘗試的副檔名
        if platform.system() == "Darwin":
            exts = [".icns", ".ico", ".png"]
        else:
            exts = [".ico", ".png", ".icns"]

        # 方法1: 打包後的資源路徑 (PyInstaller)
        if hasattr(sys, "_MEIPASS"):
            for ext in exts:
                icon_path = os.path.join(sys._MEIPASS, "app_icon" + ext)
                if os.path.exists(icon_path):
                    return icon_path

        # 方法2: 相對於執行檔的路徑
        script_dir = os.path.dirname(os.path.abspath(__file__))
        for ext in exts:
            icon_path = os.path.join(script_dir, "app_icon" + ext)
            if os.path.exists(icon_path):
                return icon_path

        # 方法3: 相對於當前工作目錄
        for ext in exts:
            icon_path = os.path.join(os.getcwd(), "app_icon" + ext)
            if os.path.exists(icon_path):
                return icon_path

        # 方法4: 在應用程式資料目錄中尋找
        for ext in exts:
            icon_path = os.path.join(self.data_dir, "app_icon" + ext)
            if os.path.exists(icon_path):
                return icon_path

        return None

    def _get_app_data_dir(self):
        """取得應用程式資料目錄"""
        import platform

        if platform.system() == "Windows":
            app_data = os.environ.get("LOCALAPPDATA", os.path.expanduser("~"))
            data_dir = os.path.join(app_data, APP_NAME)
        elif platform.system() == "Darwin":
            data_dir = os.path.join(
                os.path.expanduser("~"), "Library", "Application Support", APP_NAME
            )
        else:
            data_dir = os.path.join(
                os.path.expanduser("~"), ".local", "share", APP_NAME
            )
        os.makedirs(data_dir, exist_ok=True)
        return data_dir

    def _get_safe_db_path(self):
        try:
            with open(
                os.path.join(self._get_app_data_dir(), SETTINGS_FILENAME),
                "r",
                encoding="utf-8",
            ) as f:
                s = json.load(f)
                if s.get("db_path"):
                    return s["db_path"]
        except Exception:
            pass
        return os.path.join(self._get_app_data_dir(), "file_index.db")

    def init_ui(self):
        self.setWindowTitle("SearchingPro - 進階檔案搜尋工具")
        self.setGeometry(100, 100, 1280, 800)

        self._init_menu_bar()
        self._apply_theme()

        main_layout = QVBoxLayout()
        main_layout.setContentsMargins(12, 12, 12, 12)
        main_layout.setSpacing(8)

        self.tabs = QTabWidget()
        self.search_tab = QWidget()
        self.index_tab = QWidget()
        self.settings_tab = QWidget()
        self.statistics_tab = QWidget()

        self.tabs.addTab(self.search_tab, " 搜尋 ")
        self.tabs.addTab(self.index_tab, " 索引管理 ")
        self.tabs.addTab(self.statistics_tab, " 統計 ")
        self.tabs.addTab(self.settings_tab, " 設定 ")
        self.tabs.tabBar().setExpanding(True)

        self._init_search_tab()
        self._init_index_tab()
        self._init_statistics_tab()
        self._init_settings_tab()

        main_layout.addWidget(self.tabs)

        self.statusBar = QStatusBar()
        self.setStatusBar(self.statusBar)
        self.statusBar.showMessage(" 就緒")

        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        self.progress_bar.setMaximumWidth(200)
        self.statusBar.addPermanentWidget(self.progress_bar)

        self.theme_status_label = QLabel(" 淺色模式 ")
        self.theme_status_label.setStyleSheet("color: #8e8e93; font-size: 11px;")
        self.statusBar.addPermanentWidget(self.theme_status_label)

        central_widget = QWidget()
        central_widget.setLayout(main_layout)
        self.setCentralWidget(central_widget)

    def _init_menu_bar(self):
        """初始化選單列"""
        menubar = self.menuBar()

        file_menu = menubar.addMenu("檔案")

        export_action = file_menu.addAction("匯出搜尋結果為 CSV...")
        export_action.setShortcut(_platform_shortcut("Cmd+E"))
        export_action.triggered.connect(self.export_results_csv)

        import_action = file_menu.addAction("匯入搜尋結果...")
        import_action.setShortcut(_platform_shortcut("Cmd+I"))
        import_action.triggered.connect(self._import_search_results)

        file_menu.addSeparator()

        exit_action = file_menu.addAction("離開")
        exit_action.setShortcut(_platform_shortcut("Cmd+Q"))
        exit_action.triggered.connect(self.close)

        edit_menu = menubar.addMenu("編輯")

        clear_search_action = edit_menu.addAction("清除搜尋條件")
        clear_search_action.setShortcut(_platform_shortcut("Cmd+Shift+C"))
        clear_search_action.triggered.connect(self._clear_search_fields)

        edit_menu.addSeparator()

        toggle_preview_action = edit_menu.addAction("切換檔案預覽")
        toggle_preview_action.setShortcut("F3")
        toggle_preview_action.triggered.connect(self._toggle_preview)

        view_menu = menubar.addMenu("檢視")

        toggle_theme_action = view_menu.addAction("切換深色/淺色模式")
        toggle_theme_action.setShortcut(_platform_shortcut("Cmd+T"))
        toggle_theme_action.triggered.connect(self._toggle_theme)

        view_menu.addSeparator()

        toggle_realtime_action = view_menu.addAction("切換即時搜尋")
        toggle_realtime_action.setShortcut(_platform_shortcut("Cmd+R"))
        toggle_realtime_action.triggered.connect(self._toggle_realtime_search)

        help_menu = menubar.addMenu("說明")

        shortcuts_action = help_menu.addAction("快捷鍵列表")
        shortcuts_action.setShortcut(_platform_shortcut("Cmd+/"))
        shortcuts_action.triggered.connect(self._show_shortcuts)

        about_action = help_menu.addAction("關於 SearchingPro")
        about_action.triggered.connect(self._show_about)

    def _apply_theme(self):
        """套用當前主題樣式"""
        is_dark = getattr(self, "dark_mode", False)
        self.setStyleSheet(DARK_STYLE if is_dark else LIGHT_STYLE)
        if hasattr(self, "theme_status_label"):
            self.theme_status_label.setText(" 深色模式 " if is_dark else " 淺色模式 ")

    def _toggle_theme(self):
        """切換深色/淺色模式"""
        self.dark_mode = not getattr(self, "dark_mode", False)
        self._apply_theme()
        self._save_theme_preference()
        if hasattr(self, "theme_toggle_btn"):
            self.theme_toggle_btn.setText(
                "切換為淺色模式" if self.dark_mode else "切換為深色模式"
            )
        if hasattr(self, "statistics_widget"):
            self.statistics_widget.set_dark_mode(self.dark_mode)

    def _save_theme_preference(self):
        """儲存主題偏好到設定檔"""
        try:
            if os.path.exists(self.settings_path):
                with open(self.settings_path, "r", encoding="utf-8") as f:
                    s = json.load(f)
            else:
                s = {}
            s["dark_mode"] = getattr(self, "dark_mode", False)
            with open(self.settings_path, "w", encoding="utf-8") as f:
                json.dump(s, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"儲存主題偏好失敗: {e}")

    def _load_theme_preference(self):
        """載入主題偏好"""
        try:
            if os.path.exists(self.settings_path):
                with open(self.settings_path, "r", encoding="utf-8") as f:
                    s = json.load(f)
                self.dark_mode = s.get("dark_mode", False)
        except Exception:
            self.dark_mode = False

    def export_results_csv(self):
        """匯出搜尋結果"""
        if self.result_table.rowCount() == 0:
            self.statusBar.showMessage("沒有可匯出的搜尋結果")
            return

        file_path, selected_filter = QFileDialog.getSaveFileName(
            self, "匯出結果", "", "CSV 檔案 (*.csv);;JSON 檔案 (*.json);;Excel 檔案 (*.xlsx)"
        )
        if not file_path:
            return

        if selected_filter == "JSON 檔案 (*.json)":
            self._export_to_json(file_path)
        elif selected_filter == "Excel 檔案 (*.xlsx)":
            self._export_to_excel(file_path)
        else:
            self._export_to_csv(file_path)

    def _export_to_csv(self, file_path):
        try:
            import csv
            with open(file_path, "w", newline="", encoding="utf-8-sig") as f:
                writer = csv.writer(f)
                headers = [
                    self.result_table.horizontalHeaderItem(i).text()
                    for i in range(self.result_table.columnCount())
                ]
                writer.writerow(headers)

                for row in range(self.result_table.rowCount()):
                    row_data = []
                    for col in range(self.result_table.columnCount()):
                        item = self.result_table.item(row, col)
                        row_data.append(item.text() if item else "")
                    writer.writerow(row_data)

            self.statusBar.showMessage(f"已匯出 {self.result_table.rowCount()} 個結果到 {file_path}")
        except Exception as e:
            self.statusBar.showMessage(f"匯出失敗: {str(e)}")

    def _export_to_json(self, file_path):
        try:
            import json
            results = []
            for row in range(self.result_table.rowCount()):
                row_data = {}
                for col in range(self.result_table.columnCount()):
                    header = self.result_table.horizontalHeaderItem(col).text()
                    item = self.result_table.item(row, col)
                    row_data[header] = item.text() if item else ""
                results.append(row_data)

            with open(file_path, "w", encoding="utf-8") as f:
                json.dump(results, f, ensure_ascii=False, indent=2)

            self.statusBar.showMessage(f"已匯出 {len(results)} 個結果到 {file_path}")
        except Exception as e:
            self.statusBar.showMessage(f"匯出失敗: {str(e)}")

    def _export_to_excel(self, file_path):
        try:
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            ws.title = "搜尋結果"

            headers = [
                self.result_table.horizontalHeaderItem(i).text()
                for i in range(self.result_table.columnCount())
            ]
            ws.append(headers)

            for row in range(self.result_table.rowCount()):
                row_data = []
                for col in range(self.result_table.columnCount()):
                    item = self.result_table.item(row, col)
                    row_data.append(item.text() if item else "")
                ws.append(row_data)

            wb.save(file_path)
            self.statusBar.showMessage(f"已匯出 {self.result_table.rowCount()} 個結果到 {file_path}")
        except Exception as e:
            self.statusBar.showMessage(f"匯出失敗: {str(e)}")

    def _import_search_results(self):
        """從 CSV 或 JSON 匯入搜尋結果"""
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "匯入搜尋結果",
            "",
            "CSV 檔案 (*.csv);;JSON 檔案 (*.json);;所有檔案 (*)",
        )
        if not file_path:
            return

        records = []
        try:
            if file_path.endswith(".csv"):
                records = self._parse_csv(file_path)
            elif file_path.endswith(".json"):
                records = self._parse_json(file_path)
            else:
                self.statusBar.showMessage("不支援的檔案格式")
                return
        except Exception as e:
            QMessageBox.warning(self, "匯入失敗", f"解析檔案時出錯: {e}")
            return

        if not records:
            self.statusBar.showMessage("匯入的檔案為空")
            return

        if self.result_table.rowCount() > 0:
            reply = QMessageBox.question(
                self,
                "匯入搜尋結果",
                f"目前表格中有 {self.result_table.rowCount()} 筆結果。\n"
                "匯入將覆蓋現有結果，確定要繼續嗎？",
                QMessageBox.Yes | QMessageBox.No,
            )
            if reply != QMessageBox.Yes:
                return

        self._populate_table_from_import(records)
        exists_count = self._check_imported_files_existence()
        total = self.result_table.rowCount()
        self.statusBar.showMessage(
            f"已匯入 {total} 筆結果，其中 {exists_count} 筆檔案仍存在"
        )

    def _parse_csv(self, file_path):
        """解析 CSV 格式的搜尋結果"""
        import csv
        records = []
        with open(file_path, "r", encoding="utf-8-sig", errors="replace") as f:
            reader = csv.DictReader(f)
            for row in reader:
                # Map both Chinese and English column headers
                record = {
                    "filename": row.get("檔名", row.get("Filename", "")),
                    "path": row.get("路徑", row.get("Path", "")),
                    "size": row.get("大小", row.get("Size", "")),
                    "modified": row.get("修改日期", row.get("Modified", "")),
                    "ext": row.get("類型", row.get("Type", "")),
                }
                if record["filename"] or record["path"]:
                    records.append(record)
        return records

    def _parse_json(self, file_path):
        """解析 JSON 格式的搜尋結果"""
        import json
        records = []
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, list):
            raise ValueError("JSON 檔案應包含一個物件陣列")
        for item in data:
            record = {
                "filename": item.get("檔名", item.get("Filename", "")),
                "path": item.get("路徑", item.get("Path", "")),
                "size": item.get("大小", item.get("Size", "")),
                "modified": item.get("修改日期", item.get("Modified", "")),
                "ext": item.get("類型", item.get("Type", "")),
            }
            if record["filename"] or record["path"]:
                records.append(record)
        return records

    def _populate_table_from_import(self, records):
        """將匯入的資料填入結果表格"""
        self.result_table.setRowCount(len(records))
        for i, rec in enumerate(records):
            self.result_table.setItem(i, 0, QTableWidgetItem(rec["filename"]))
            self.result_table.setItem(
                i, 1, QTableWidgetItem(os.path.normpath(rec["path"])) if rec["path"] else QTableWidgetItem("")
            )
            self.result_table.setItem(i, 2, QTableWidgetItem(rec["size"]))
            self.result_table.setItem(i, 3, QTableWidgetItem(rec["modified"]))
            self.result_table.setItem(i, 4, QTableWidgetItem(rec["ext"]))
        self.result_table.resizeColumnsToContents()

    def _check_imported_files_existence(self):
        """檢查匯入的檔案是否存在，對不存在的加標記"""
        exists_count = 0
        for row in range(self.result_table.rowCount()):
            path_item = self.result_table.item(row, 1)
            name_item = self.result_table.item(row, 0)
            if path_item and path_item.text().strip():
                path = os.path.normpath(path_item.text())
                if os.path.exists(path):
                    exists_count += 1
                elif name_item:
                    current_name = name_item.text()
                    if not current_name.startswith("⚠️ "):
                        name_item.setText(f"⚠️ {current_name}")
            elif name_item:
                # No path provided, mark as missing
                current_name = name_item.text()
                if not current_name.startswith("⚠️ "):
                    name_item.setText(f"⚠️ {current_name}")
        return exists_count

    def _clear_search_fields(self):
        """清除所有搜尋條件"""
        self.filename_input.clear()
        self.content_input.clear()
        self.file_type_input.clear()
        self.boolean_input.clear()
        self.path_filter.clear()
        self.min_size_input.setValue(0)
        self.max_size_input.setValue(0)
        self.use_regex_filename.setChecked(False)
        self.use_fuzzy_filename.setChecked(False)
        self.use_regex_content.setChecked(False)
        self.use_fuzzy_content.setChecked(False)
        self.use_boolean.setChecked(False)
        self.use_date_filter.setChecked(False)
        self.statusBar.showMessage("已清除搜尋條件")

    def _toggle_preview(self):
        """切換檔案預覽顯示"""
        if self.preview_widget.isVisible():
            self.preview_widget.hide()
            self.preview_visible = False
        else:
            self.preview_widget.show()
            self.preview_visible = True

    def _toggle_realtime_search(self):
        """切換即時搜尋"""
        self.use_realtime_search.setChecked(not self.use_realtime_search.isChecked())

    def _show_shortcuts(self):
        """顯示快捷鍵列表"""
        ctrl = _shortcut_label("Cmd")
        shortcuts_text = f"""
<b>快捷鍵列表</b><br><br>
<table>
<tr><td><b>{ctrl}+F</b></td><td>聚焦搜尋框</td></tr>
<tr><td><b>{ctrl}+E</b></td><td>匯出搜尋結果</td></tr>
<tr><td><b>{ctrl}+I</b></td><td>匯入搜尋結果</td></tr>
<tr><td><b>{ctrl}+T</b></td><td>切換深色/淺色模式</td></tr>
<tr><td><b>{ctrl}+R</b></td><td>切換即時搜尋</td></tr>
<tr><td><b>{ctrl}+Shift+C</b></td><td>清除搜尋條件</td></tr>
<tr><td><b>F3</b></td><td>切換檔案預覽</td></tr>
<tr><td><b>{ctrl}+Q</b></td><td>離開程式</td></tr>
<tr><td><b>{ctrl}+/</b></td><td>顯示此對話框</td></tr>
<tr><td><b>Enter</b></td><td>執行搜尋</td></tr>
<tr><td><b>Esc</b></td><td>停止搜尋</td></tr>
<tr><td><b>Delete</b></td><td>刪除選中索引</td></tr>
</table>
        """
        msg = QMessageBox(self)
        msg.setWindowTitle("快捷鍵列表")
        msg.setText(shortcuts_text)
        msg.setTextFormat(Qt.RichText)
        msg.exec_()

    def _show_about(self):
        """顯示關於對話框"""
        msg = QMessageBox(self)
        msg.setWindowTitle("關於 SearchingPro")
        msg.setText(
            "<b>SearchingPro</b><br>進階檔案搜尋工具 v2.0<br><br>"
            "功能：模糊搜尋、即時搜尋、布林邏輯、搜尋歷史<br>"
            "技術：Python + PyQt5 + SQLite<br><br>"
            "© 2026 SearchingPro"
        )
        msg.exec_()

    def _init_search_tab(self):
        main_layout = QVBoxLayout()
        main_layout.setContentsMargins(8, 8, 8, 8)
        main_layout.setSpacing(8)

        basic_section = CollapsibleGroupBox(" 基本搜尋")
        basic_layout = QGridLayout()
        basic_layout.setContentsMargins(8, 8, 8, 8)
        basic_layout.setHorizontalSpacing(12)
        basic_layout.setVerticalSpacing(10)

        basic_layout.addWidget(QLabel("檔案名稱:"), 0, 0)
        self.filename_input = QLineEdit()
        self.filename_input.setPlaceholderText("輸入檔案名稱關鍵字...")
        basic_layout.addWidget(self.filename_input, 0, 1)
        self.filename_suggestion = SearchSuggestionWidget(self)
        self.filename_suggestion.suggestion_selected.connect(self._on_filename_suggestion_selected)
        self.filename_input.textChanged.connect(self._on_filename_text_changed)
        mode_layout1 = QHBoxLayout()
        self.use_regex_filename = QCheckBox("正則")
        self.use_fuzzy_filename = QCheckBox("模糊")
        mode_layout1.addWidget(self.use_regex_filename)
        mode_layout1.addWidget(self.use_fuzzy_filename)
        mode_layout1.addStretch()
        basic_layout.addLayout(mode_layout1, 0, 2)

        basic_layout.addWidget(QLabel("檔案內容:"), 1, 0)
        self.content_input = QLineEdit()
        self.content_input.setPlaceholderText("輸入檔案內容關鍵字...")
        basic_layout.addWidget(self.content_input, 1, 1)
        self.content_suggestion = SearchSuggestionWidget(self)
        self.content_suggestion.suggestion_selected.connect(self._on_content_suggestion_selected)
        self.content_input.textChanged.connect(self._on_content_text_changed)
        mode_layout2 = QHBoxLayout()
        self.use_regex_content = QCheckBox("正則")
        self.use_fuzzy_content = QCheckBox("模糊")
        mode_layout2.addWidget(self.use_regex_content)
        mode_layout2.addWidget(self.use_fuzzy_content)
        mode_layout2.addStretch()
        basic_layout.addLayout(mode_layout2, 1, 2)

        action_bar = QHBoxLayout()
        self.search_button = QPushButton(" 搜尋 ")
        self.search_button.setObjectName("successButton")
        self.search_button.clicked.connect(self.start_search)
        self.search_button.setMinimumWidth(100)
        self.stop_search_button = QPushButton(" 停止 ")
        self.stop_search_button.clicked.connect(self.stop_search)
        self.stop_search_button.setEnabled(False)
        self.stop_search_button.setObjectName("deleteButton")
        action_bar.addWidget(self.search_button)
        action_bar.addWidget(self.stop_search_button)
        action_bar.addSpacing(16)
        self.use_realtime_search = QCheckBox("即時搜尋")
        action_bar.addWidget(self.use_realtime_search)
        self.realtime_delay_label = QLabel("(300ms)")
        self.realtime_delay_label.setStyleSheet("color: #8e8e93; font-size: 11px;")
        action_bar.addWidget(self.realtime_delay_label)
        action_bar.addStretch()
        self.use_boolean = QCheckBox("布爾搜尋")
        action_bar.addWidget(self.use_boolean)
        basic_layout.addLayout(action_bar, 2, 0, 1, 3)

        basic_section.setContentLayout(basic_layout)
        main_layout.addWidget(basic_section)

        advanced_section = CollapsibleGroupBox(" 進階過濾")
        advanced_layout = QGridLayout()
        advanced_layout.setContentsMargins(8, 8, 8, 8)
        advanced_layout.setHorizontalSpacing(12)
        advanced_layout.setVerticalSpacing(10)

        advanced_layout.addWidget(QLabel("檔案類型:"), 0, 0)
        self.file_type_input = QLineEdit()
        self.file_type_input.setPlaceholderText(".pdf, .docx, .xlsx (逗號分隔)")
        advanced_layout.addWidget(self.file_type_input, 0, 1)

        size_layout = QHBoxLayout()
        size_layout.addWidget(QLabel("最小:"))
        self.min_size_input = QSpinBox()
        self.min_size_input.setRange(0, 1000000)
        self.min_size_input.setSuffix(" KB")
        size_layout.addWidget(self.min_size_input)
        size_layout.addWidget(QLabel("最大:"))
        self.max_size_input = QSpinBox()
        self.max_size_input.setRange(0, 1000000)
        self.max_size_input.setSuffix(" KB")
        self.max_size_input.setValue(0)
        self.max_size_input.setSpecialValueText("不限")
        size_layout.addWidget(self.max_size_input)
        size_layout.addStretch()
        advanced_layout.addLayout(size_layout, 0, 2)

        advanced_layout.addWidget(QLabel("修改日期:"), 1, 0)
        date_layout = QHBoxLayout()
        self.date_after = QDateEdit()
        self.date_after.setDate(QDate.currentDate().addDays(-30))
        self.date_after.setCalendarPopup(True)
        date_layout.addWidget(self.date_after)
        date_layout.addWidget(QLabel("至"))
        self.date_before = QDateEdit()
        self.date_before.setDate(QDate.currentDate())
        self.date_before.setCalendarPopup(True)
        date_layout.addWidget(self.date_before)
        self.use_date_filter = QCheckBox("啟用")
        date_layout.addWidget(self.use_date_filter)
        date_layout.addStretch()
        advanced_layout.addLayout(date_layout, 1, 1, 1, 2)

        advanced_layout.addWidget(QLabel("路徑過濾:"), 2, 0)
        self.path_filter = QLineEdit()
        self.path_filter.setPlaceholderText("過濾包含此路徑的檔案...")
        advanced_layout.addWidget(self.path_filter, 2, 1)
        advanced_layout.addWidget(QLabel(""), 2, 2)

        advanced_section.setContentLayout(advanced_layout)
        main_layout.addWidget(advanced_section)

        boolean_section = CollapsibleGroupBox(" 布爾搜尋")
        boolean_layout = QVBoxLayout()
        boolean_layout.setContentsMargins(8, 8, 8, 8)
        self.boolean_input = QLineEdit()
        self.boolean_input.setPlaceholderText(
            "支援: AND, OR, NOT, 括號 ()  例如: (report AND budget) OR (meeting AND NOT draft)"
        )
        boolean_layout.addWidget(self.boolean_input)
        hint_label = QLabel(
            "提示: 使用括號分組複雜條件，如 (A AND B) OR (C AND NOT D)"
        )
        hint_label.setStyleSheet("color: #8e8e93; font-size: 11px;")
        boolean_layout.addWidget(hint_label)
        boolean_section.setContentLayout(boolean_layout)
        main_layout.addWidget(boolean_section)

        history_section = CollapsibleGroupBox(" 搜尋歷史")
        history_layout = QHBoxLayout()
        history_layout.setContentsMargins(8, 8, 8, 8)
        history_layout.addWidget(QLabel("歷史記錄:"))
        self.search_history_combo = QComboBox()
        self.search_history_combo.setMinimumWidth(250)
        self.search_history_combo.addItem("(無歷史記錄)")
        history_layout.addWidget(self.search_history_combo)
        self.clear_history_button = QPushButton("清除歷史")
        self.clear_history_button.setObjectName("secondaryButton")
        self.clear_history_button.setMaximumWidth(90)
        history_layout.addWidget(self.clear_history_button)
        history_layout.addStretch()
        history_section.setContentLayout(history_layout)
        main_layout.addWidget(history_section)

        results_toolbar = QHBoxLayout()
        self.result_info = QLabel(" 尚無搜尋結果")
        self.result_info.setStyleSheet("font-weight: 600; font-size: 13px;")
        results_toolbar.addWidget(self.result_info)
        results_toolbar.addStretch()
        export_btn = QPushButton(" 匯出 CSV ")
        export_btn.setObjectName("secondaryButton")
        export_btn.clicked.connect(self.export_results_csv)
        results_toolbar.addWidget(export_btn)
        preview_toggle_btn = QPushButton(" 預覽 F3 ")
        preview_toggle_btn.setObjectName("secondaryButton")
        preview_toggle_btn.clicked.connect(self._toggle_preview)
        results_toolbar.addWidget(preview_toggle_btn)
        main_layout.addLayout(results_toolbar)

        self.result_table = QTableWidget(0, 5)
        self.result_table.setHorizontalHeaderLabels(
            ["檔名", "路徑", "大小", "修改日期", "類型"]
        )
        self.result_table.horizontalHeader().setSectionResizeMode(
            1, QHeaderView.Stretch
        )
        self.result_table.verticalHeader().setVisible(False)
        self.result_table.setSelectionBehavior(QTableWidget.SelectRows)
        self.result_table.setSelectionMode(QTableWidget.ExtendedSelection)
        self.result_table.setEditTriggers(QTableWidget.NoEditTriggers)
        self.result_table.setAlternatingRowColors(True)
        self.result_table.doubleClicked.connect(self.open_file)
        self.result_table.setContextMenuPolicy(Qt.CustomContextMenu)
        self.result_table.customContextMenuRequested.connect(self.show_context_menu)
        main_layout.addWidget(self.result_table)

        self.pagination_widget = PaginationWidget()
        self.pagination_widget.page_changed.connect(self._on_page_changed)
        main_layout.addWidget(self.pagination_widget)

        preview_layout = QVBoxLayout()
        preview_layout.setContentsMargins(4, 4, 4, 4)
        preview_toolbar = QHBoxLayout()
        preview_label = QLabel(" 檔案預覽")
        preview_label.setStyleSheet("font-weight: 600;")
        preview_toolbar.addWidget(preview_label)
        preview_toolbar.addStretch()
        preview_layout.addLayout(preview_toolbar)
        self.content_preview = QTextEdit()
        self.content_preview.setReadOnly(True)
        self.content_preview.setFont(_monospace_font())
        preview_layout.addWidget(self.content_preview)

        splitter = QSplitter(Qt.Vertical)
        results_container = QWidget()
        results_container.setLayout(QVBoxLayout())
        results_container.layout().setContentsMargins(0, 0, 0, 0)
        results_container.layout().addWidget(self.result_table)
        splitter.addWidget(results_container)
        self.preview_widget = QWidget()
        self.preview_widget.setLayout(preview_layout)
        splitter.addWidget(self.preview_widget)
        splitter.setSizes([500, 0])

        self.preview_visible = False
        self.preview_widget.hide()

        main_layout.addWidget(splitter)
        self.search_tab.setLayout(main_layout)

    def _init_index_tab(self):
        main_layout = QVBoxLayout()
        main_layout.setContentsMargins(8, 8, 8, 8)
        main_layout.setSpacing(8)

        action_section = CollapsibleGroupBox(" 索引操作")
        action_layout = QVBoxLayout()
        action_layout.setContentsMargins(8, 8, 8, 8)

        primary_buttons = QHBoxLayout()
        self.start_index_button = QPushButton(" 開始全量索引 ")
        self.start_index_button.setObjectName("successButton")
        self.start_index_button.clicked.connect(self.start_indexing_full)
        primary_buttons.addWidget(self.start_index_button)

        self.update_index_button = QPushButton(" 增量更新 ")
        self.update_index_button.clicked.connect(self.start_indexing_update)
        primary_buttons.addWidget(self.update_index_button)

        self.stop_index_button = QPushButton(" 停止 ")
        self.stop_index_button.clicked.connect(self.stop_indexing)
        self.stop_index_button.setEnabled(False)
        self.stop_index_button.setObjectName("deleteButton")
        primary_buttons.addWidget(self.stop_index_button)

        primary_buttons.addSpacing(16)
        self.last_index_label = QLabel(" 最後索引: 尚未建立")
        self.last_index_label.setStyleSheet("color: #8e8e93; font-size: 12px;")
        primary_buttons.addWidget(self.last_index_label)
        primary_buttons.addStretch()

        action_layout.addLayout(primary_buttons)

        progress_layout = QHBoxLayout()
        progress_layout.addWidget(QLabel("進度:"))
        self.index_progress_info = QLabel("")
        progress_layout.addWidget(self.index_progress_info)
        progress_layout.addStretch()
        action_layout.addLayout(progress_layout)

        action_section.setContentLayout(action_layout)
        main_layout.addWidget(action_section)

        dir_section = CollapsibleGroupBox(" 索引目錄管理")
        dir_layout = QVBoxLayout()
        dir_layout.setContentsMargins(8, 8, 8, 8)

        dir_input_layout = QHBoxLayout()
        self.directory_list = QTextEdit()
        self.directory_list.setPlaceholderText("索引目錄列表 (每行一個路徑)")
        self.directory_list.setMaximumHeight(100)
        dir_input_layout.addWidget(self.directory_list)

        dir_buttons = QVBoxLayout()
        dir_buttons.setSpacing(6)
        add_dir_btn = QPushButton("新增目錄")
        add_dir_btn.clicked.connect(self.add_directory)
        dir_buttons.addWidget(add_dir_btn)
        add_net_btn = QPushButton("網路路徑")
        add_net_btn.clicked.connect(self.add_network_path)
        dir_buttons.addWidget(add_net_btn)
        clear_btn = QPushButton("清空")
        clear_btn.setObjectName("secondaryButton")
        clear_btn.clicked.connect(self.clear_directories)
        dir_buttons.addWidget(clear_btn)
        dir_buttons.addStretch()
        dir_input_layout.addLayout(dir_buttons)

        dir_layout.addLayout(dir_input_layout)

        exclude_row = QHBoxLayout()
        exclude_row.addWidget(QLabel("排除模式:"))
        self.exclude_patterns_input = QLineEdit()
        self.exclude_patterns_input.setPlaceholderText("正則表達式，逗號分隔 (例如: .*\\.git, .*\\.svn)")
        exclude_row.addWidget(self.exclude_patterns_input)
        adv_exclude_btn = QPushButton("進階排除")
        adv_exclude_btn.setObjectName("secondaryButton")
        adv_exclude_btn.clicked.connect(self._show_advanced_exclude_dialog)
        exclude_row.addWidget(adv_exclude_btn)
        dir_layout.addLayout(exclude_row)

        batch_row = QHBoxLayout()
        batch_row.addWidget(QLabel("批次大小:"))
        self.index_batch_size = QSpinBox()
        self.index_batch_size.setRange(1000, 100000)
        self.index_batch_size.setSingleStep(1000)
        self.index_batch_size.setValue(10000)
        batch_row.addWidget(self.index_batch_size)
        batch_row.addStretch()
        dir_layout.addLayout(batch_row)

        watcher_row = QHBoxLayout()
        watcher_row.addWidget(QLabel("即時監控:"))
        self.enable_watcher_btn = QPushButton("啟動監控")
        self.enable_watcher_btn.setObjectName("successButton")
        self.enable_watcher_btn.clicked.connect(self.start_file_watching)
        watcher_row.addWidget(self.enable_watcher_btn)
        self.disable_watcher_btn = QPushButton("停止監控")
        self.disable_watcher_btn.setObjectName("secondaryButton")
        self.disable_watcher_btn.clicked.connect(self.stop_file_watching)
        self.disable_watcher_btn.setEnabled(False)
        watcher_row.addWidget(self.disable_watcher_btn)
        watcher_row.addStretch()
        dir_layout.addLayout(watcher_row)

        dir_section.setContentLayout(dir_layout)
        main_layout.addWidget(dir_section)

        split_layout = QHBoxLayout()
        split_layout.setContentsMargins(0, 0, 0, 0)
        split_layout.setSpacing(8)

        left_panel = QVBoxLayout()
        left_panel.setSpacing(0)

        tree_header = QHBoxLayout()
        tree_header.addWidget(QLabel(" 索引內容"))
        tree_header.addStretch()
        refresh_tree_btn = QPushButton("重新整理")
        refresh_tree_btn.setObjectName("secondaryButton")
        refresh_tree_btn.setMaximumWidth(80)
        refresh_tree_btn.clicked.connect(self.refresh_index_tree)
        tree_header.addWidget(refresh_tree_btn)
        left_panel.addLayout(tree_header)

        self.index_tree = QTreeWidget()
        self.index_tree.setHeaderLabels(["路徑", "檔案數", "總大小", "狀態"])
        self.index_tree.setAlternatingRowColors(True)
        left_panel.addWidget(self.index_tree)

        split_layout.addLayout(left_panel, 2)

        right_panel = QVBoxLayout()
        right_panel.setSpacing(8)

        targets_header = QHBoxLayout()
        targets_header.addWidget(QLabel(" 已索引目標"))
        targets_header.addStretch()
        right_panel.addLayout(targets_header)

        self.update_targets_list = QListWidget()
        self.update_targets_list.setSelectionMode(QListWidget.MultiSelection)
        self.update_targets_list.setToolTip("可勾選要更新的已索引目標資料夾")
        self.update_targets_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.update_targets_list.customContextMenuRequested.connect(
            self._show_targets_context_menu
        )
        right_panel.addWidget(self.update_targets_list)

        stats_box = QGroupBox(" 索引統計")
        stats_layout = QVBoxLayout()
        stats_layout.setContentsMargins(8, 12, 8, 8)
        self.index_stats = QLabel("尚未建立索引")
        self.index_stats.setWordWrap(True)
        stats_layout.addWidget(self.index_stats)
        stats_btn_layout = QHBoxLayout()
        refresh_stats_btn = QPushButton("刷新")
        refresh_stats_btn.setObjectName("secondaryButton")
        refresh_stats_btn.setMaximumWidth(60)
        refresh_stats_btn.clicked.connect(self.refresh_index_stats)
        stats_btn_layout.addWidget(refresh_stats_btn)
        view_history_btn = QPushButton("更新歷史")
        view_history_btn.setObjectName("secondaryButton")
        view_history_btn.setMaximumWidth(80)
        view_history_btn.clicked.connect(self.show_update_history)
        stats_btn_layout.addWidget(view_history_btn)
        stats_btn_layout.addStretch()
        stats_layout.addLayout(stats_btn_layout)
        stats_box.setLayout(stats_layout)
        right_panel.addWidget(stats_box)

        split_layout.addLayout(right_panel, 1)

        main_layout.addLayout(split_layout)

        maint_toolbar = QHBoxLayout()
        maint_toolbar.addWidget(QLabel(" 維護工具:"))
        maint_toolbar.addSpacing(8)

        verify_btn = QPushButton(" 驗證索引 ")
        verify_btn.setObjectName("secondaryButton")
        verify_btn.clicked.connect(self.verify_index)
        maint_toolbar.addWidget(verify_btn)

        cleanup_btn = QPushButton(" 清理索引 ")
        cleanup_btn.setObjectName("secondaryButton")
        cleanup_btn.clicked.connect(self.cleanup_index)
        maint_toolbar.addWidget(cleanup_btn)

        optimize_btn = QPushButton(" 優化資料庫 ")
        optimize_btn.setObjectName("secondaryButton")
        optimize_btn.clicked.connect(self.optimize_database)
        maint_toolbar.addWidget(optimize_btn)

        maint_toolbar.addSpacing(16)
        self.find_dup_btn = QPushButton(" 偵測重複檔案 ")
        self.find_dup_btn.setObjectName("secondaryButton")
        self.find_dup_btn.clicked.connect(self.find_duplicate_files)
        maint_toolbar.addWidget(self.find_dup_btn)

        self.find_empty_btn = QPushButton(" 偵測空檔案 ")
        self.find_empty_btn.setObjectName("secondaryButton")
        self.find_empty_btn.clicked.connect(self.find_empty_files)
        maint_toolbar.addWidget(self.find_empty_btn)

        maint_toolbar.addSpacing(16)
        delete_btn = QPushButton(" 刪除選中 ")
        delete_btn.setObjectName("deleteButton")
        delete_btn.clicked.connect(self.delete_selected_index)
        maint_toolbar.addWidget(delete_btn)

        rebuild_btn = QPushButton(" 重建索引 ")
        rebuild_btn.setObjectName("warningButton")
        rebuild_btn.clicked.connect(self.rebuild_index)
        maint_toolbar.addWidget(rebuild_btn)

        maint_toolbar.addStretch()
        main_layout.addLayout(maint_toolbar)

        self.maintenance_progress = QProgressBar()
        self.maintenance_progress.setVisible(False)
        main_layout.addWidget(self.maintenance_progress)

        self.maintenance_status = QLabel("")
        self.maintenance_status.setStyleSheet("color: #8e8e93; font-size: 11px;")
        main_layout.addWidget(self.maintenance_status)

        self.index_tab.setLayout(main_layout)

    def _init_statistics_tab(self):
        layout = QVBoxLayout()
        layout.setContentsMargins(0, 0, 0, 0)
        self.statistics_widget = StatisticsWidget(self.db_path, self)
        layout.addWidget(self.statistics_widget)
        self.statistics_tab.setLayout(layout)

    def _init_settings_tab(self):
        main_layout = QHBoxLayout()
        main_layout.setContentsMargins(8, 8, 8, 8)
        main_layout.setSpacing(0)

        nav_panel = QWidget()
        nav_panel.setMaximumWidth(180)
        nav_panel.setMinimumWidth(140)
        nav_layout = QVBoxLayout(nav_panel)
        nav_layout.setContentsMargins(0, 0, 0, 0)
        nav_layout.setSpacing(2)

        nav_label = QLabel(" 設定分類")
        nav_label.setStyleSheet("font-weight: 700; font-size: 14px; padding: 8px 4px;")
        nav_layout.addWidget(nav_label)

        nav_sep = QFrame()
        nav_sep.setFrameShape(QFrame.HLine)
        nav_sep.setStyleSheet("background-color: #d1d1d6;")
        nav_layout.addWidget(nav_sep)

        self.settings_nav = QListWidget()
        self.settings_nav.setViewMode(QListWidget.ListMode)
        self.settings_nav.setSpacing(2)
        self.settings_nav.addItems([
            " 搜尋",
            " 索引",
            " 自動更新",
            " 外觀",
            " 系統整合",
            " 資料庫",
        ])
        self.settings_nav.setCurrentRow(0)
        self.settings_nav.setMaximumWidth(170)
        nav_layout.addWidget(self.settings_nav)

        self.settings_stack = QStackedWidget()
        self.settings_stack.setContentsMargins(12, 0, 0, 0)

        self._init_search_settings()
        self._init_index_settings()
        self._init_auto_update_settings()
        self._init_appearance_settings()
        self._init_system_settings()
        self._init_db_settings()

        self.settings_nav.currentRowChanged.connect(self.settings_stack.setCurrentIndex)

        main_layout.addWidget(nav_panel)
        main_layout.addWidget(self.settings_stack, 1)
        self.settings_tab.setLayout(main_layout)

    def _init_search_settings(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        layout.setContentsMargins(0, 12, 12, 12)
        layout.setSpacing(12)

        search_box = QGroupBox(" 搜尋設定")
        search_layout = QGridLayout()
        search_layout.setContentsMargins(12, 16, 12, 12)
        search_layout.setHorizontalSpacing(12)
        search_layout.setVerticalSpacing(10)

        search_layout.addWidget(QLabel("預設排序:"), 0, 0)
        self.sort_by = QComboBox()
        self.sort_by.addItems(["檔名", "大小", "修改日期", "路徑"])
        search_layout.addWidget(self.sort_by, 0, 1)
        self.sort_desc = QCheckBox("降序")
        search_layout.addWidget(self.sort_desc, 0, 2)

        search_layout.addWidget(QLabel("搜尋執行緒:"), 1, 0)
        self.search_threads = QSpinBox()
        self.search_threads.setRange(1, 16)
        self.search_threads.setValue(4)
        search_layout.addWidget(self.search_threads, 1, 1)

        search_layout.addWidget(QLabel("最大記憶體 (MB):"), 2, 0)
        self.max_memory = QSpinBox()
        self.max_memory.setRange(100, 4096)
        self.max_memory.setValue(1024)
        search_layout.addWidget(self.max_memory, 2, 1)

        search_box.setLayout(search_layout)
        layout.addWidget(search_box)
        layout.addStretch()

        self.settings_stack.addWidget(panel)

    def _init_index_settings(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        layout.setContentsMargins(0, 12, 12, 12)
        layout.setSpacing(12)

        index_box = QGroupBox(" 索引設定")
        index_layout = QGridLayout()
        index_layout.setContentsMargins(12, 16, 12, 12)
        index_layout.setHorizontalSpacing(12)
        index_layout.setVerticalSpacing(10)

        index_layout.addWidget(QLabel("索引執行緒:"), 0, 0)
        self.index_threads = QSpinBox()
        self.index_threads.setRange(1, 16)
        self.index_threads.setValue(4)
        index_layout.addWidget(self.index_threads, 0, 1)

        self.index_content = QCheckBox("索引檔案內容")
        self.index_content.setChecked(True)
        index_layout.addWidget(self.index_content, 1, 0)

        ct_layout = QHBoxLayout()
        ct_layout.addWidget(QLabel("內容檔案類型:"))
        self.content_types = QLineEdit()
        self.content_types.setText(".txt, .py, .java, .c, .cpp, .html, .xml, .json, .csv, .md, .pdf, .docx, .xlsx, .pptx")
        ct_layout.addWidget(self.content_types)
        index_layout.addLayout(ct_layout, 2, 0, 1, 2)

        index_box.setLayout(index_layout)
        layout.addWidget(index_box)
        layout.addStretch()

        self.settings_stack.addWidget(panel)

    def _init_auto_update_settings(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        layout.setContentsMargins(0, 12, 12, 12)
        layout.setSpacing(12)

        update_box = QGroupBox(" 自動更新設定")
        update_layout = QGridLayout()
        update_layout.setContentsMargins(12, 16, 12, 12)
        update_layout.setHorizontalSpacing(12)
        update_layout.setVerticalSpacing(10)

        update_layout.addWidget(QLabel("更新頻率:"), 0, 0)
        self.auto_update = QComboBox()
        self.auto_update.addItems(["不自動更新", "每天", "每週", "自定義間隔", "手動"])
        self.auto_update.currentIndexChanged.connect(self._on_auto_update_changed)
        update_layout.addWidget(self.auto_update, 0, 1)

        self.custom_interval_widget = QWidget()
        cil = QHBoxLayout()
        cil.setContentsMargins(0, 0, 0, 0)
        cil.addWidget(QLabel("間隔:"))
        self.custom_interval_hours = QSpinBox()
        self.custom_interval_hours.setRange(1, 168)
        self.custom_interval_hours.setValue(6)
        self.custom_interval_hours.setSuffix(" 小時")
        cil.addWidget(self.custom_interval_hours)
        cil.addStretch()
        self.custom_interval_widget.setLayout(cil)
        update_layout.addWidget(self.custom_interval_widget, 1, 0, 1, 2)

        self.daily_time_widget = QWidget()
        dtl = QHBoxLayout()
        dtl.setContentsMargins(0, 0, 0, 0)
        from PyQt5.QtCore import QTime
        dtl.addWidget(QLabel("執行時間:"))
        self.daily_update_time = QTimeEdit()
        self.daily_update_time.setTime(QTime(2, 0))
        self.daily_update_time.setDisplayFormat("HH:mm")
        dtl.addWidget(self.daily_update_time)
        dtl.addStretch()
        self.daily_time_widget.setLayout(dtl)
        update_layout.addWidget(self.daily_time_widget, 2, 0, 1, 2)

        self.manual_auto_update_button = QPushButton("立即執行")
        self.manual_auto_update_button.setObjectName("successButton")
        self.manual_auto_update_button.clicked.connect(self._manual_trigger_auto_update)
        update_layout.addWidget(self.manual_auto_update_button, 3, 0, 1, 2)

        self.auto_update_status = QLabel("狀態: 未設定")
        self.auto_update_status.setStyleSheet("color: #8e8e93; font-size: 12px;")
        update_layout.addWidget(self.auto_update_status, 4, 0, 1, 2)

        self.next_run_preview = QLabel("下次執行: 未排程")
        self.next_run_preview.setStyleSheet("color: #007aff; font-size: 12px; font-weight: 600;")
        update_layout.addWidget(self.next_run_preview, 5, 0, 1, 2)

        self.preview_update_timer = QTimer()
        self.preview_update_timer.timeout.connect(self._update_next_run_preview)
        self.preview_update_timer.start(1000)

        update_box.setLayout(update_layout)
        layout.addWidget(update_box)
        layout.addStretch()

        self.settings_stack.addWidget(panel)

    def _init_appearance_settings(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        layout.setContentsMargins(0, 12, 12, 12)
        layout.setSpacing(12)

        appearance_box = QGroupBox(" 外觀設定")
        appearance_layout = QVBoxLayout()
        appearance_layout.setContentsMargins(12, 16, 12, 12)
        appearance_layout.setSpacing(12)

        theme_row = QHBoxLayout()
        theme_row.addWidget(QLabel("主題模式:"))
        self.theme_toggle_btn = QPushButton("切換為深色模式")
        self.theme_toggle_btn.setObjectName("secondaryButton")
        self.theme_toggle_btn.clicked.connect(self._toggle_theme)
        theme_row.addWidget(self.theme_toggle_btn)
        theme_row.addStretch()
        appearance_layout.addLayout(theme_row)

        theme_preview = QLabel("預覽: 當前為淺色主題，深色模式適合夜間使用")
        theme_preview.setStyleSheet("color: #8e8e93; font-size: 12px;")
        appearance_layout.addWidget(theme_preview)

        language_row = QHBoxLayout()
        language_row.addWidget(QLabel("語言:"))
        self.language_combo = QComboBox()
        self.language_combo.addItem("繁體中文", "zh-TW")
        self.language_combo.addItem("簡體中文", "zh-CN")
        self.language_combo.addItem("英文", "en")
        self.language_combo.setCurrentIndex(0)
        self.language_combo.currentIndexChanged.connect(self._on_language_changed)
        self.language_combo.setFixedWidth(150)
        language_row.addWidget(self.language_combo)
        language_row.addStretch()
        appearance_layout.addLayout(language_row)

        language_preview = QLabel("切換語言後部分介面文字將更新")
        language_preview.setStyleSheet("color: #8e8e93; font-size: 12px;")
        appearance_layout.addWidget(language_preview)

        appearance_box.setLayout(appearance_layout)
        layout.addWidget(appearance_box)
        layout.addStretch()

        self.settings_stack.addWidget(panel)

    def _on_language_changed(self, index):
        """語言切換"""
        language_code = self.language_combo.itemData(index)
        if i18n.load_language(language_code):
            self._update_ui_language()
            self._save_language_preference(language_code)
            self.statusBar.showMessage(f"語言已切換為 {i18n.get_language_name(language_code)}")

    def _update_ui_language(self):
        """更新介面文字"""
        self.setWindowTitle(_tr("app_title"))
        self.tabs.setTabText(0, f" {_tr('tabs.search')} ")
        self.tabs.setTabText(1, f" {_tr('tabs.index')} ")
        self.tabs.setTabText(2, f" {_tr('tabs.statistics')} ")
        self.tabs.setTabText(3, f" {_tr('tabs.settings')} ")
        if hasattr(self, "statistics_widget"):
            self.statistics_widget.title_label.setText(_tr("statistics.title"))
            self.statistics_widget.refresh_btn.setText(_tr("statistics.refresh"))
        self.statusBar.showMessage(f" {_tr('status.ready')} ")

    def _save_language_preference(self, language_code):
        """儲存語言偏好"""
        try:
            if os.path.exists(self.settings_path):
                with open(self.settings_path, "r", encoding="utf-8") as f:
                    s = json.load(f)
            else:
                s = {}
            s["language"] = language_code
            with open(self.settings_path, "w", encoding="utf-8") as f:
                json.dump(s, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"儲存語言偏好失敗: {e}")

    def _load_language_preference(self):
        """載入語言偏好"""
        try:
            if os.path.exists(self.settings_path):
                with open(self.settings_path, "r", encoding="utf-8") as f:
                    s = json.load(f)
                language_code = s.get("language", "zh-TW")
                if language_code in I18nManager.SUPPORTED_LANGUAGES:
                    i18n.load_language(language_code)
                    return language_code
        except Exception:
            pass
        return "zh-TW"

    def _init_system_settings(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        layout.setContentsMargins(0, 12, 12, 12)
        layout.setSpacing(12)

        system_box = QGroupBox(" 系統整合")
        system_layout = QVBoxLayout()
        system_layout.setContentsMargins(12, 16, 12, 12)
        system_layout.setSpacing(10)

        self.start_with_system = QCheckBox("系統啟動時自動啟動")
        self.start_with_system.stateChanged.connect(self._on_start_with_system_changed)
        system_layout.addWidget(self.start_with_system)

        self.minimize_to_tray = QCheckBox("最小化到系統匣")
        system_layout.addWidget(self.minimize_to_tray)

        self.no_ask_close_to_tray = QCheckBox("關閉視窗時直接縮到系統匣（不再提醒）")
        system_layout.addWidget(self.no_ask_close_to_tray)

        system_box.setLayout(system_layout)
        layout.addWidget(system_box)
        layout.addStretch()

        self.settings_stack.addWidget(panel)

    def _init_db_settings(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        layout.setContentsMargins(0, 12, 12, 12)
        layout.setSpacing(12)

        db_box = QGroupBox(" 資料庫設定")
        db_layout = QGridLayout()
        db_layout.setContentsMargins(12, 16, 12, 12)
        db_layout.setHorizontalSpacing(12)
        db_layout.setVerticalSpacing(10)

        db_layout.addWidget(QLabel("資料庫位置:"), 0, 0)
        self.db_path_display = QLineEdit()
        self.db_path_display.setText(self.db_path)
        self.db_path_display.setReadOnly(True)
        db_layout.addWidget(self.db_path_display, 0, 1)

        change_db_path_button = QPushButton("更改位置")
        change_db_path_button.setObjectName("secondaryButton")
        change_db_path_button.clicked.connect(self.change_db_path)
        db_layout.addWidget(change_db_path_button, 0, 2)

        db_box.setLayout(db_layout)
        layout.addWidget(db_box)

        save_row = QHBoxLayout()
        save_row.addStretch()
        save_button = QPushButton(" 保存所有設定 ")
        save_button.setObjectName("successButton")
        save_button.setMinimumWidth(120)
        save_button.clicked.connect(self.save_settings)
        save_row.addWidget(save_button)
        layout.addLayout(save_row)

        layout.addStretch()

        self.settings_stack.addWidget(panel)

    # ============ 新增的索引管理功能 ============

    def refresh_index_tree(self):
        """重新整理索引樹狀顯示"""
        self.index_tree.clear()

        if not os.path.exists(self.db_path):
            return

        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            # 取得所有不同的根目錄
            cursor.execute("""
                SELECT 
                    SUBSTR(path, 1, LENGTH(path) - LENGTH(REPLACE(path, '/', ''))) as root_path,
                    COUNT(*) as file_count,
                    SUM(size) as total_size
                FROM files
                GROUP BY root_path
                ORDER BY root_path
            """)

            # 建立根目錄對應的樹狀項目
            root_items = {}

            # 首先取得所有檔案路徑以建立目錄結構
            cursor.execute("SELECT path, size FROM files ORDER BY path")
            all_files = cursor.fetchall()

            # 建立目錄統計
            dir_stats = {}
            for file_path, size in all_files:
                # 取得目錄路徑
                dir_path = os.path.dirname(file_path)

                # 累計每個目錄的統計
                while dir_path:
                    if dir_path not in dir_stats:
                        dir_stats[dir_path] = {"count": 0, "size": 0}
                    dir_stats[dir_path]["count"] += 1
                    dir_stats[dir_path]["size"] += size

                    # 往上一層目錄
                    parent_dir = os.path.dirname(dir_path)
                    if parent_dir == dir_path:  # 到達根目錄
                        break
                    dir_path = parent_dir

            # 建立樹狀結構
            added_dirs = set()
            for dir_path in sorted(dir_stats.keys()):
                if dir_path in added_dirs:
                    continue

                stats = dir_stats[dir_path]

                # 格式化大小顯示
                size_mb = stats["size"] / (1024 * 1024)
                if size_mb < 1024:
                    size_str = f"{size_mb:.2f} MB"
                else:
                    size_gb = size_mb / 1024
                    size_str = f"{size_gb:.2f} GB"

                # 檢查目錄是否存在
                if os.path.exists(dir_path):
                    status = "正常"
                    color = QColor(0, 128, 0)  # 綠色
                else:
                    status = "不存在"
                    color = QColor(255, 0, 0)  # 紅色

                # 建立樹狀項目
                item = QTreeWidgetItem(
                    [dir_path, str(stats["count"]), size_str, status]
                )

                # 設定狀態欄顏色
                item.setForeground(3, QBrush(color))

                self.index_tree.addTopLevelItem(item)
                added_dirs.add(dir_path)

            conn.close()

            # 調整欄位寬度
            for i in range(4):
                self.index_tree.resizeColumnToContents(i)

        except Exception as e:
            print(f"重新整理索引樹狀顯示時出錯: {e}")

    def verify_index(self):
        """驗證索引完整性"""
        reply = QMessageBox.question(
            self,
            "確認驗證",
            "要驗證索引的完整性嗎？\n這可能需要一些時間。",
            QMessageBox.Yes | QMessageBox.No,
        )

        if reply != QMessageBox.Yes:
            return

        self.maintenance_progress.setVisible(True)
        self.maintenance_status.setText("正在驗證索引...")

        self.maintenance_thread = IndexMaintenanceThread(self.db_path, "verify")
        self.maintenance_thread.progress_update.connect(
            self.update_maintenance_progress
        )
        self.maintenance_thread.maintenance_complete.connect(self.maintenance_complete)
        self.maintenance_thread.start()

    def cleanup_index(self):
        """清理索引"""
        dialog = IndexCleanupDialog(self)
        if dialog.exec_() == QDialog.Accepted:
            options = dialog.get_options()

            self.maintenance_progress.setVisible(True)
            self.maintenance_status.setText("正在清理索引...")

            self.maintenance_thread = IndexMaintenanceThread(
                self.db_path, "cleanup", options
            )
            self.maintenance_thread.progress_update.connect(
                self.update_maintenance_progress
            )
            self.maintenance_thread.maintenance_complete.connect(
                self.maintenance_complete
            )
            self.maintenance_thread.start()

    def optimize_database(self):
        """優化資料庫"""
        reply = QMessageBox.question(
            self,
            "確認優化",
            "要優化資料庫嗎？\n這將重建索引並清理未使用的空間。",
            QMessageBox.Yes | QMessageBox.No,
        )

        if reply != QMessageBox.Yes:
            return

        self.maintenance_progress.setVisible(True)
        self.maintenance_status.setText("正在優化資料庫...")

        self.maintenance_thread = IndexMaintenanceThread(self.db_path, "optimize")
        self.maintenance_thread.progress_update.connect(
            self.update_maintenance_progress
        )
        self.maintenance_thread.maintenance_complete.connect(self.maintenance_complete)
        self.maintenance_thread.start()

    def delete_selected_index(self):
        """刪除選中的索引"""
        selected_items = self.index_tree.selectedItems()
        if not selected_items:
            QMessageBox.information(self, "提示", "請先選擇要刪除的索引項目。")
            return

        paths = [item.text(0) for item in selected_items]

        reply = QMessageBox.warning(
            self,
            "確認刪除",
            f"確定要刪除以下 {len(paths)} 個路徑的索引嗎？\n\n"
            + "\n".join(paths[:5])
            + ("\n..." if len(paths) > 5 else ""),
            QMessageBox.Yes | QMessageBox.No,
        )

        if reply != QMessageBox.Yes:
            return

        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            deleted_count = 0
            for path in paths:
                cursor.execute("DELETE FROM files WHERE path LIKE ?", (f"{path}%",))
                deleted_count += cursor.rowcount

            conn.commit()
            conn.close()

            QMessageBox.information(
                self, "刪除完成", f"已刪除 {deleted_count} 個檔案的索引。"
            )

            # 重新整理顯示
            self.refresh_index_tree()
            self.refresh_index_stats()
            self.refresh_indexed_targets_list()

        except Exception as e:
            QMessageBox.critical(self, "刪除失敗", f"刪除索引時發生錯誤：\n{str(e)}")

    def _compute_file_hash(self, path, algorithm="md5"):
        """計算檔案的 hash 值"""
        import hashlib
        h = hashlib.md5() if algorithm == "md5" else hashlib.sha256()
        try:
            with open(path, "rb") as f:
                while True:
                    chunk = f.read(8192)
                    if not chunk:
                        break
                    h.update(chunk)
            return h.hexdigest()
        except Exception:
            return None

    def find_duplicate_files(self):
        """偵測重複檔案"""
        self.maintenance_progress.setVisible(True)
        self.maintenance_progress.setValue(0)
        self.maintenance_status.setText("正在偵測重複檔案...")
        self._dup_scan_running = True

        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            # Step 1: Find files with duplicate sizes
            cursor.execute("""
                SELECT size, COUNT(*) as cnt FROM files
                WHERE size > 0 GROUP BY size HAVING cnt > 1
            """)
            size_groups = cursor.fetchall()

            if not size_groups:
                self.maintenance_progress.setVisible(False)
                self.maintenance_status.setText("未找到重複檔案")
                QMessageBox.information(self, "偵測完成", "未找到重複的檔案。")
                conn.close()
                return

            # Step 2: For each size group, compute hashes to confirm duplicates
            dup_groups = []
            total_checked = 0
            total_groups = len(size_groups)

            for size, count in size_groups:
                if not getattr(self, "_dup_scan_running", False):
                    self.maintenance_status.setText("偵測已取消")
                    self.maintenance_progress.setVisible(False)
                    conn.close()
                    return

                progress_pct = int((total_checked / total_groups) * 100)
                self.maintenance_progress.setValue(progress_pct)
                self.maintenance_status.setText(
                    f"偵測中... ({total_checked}/{total_groups} 組)"
                )

                cursor.execute("SELECT path, filename FROM files WHERE size = ?", (size,))
                candidates = cursor.fetchall()

                hash_map = {}
                for path, filename in candidates:
                    file_hash = self._compute_file_hash(path)
                    if file_hash:
                        if file_hash not in hash_map:
                            hash_map[file_hash] = []
                        hash_map[file_hash].append((path, filename, size))

                for h, files in hash_map.items():
                    if len(files) > 1:
                        dup_groups.append((h, files))
                total_checked += 1

            conn.close()
            self.maintenance_progress.setVisible(False)

            if not dup_groups:
                self.maintenance_status.setText("未找到重複檔案")
                QMessageBox.information(self, "偵測完成", "未找到內容相同的重複檔案。")
                return

            dialog = DuplicateFilesDialog(dup_groups, self)
            dialog.exec_()

        except Exception as e:
            self.maintenance_progress.setVisible(False)
            QMessageBox.critical(self, "偵測失敗", f"偵測重複檔案時出錯: {e}")

    def find_empty_files(self):
        """偵測空檔案"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute(
                "SELECT path, filename FROM files WHERE size = 0 ORDER BY path"
            )
            empty_files = cursor.fetchall()
            conn.close()

            if not empty_files:
                QMessageBox.information(self, "偵測完成", "未找到空檔案。")
                return

            dialog = EmptyFilesDialog(empty_files, self)
            dialog.exec_()

        except Exception as e:
            QMessageBox.critical(self, "偵測失敗", f"偵測空檔案時出錯: {e}")

    def rebuild_index(self):
        """重建索引結構"""
        reply = QMessageBox.warning(
            self,
            "確認重建",
            "重建索引將刪除並重新建立所有索引表格。\n"
            "現有資料將盡可能保留，但此操作有風險。\n\n"
            "確定要繼續嗎？",
            QMessageBox.Yes | QMessageBox.No,
        )

        if reply != QMessageBox.Yes:
            return

        self.maintenance_progress.setVisible(True)
        self.maintenance_status.setText("正在重建索引...")

        self.maintenance_thread = IndexMaintenanceThread(self.db_path, "rebuild")
        self.maintenance_thread.progress_update.connect(
            self.update_maintenance_progress
        )
        self.maintenance_thread.maintenance_complete.connect(self.maintenance_complete)
        self.maintenance_thread.start()

    def update_maintenance_progress(self, message, progress):
        """更新維護進度"""
        self.maintenance_status.setText(message)
        self.maintenance_progress.setValue(progress)

    def maintenance_complete(self, stats):
        """維護操作完成"""
        self.maintenance_progress.setVisible(False)

        # 根據不同操作顯示結果
        if "deleted_files" in stats:  # 清理操作
            message = f"清理完成：\n"
            message += f"- 刪除不存在的檔案: {stats['deleted_files']}\n"
            message += f"- 移除指定路徑: {stats['removed_by_path']}\n"
            message += f"- 移除指定類型: {stats['removed_by_type']}\n"
            message += f"- 移除大型檔案: {stats['removed_by_size']}\n"
            message += f"- 移除重複項目: {stats['removed_duplicates']}\n"
            message += f"總計移除: {stats['total_removed']} 個項目"

        elif "missing_files" in stats:  # 驗證操作
            message = f"驗證完成：\n"
            message += f"- 總檔案數: {stats['total_files']}\n"
            message += f"- 不存在的檔案: {stats['missing_files']}\n"
            message += f"- 重複項目: {stats['duplicate_entries']}"

        elif "space_saved" in stats:  # 優化操作
            message = f"優化完成：\n"
            message += f"- 原始大小: {stats['original_size'] / (1024 * 1024):.2f} MB\n"
            message += (
                f"- 優化後大小: {stats['optimized_size'] / (1024 * 1024):.2f} MB\n"
            )
            message += f"- 節省空間: {stats['space_saved'] / (1024 * 1024):.2f} MB"

        elif "tables_rebuilt" in stats:  # 重建操作
            message = f"重建完成：\n"
            message += f"- 重建表格: {stats['tables_rebuilt']}\n"
            message += f"- 重建索引: {stats['indexes_rebuilt']}"
        else:
            message = "操作完成"

        if stats.get("errors"):
            message += f"\n\n錯誤：\n" + "\n".join(stats["errors"][:5])

        self.maintenance_status.setText("操作完成")
        QMessageBox.information(self, "操作完成", message)

        # 重新整理顯示
        self.refresh_index_tree()
        self.refresh_index_stats()
        self.refresh_indexed_targets_list()

    # ---------------- 系統匣 ----------------
    def _init_tray_icon(self):
        if not QSystemTrayIcon.isSystemTrayAvailable():
            return
        self.tray_icon = QSystemTrayIcon(self)

        # 使用與主視窗相同的圖示
        icon_set = False
        try:
            icon_path = self._get_icon_path()
            if icon_path and os.path.exists(icon_path):
                icon = QIcon(icon_path)
                if not icon.isNull():
                    self.tray_icon.setIcon(icon)
                    icon_set = True
                else:
                    print(f"圖示檔案無效: {icon_path}")
            else:
                print(f"圖示檔案不存在: {icon_path}")
        except Exception as e:
            print(f"設定系統匣圖示時發生錯誤: {e}")

        # 如果主圖示載入失敗，使用備用圖示
        if not icon_set:
            try:
                backup_icon = self.style().standardIcon(QStyle.SP_ComputerIcon)
                self.tray_icon.setIcon(backup_icon)
            except Exception as e:
                print(f"設定備用系統匣圖示時發生錯誤: {e}")

        self.tray_icon.setToolTip("進階檔案搜尋工具")

        menu = QMenu()

        toggle_action = QAction("顯示/隱藏視窗", self)
        toggle_action.triggered.connect(self.toggle_main_window)
        menu.addAction(toggle_action)

        start_full_action = QAction("開始索引（全量）", self)
        start_full_action.triggered.connect(self.start_indexing_full)
        menu.addAction(start_full_action)

        refresh_action = QAction("刷新統計", self)
        refresh_action.triggered.connect(self.refresh_index_stats)
        menu.addAction(refresh_action)

        menu.addSeparator()

        exit_action = QAction("退出", self)
        exit_action.triggered.connect(self.exit_app)
        menu.addAction(exit_action)

        self.tray_icon.setContextMenu(menu)
        self.tray_icon.activated.connect(self._on_tray_activated)
        self.tray_icon.show()

    def _on_tray_activated(self, reason):
        # 左鍵雙擊或單擊（依平台行為）顯示主視窗
        if reason in (QSystemTrayIcon.DoubleClick, QSystemTrayIcon.Trigger):
            self.showNormal()
            self.activateWindow()
            self.raise_()

    def toggle_main_window(self):
        if self.isVisible() and not self.isMinimized():
            self.hide()
        else:
            self.showNormal()
            self.activateWindow()
            self.raise_()

    def exit_app(self):
        self.tray_icon.hide() if self.tray_icon else None
        QApplication.instance().quit()

    def changeEvent(self, event):
        # 監聽最小化事件，若勾選最小化到匣則隱藏
        if event.type() == QEvent.WindowStateChange:
            if self.windowState() & Qt.WindowMinimized:
                if self.minimize_to_tray.isChecked():
                    QTimer.singleShot(0, self.hide)  # 讓最小化先完成，再隱藏
                    self.statusBar.showMessage("程式已縮到系統匣")
        super().changeEvent(event)

    def closeEvent(self, event):
        # 關閉時處理縮到匣
        if self.minimize_to_tray.isChecked():
            if self.no_ask_close_to_tray.isChecked() or self.close_to_tray_no_ask:
                event.ignore()
                self.hide()
                self.statusBar.showMessage("程式已縮到系統匣（關閉視窗不會退出）")
                return
            # 詢問一次
            box = QMessageBox(self)
            box.setIcon(QMessageBox.Question)
            box.setWindowTitle("最小化到系統匣")
            box.setText("要關閉視窗並縮到系統匣嗎？\n選擇「否」將直接退出程式。")
            box.setStandardButtons(QMessageBox.Yes | QMessageBox.No)
            dont_ask = QCheckBox("不再提醒（之後直接縮到系統匣）", box)
            box.setCheckBox(dont_ask)
            ret = box.exec_()
            if ret == QMessageBox.Yes:
                event.ignore()
                self.hide()
                if dont_ask.isChecked():
                    self.close_to_tray_no_ask = True
                    self.save_settings_to_disk()
                self.statusBar.showMessage("程式已縮到系統匣")
                return
        # 確保完全退出程序
        self._cleanup_on_exit()
        self.tray_icon.hide() if self.tray_icon else None
        QApplication.instance().quit()
        super().closeEvent(event)

    def _cleanup_on_exit(self):
        """清理所有執行中的執行緒、計時器和監控"""
        # 停止檔案監控
        if hasattr(self, "file_watcher") and self.file_watcher:
            try:
                self.file_watcher.stop_watching()
            except Exception:
                pass

        # 停止索引執行緒
        if hasattr(self, "indexer") and self.indexer and self.indexer.isRunning():
            self.indexer.is_running = False
            self.indexer.quit()
            self.indexer.wait()

        # 停止搜尋執行緒
        if hasattr(self, "searcher") and self.searcher and self.searcher.isRunning():
            self.searcher.is_running = False
            self.searcher.quit()
            self.searcher.wait()

        # 停止維護執行緒
        if hasattr(self, "maintenance_thread") and self.maintenance_thread and self.maintenance_thread.isRunning():
            self.maintenance_thread.is_running = False
            self.maintenance_thread.quit()
            self.maintenance_thread.wait()

        # 停止所有計時器
        for timer_attr in ["preview_update_timer", "auto_update_timer", "daily_timer", "weekly_timer"]:
            timer = getattr(self, timer_attr, None)
            if timer:
                try:
                    timer.stop()
                except Exception:
                    pass

    # ------------- 系統啟動設定 -------------
    def _get_registry_run_key(self):
        """取得Windows註冊表Run鍵"""
        if winreg is None:
            return None
        try:
            return winreg.OpenKey(
                winreg.HKEY_CURRENT_USER,
                r"Software\Microsoft\Windows\CurrentVersion\Run",
                0,
                winreg.KEY_ALL_ACCESS,
            )
        except Exception as e:
            print(f"無法開啟註冊表鍵: {e}")
            return None

    def _is_auto_start_enabled(self):
        """檢查是否已設定開機啟動"""
        import platform

        if platform.system() == "Windows":
            if winreg is None:
                return False

            key = self._get_registry_run_key()
            if key is None:
                return False

            try:
                value, _ = winreg.QueryValueEx(key, APP_NAME)
                winreg.CloseKey(key)
                return True
            except FileNotFoundError:
                winreg.CloseKey(key)
                return False
            except Exception as e:
                print(f"檢查開機啟動狀態時發生錯誤: {e}")
                winreg.CloseKey(key)
                return False

        elif platform.system() == "Darwin":
            plist_path = os.path.expanduser(f"~/Library/LaunchAgents/{APP_NAME}.plist")
            return os.path.exists(plist_path)

        return False

    def _set_auto_start(self, enable):
        """設定或取消開機啟動"""
        import platform

        if platform.system() == "Windows":
            if winreg is None:
                QMessageBox.warning(self, "不支援的系統", "無法存取 Windows 註冊表")
                return False

            key = self._get_registry_run_key()
            if key is None:
                QMessageBox.warning(self, "設定失敗", "無法開啟註冊表")
                return False

            try:
                if enable:
                    exe_path = sys.executable
                    if getattr(sys, "frozen", False):
                        # PyInstaller 打包後的路徑
                        exe_path = sys.executable
                    winreg.SetValueEx(key, APP_NAME, 0, winreg.REG_SZ, exe_path)
                else:
                    try:
                        winreg.DeleteValue(key, APP_NAME)
                    except FileNotFoundError:
                        pass  # 值不存在，忽略
                winreg.CloseKey(key)
                return True
            except Exception as e:
                QMessageBox.warning(
                    self, "設定失敗", f"設定開機啟動時發生錯誤: {str(e)}"
                )
                try:
                    winreg.CloseKey(key)
                except Exception:
                    pass
                return False

        elif platform.system() == "Darwin":  # macOS
            import plistlib

            launch_agents_dir = os.path.expanduser("~/Library/LaunchAgents")
            plist_path = os.path.join(launch_agents_dir, f"{APP_NAME}.plist")

            try:
                if enable:
                    os.makedirs(launch_agents_dir, exist_ok=True)

                    # 取得執行檔路徑
                    if getattr(sys, "frozen", False):
                        # PyInstaller 打包的 .app
                        exe_path = os.path.join(
                            os.path.dirname(os.path.dirname(sys.executable)),
                            "MacOS",
                            os.path.basename(sys.executable),
                        )
                        if not os.path.exists(exe_path):
                            exe_path = sys.executable
                        plist_data = {
                            "Label": APP_NAME,
                            "ProgramArguments": [exe_path],
                            "RunAtLoad": True,
                            "KeepAlive": False,
                        }
                    else:
                        # 開發模式：python SearchingPro.py
                        exe_path = sys.executable
                        script_path = os.path.abspath(
                            os.path.join(os.path.dirname(__file__), "SearchingPro.py")
                        ) if "__file__" in dir() else os.path.join(
                            os.getcwd(), "SearchingPro.py"
                        )
                        plist_data = {
                            "Label": APP_NAME,
                            "ProgramArguments": [exe_path, script_path],
                            "RunAtLoad": True,
                            "KeepAlive": False,
                        }

                    with open(plist_path, "wb") as f:
                        plistlib.dump(plist_data, f)

                    os.system(f"launchctl load {plist_path}")
                    return True
                else:
                    if os.path.exists(plist_path):
                        os.system(f"launchctl unload {plist_path}")
                        os.remove(plist_path)
                    return True

            except Exception as e:
                QMessageBox.warning(
                    self, "設定失敗", f"設定開機啟動時發生錯誤: {str(e)}"
                )
                return False

        else:
            QMessageBox.information(
                self, "不支援的系統", "此功能目前只支援 Windows 和 macOS"
            )
            return False

    def _on_start_with_system_changed(self):
        """系統啟動設定變更時的處理"""
        enable = self.start_with_system.isChecked()
        if self._set_auto_start(enable):
            if enable:
                self.statusBar.showMessage("已設定開機啟動", 3000)
            else:
                self.statusBar.showMessage("已取消開機啟動", 3000)
        else:
            # 設定失敗，恢復checkbox狀態
            self.start_with_system.blockSignals(True)
            self.start_with_system.setChecked(not enable)
            self.start_with_system.blockSignals(False)

    # ------------- 設定檔 讀寫 與 UI 還原 -------------
    def _update_last_index_label(self, ts):
        if ts:
            self.last_index_label.setText(ts)
        else:
            self.last_index_label.setText("尚未建立")

    def load_settings_into_ui(self):
        try:
            if not os.path.exists(self.settings_path):
                return
            with open(self.settings_path, "r", encoding="utf-8") as f:
                s = json.load(f)
        except Exception as e:
            print(f"讀取設定檔失敗: {e}")
            return

        dirs = s.get("index_directories", [])
        self.directory_list.setPlainText("\n".join(dirs))
        self.exclude_patterns_input.setText(", ".join(s.get("exclude_patterns", [])))

        self.index_threads.setValue(s.get("index_threads", 4))
        self.index_batch_size.setValue(s.get("index_batch_size", 10000))

        self.index_content.setChecked(s.get("index_content", True))
        self.content_types.setText(s.get("content_types", self.content_types.text()))

        self.sort_by.setCurrentIndex(s.get("search_sort_by_index", 0))
        self.sort_desc.setChecked(s.get("search_sort_desc", False))
        self.search_threads.setValue(s.get("search_threads", 4))
        self.max_memory.setValue(s.get("max_memory", 1024))

        dbp = s.get("db_path")
        if dbp:
            self.db_path = dbp
            self.db_path_display.setText(dbp)

        # 系統匣設定
        self.minimize_to_tray.setChecked(s.get("minimize_to_tray", False))
        self.close_to_tray_no_ask = bool(s.get("close_to_tray_no_ask", False))
        self.no_ask_close_to_tray.setChecked(self.close_to_tray_no_ask)

        # 系統啟動設定
        # 檢查實際的註冊表狀態，而不是依賴設定檔
        self.start_with_system.blockSignals(True)
        self.start_with_system.setChecked(self._is_auto_start_enabled())
        self.start_with_system.blockSignals(False)

        self._update_last_index_label(s.get("last_index_time"))

        # 載入自動更新設定
        auto_update_index = s.get("auto_update", 0)
        if 0 <= auto_update_index < self.auto_update.count():
            self.auto_update.setCurrentIndex(auto_update_index)

        # 載入自定義間隔設定
        custom_interval = s.get("custom_interval_hours", 6)
        self.custom_interval_hours.setValue(custom_interval)

        # 載入每日執行時間設定
        from PyQt5.QtCore import QTime

        daily_time_str = s.get("daily_update_time", "02:00")
        try:
            hour, minute = map(int, daily_time_str.split(":"))
            self.daily_update_time.setTime(QTime(hour, minute))
        except (ValueError, TypeError):
            self.daily_update_time.setTime(QTime(2, 0))

        # 初始化控件可見性
        self._on_auto_update_changed()

        self.search_directories = dirs
        self.exclude_patterns = s.get("exclude_patterns", [])

        # 載入進階排除規則
        advanced_rules = s.get("advanced_exclude_rules", {})
        if advanced_rules:
            self.advanced_exclude_rules.from_dict(advanced_rules)

        self.refresh_indexed_targets_list()
        self.refresh_index_tree()

        if hasattr(self, "theme_toggle_btn"):
            self.theme_toggle_btn.setText(
                "切換為淺色模式" if self.dark_mode else "切換為深色模式"
            )

    def save_settings_to_disk(self):
        old = {}
        try:
            if os.path.exists(self.settings_path):
                with open(self.settings_path, "r", encoding="utf-8") as f:
                    old = json.load(f)
        except Exception:
            pass

        s = {
            "db_path": self.db_path,
            "index_directories": [
                d for d in self.directory_list.toPlainText().split("\n") if d.strip()
            ],
            "exclude_patterns": [
                p.strip()
                for p in self.exclude_patterns_input.text().split(",")
                if p.strip()
            ],
            "index_threads": self.index_threads.value(),
            "index_batch_size": self.index_batch_size.value(),
            "index_content": self.index_content.isChecked(),
            "content_types": self.content_types.text(),
            "search_sort_by_index": self.sort_by.currentIndex(),
            "search_sort_desc": self.sort_desc.isChecked(),
            "search_threads": self.search_threads.value(),
            "max_memory": self.max_memory.value(),
            "last_index_time": old.get("last_index_time"),
            "minimize_to_tray": self.minimize_to_tray.isChecked(),
            "close_to_tray_no_ask": self.no_ask_close_to_tray.isChecked()
            or self.close_to_tray_no_ask,
            "auto_update": self.auto_update.currentIndex(),
            "custom_interval_hours": self.custom_interval_hours.value(),
            "daily_update_time": self.daily_update_time.time().toString("HH:mm"),
            "advanced_exclude_rules": self.advanced_exclude_rules.to_dict(),
            "dark_mode": getattr(self, "dark_mode", False),
        }
        try:
            with open(self.settings_path, "w", encoding="utf-8") as f:
                json.dump(s, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"寫入設定檔失敗: {e}")

    def save_settings(self):
        self.save_settings_to_disk()
        # 重新設定自動更新定時器
        self._setup_auto_update_timer()
        QMessageBox.information(self, "設定已保存", "您的設定已成功保存")

    def _init_auto_update_timer(self):
        """初始化自動更新定時器"""
        self.auto_update_timer = QTimer()
        self.auto_update_timer.timeout.connect(self._auto_update_index)
        self._setup_auto_update_timer()

    def _setup_auto_update_timer(self):
        """設定自動更新定時器"""
        if not hasattr(self, "auto_update_timer"):
            return

        self.auto_update_timer.stop()

        auto_update_option = self.auto_update.currentIndex()

        if auto_update_option == 1:  # 每天
            # 計算到下次執行時間的間隔
            self._setup_daily_timer()
            print(
                f"自動更新已設定為每天 {self.daily_update_time.time().toString('HH:mm')} 執行"
            )
        elif auto_update_option == 2:  # 每週
            # 每7天執行一次，在指定時間
            self._setup_weekly_timer()
            print(
                f"自動更新已設定為每週 {self.daily_update_time.time().toString('HH:mm')} 執行"
            )
        elif auto_update_option == 3:  # 自定義間隔
            # 使用自定義間隔時間
            interval_ms = self.custom_interval_hours.value() * 60 * 60 * 1000
            self.auto_update_timer.start(interval_ms)
            print(
                f"自動更新已設定為每 {self.custom_interval_hours.value()} 小時執行一次"
            )
        elif auto_update_option == 4:  # 手動
            # 手動模式不啟動定時器，但可以通過其他方式觸發
            print("自動更新設定為手動模式")
        else:  # 不自動更新
            print("自動更新已停用")

        # 更新狀態顯示
        if hasattr(self, "auto_update_status"):
            self._update_auto_update_status()

    def _auto_update_index(self):
        """自動更新索引"""
        # 檢查是否有正在進行的索引操作
        if hasattr(self, "indexer") and self.indexer.isRunning():
            print("索引操作正在進行中，跳過自動更新")
            return

        # 檢查是否有索引目錄
        directories = [
            d.strip()
            for d in self.directory_list.toPlainText().split("\n")
            if d.strip()
        ]
        if not directories:
            print("沒有設定索引目錄，跳過自動更新")
            return

        print(
            f"開始自動更新索引 - {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        )

        # 使用增量更新模式
        self.start_indexing_update()

    def _setup_daily_timer(self):
        """設定每日定時器"""
        from PyQt5.QtCore import QDateTime, QTime
        import datetime

        now = QDateTime.currentDateTime()
        target_time = self.daily_update_time.time()

        # 計算今天的目標時間
        today_target = QDateTime(now.date(), target_time)

        # 如果今天的目標時間已過，設定為明天
        if today_target <= now:
            today_target = today_target.addDays(1)

        # 計算到目標時間的毫秒數
        ms_to_target = now.msecsTo(today_target)

        # 設定單次定時器到目標時間
        if hasattr(self, "daily_timer"):
            self.daily_timer.stop()
        else:
            self.daily_timer = QTimer()
            self.daily_timer.timeout.connect(self._daily_timer_triggered)

        self.daily_timer.setSingleShot(True)
        self.daily_timer.start(ms_to_target)

    def _setup_weekly_timer(self):
        """設定每週定時器"""
        from PyQt5.QtCore import QDateTime, QTime
        import datetime

        now = QDateTime.currentDateTime()
        target_time = self.daily_update_time.time()

        # 計算今天的目標時間
        today_target = QDateTime(now.date(), target_time)

        # 如果今天的目標時間已過，設定為下週同一天
        if today_target <= now:
            # 計算下週同一天的時間
            next_week_target = today_target.addDays(7)
        else:
            # 如果今天的時間還沒到，設定為今天
            next_week_target = today_target

        # 計算到目標時間的毫秒數
        ms_to_target = now.msecsTo(next_week_target)

        # 設定單次定時器到目標時間
        if hasattr(self, "weekly_timer"):
            self.weekly_timer.stop()
        else:
            self.weekly_timer = QTimer()
            self.weekly_timer.timeout.connect(self._weekly_timer_triggered)

        self.weekly_timer.setSingleShot(True)
        self.weekly_timer.start(ms_to_target)

        print(
            f"每週定時器已設定，下次執行時間: {next_week_target.toString('yyyy-MM-dd HH:mm:ss')}"
        )

    def _daily_timer_triggered(self):
        """每日定時器觸發"""
        self._auto_update_index()
        # 重新設定明天的定時器
        self._setup_daily_timer()

    def _weekly_timer_triggered(self):
        """每週定時器觸發"""
        self._auto_update_index()
        # 重新設定下週的定時器
        self._setup_weekly_timer()

        # 在狀態列顯示自動更新訊息
        self.statusBar.showMessage("自動更新索引中...", 3000)

    def _on_auto_update_changed(self):
        """自動更新選項變更時的處理"""
        auto_update_option = self.auto_update.currentIndex()

        # 控制時間設定容器的可見性
        if auto_update_option in [1, 2]:  # 每天或每週
            self.daily_time_widget.setVisible(True)
            self.custom_interval_widget.setVisible(False)
        elif auto_update_option == 3:  # 自定義間隔
            self.daily_time_widget.setVisible(False)
            self.custom_interval_widget.setVisible(True)
        else:  # 不自動更新或手動
            self.daily_time_widget.setVisible(False)
            self.custom_interval_widget.setVisible(False)

        self._setup_auto_update_timer()
        self._update_auto_update_status()

    def _manual_trigger_auto_update(self):
        """手動觸發自動更新"""
        if hasattr(self, "indexer") and self.indexer.isRunning():
            QMessageBox.warning(self, "無法執行", "索引操作正在進行中，請稍後再試。")
            return

        directories = [
            d.strip()
            for d in self.directory_list.toPlainText().split("\n")
            if d.strip()
        ]
        if not directories:
            QMessageBox.warning(self, "無法執行", "請先設定索引目錄。")
            return

        reply = QMessageBox.question(
            self,
            "確認執行",
            "確定要立即執行自動更新索引嗎？\n這將使用增量更新模式。",
            QMessageBox.Yes | QMessageBox.No,
        )
        if reply == QMessageBox.Yes:
            self._auto_update_index()

    def _update_auto_update_status(self):
        """更新自動更新狀態顯示"""
        auto_update_option = self.auto_update.currentIndex()

        if auto_update_option == 0:
            status_text = "自動更新狀態: 已停用"
        elif auto_update_option == 1:
            time_str = self.daily_update_time.time().toString("HH:mm")
            status_text = f"自動更新狀態: 每天 {time_str} 執行"
        elif auto_update_option == 2:
            time_str = self.daily_update_time.time().toString("HH:mm")
            status_text = f"自動更新狀態: 每週 {time_str} 執行"
        elif auto_update_option == 3:
            interval = self.custom_interval_hours.value()
            status_text = f"自動更新狀態: 每 {interval} 小時執行一次"
        elif auto_update_option == 4:
            status_text = "自動更新狀態: 手動模式（可隨時觸發）"
        else:
            status_text = "自動更新狀態: 未知"

        self.auto_update_status.setText(status_text)

    def _update_next_run_preview(self):
        """更新下次執行時間預覽"""
        if not hasattr(self, "next_run_preview"):
            return

        auto_update_option = self.auto_update.currentIndex()

        if auto_update_option == 0:  # 不自動更新
            self.next_run_preview.setText("下次執行時間: 未排程")
            return

        if auto_update_option == 4:  # 手動模式
            self.next_run_preview.setText("下次執行時間: 等待手動觸發")
            return

        try:
            from PyQt5.QtCore import QDateTime, QTime

            now = QDateTime.currentDateTime()
            target_time = self.daily_update_time.time()

            if auto_update_option == 1:  # 每天
                today_target = QDateTime(now.date(), target_time)
                if today_target <= now:
                    today_target = today_target.addDays(1)

                seconds_to_target = now.secsTo(today_target)
                hours = seconds_to_target // 3600
                minutes = (seconds_to_target % 3600) // 60
                seconds = seconds_to_target % 60

                time_str = today_target.toString("yyyy-MM-dd HH:mm:ss")
                countdown_str = f"{hours}小時 {minutes}分 {seconds}秒"
                self.next_run_preview.setText(
                    f"下次執行時間: {time_str} (倒數: {countdown_str})"
                )

            elif auto_update_option == 2:  # 每週
                today_target = QDateTime(now.date(), target_time)
                if today_target <= now:
                    next_week_target = today_target.addDays(7)
                else:
                    next_week_target = today_target

                seconds_to_target = now.secsTo(next_week_target)
                days = seconds_to_target // 86400
                hours = (seconds_to_target % 86400) // 3600
                minutes = (seconds_to_target % 3600) // 60

                time_str = next_week_target.toString("yyyy-MM-dd HH:mm:ss")
                countdown_str = f"{days}天 {hours}小時 {minutes}分"
                self.next_run_preview.setText(
                    f"下次執行時間: {time_str} (倒數: {countdown_str})"
                )

            elif auto_update_option == 3:  # 自定義間隔
                interval_hours = self.custom_interval_hours.value()

                if hasattr(self, "last_index_time_str"):
                    try:
                        last_time = QDateTime.fromString(
                            self.last_index_time_str, "yyyy-MM-dd HH:mm:ss"
                        )
                        if last_time.isValid():
                            next_time = last_time.addSecs(interval_hours * 3600)

                            if next_time <= now:
                                next_time = now.addSecs(interval_hours * 3600)

                            seconds_to_target = now.secsTo(next_time)
                            hours = seconds_to_target // 3600
                            minutes = (seconds_to_target % 3600) // 60
                            seconds = seconds_to_target % 60

                            time_str = next_time.toString("yyyy-MM-dd HH:mm:ss")
                            countdown_str = f"{hours}小時 {minutes}分 {seconds}秒"
                            self.next_run_preview.setText(
                                f"下次執行時間: {time_str} (倒數: {countdown_str})"
                            )
                        else:
                            self.next_run_preview.setText("下次執行時間: 計算中...")
                    except Exception:
                        self.next_run_preview.setText("下次執行時間: 計算中...")
                else:
                    self.next_run_preview.setText("下次執行時間: 尚未執行過索引")

        except Exception as e:
            print(f"更新下次執行時間預覽時出錯: {e}")
            self.next_run_preview.setText("下次執行時間: 計算中...")

    # ---------------- 索引管理動作 ----------------
    def change_db_path(self):
        new_path = QFileDialog.getSaveFileName(
            self,
            "選擇資料庫儲存位置",
            self.db_path,
            "SQLite 資料庫 (*.db);;所有檔案 (*.*)",
        )[0]

        if new_path:
            try:
                test_conn = sqlite3.connect(new_path)
                test_conn.execute("CREATE TABLE IF NOT EXISTS test_write (id INTEGER)")
                test_conn.execute("DROP TABLE test_write")
                test_conn.commit()
                test_conn.close()

                self.db_path = new_path
                self.db_path_display.setText(new_path)
                self.save_settings_to_disk()
                QMessageBox.information(self, "成功", "資料庫路徑已更新")

                self.refresh_indexed_targets_list()
                self.refresh_index_tree()  # 新增：更新索引樹

            except Exception as e:
                QMessageBox.warning(
                    self, "錯誤", f"無法在選擇的位置創建資料庫：{str(e)}"
                )

    def add_directory(self):
        directory = QFileDialog.getExistingDirectory(self, "選擇目錄")
        if directory:
            current_text = self.directory_list.toPlainText()
            if current_text:
                self.directory_list.setPlainText(f"{current_text}\n{directory}")
            else:
                self.directory_list.setPlainText(directory)
            self.save_settings_to_disk()

    def add_network_path(self):
        network_path, ok = QInputDialog.getText(
            self,
            "輸入網路路徑",
            "請輸入網路路徑 (UNC格式如 \\\\server\\share 或映射磁碟機路徑)",
        )
        if ok and network_path:
            current_text = self.directory_list.toPlainText()
            if current_text:
                self.directory_list.setPlainText(f"{current_text}\n{network_path}")
            else:
                self.directory_list.setPlainText(network_path)
            self.save_settings_to_disk()

    def clear_directories(self):
        self.directory_list.clear()
        self.save_settings_to_disk()

    def start_indexing_full(self):
        directories = self._collect_directories_or_warn()
        if not directories:
            return

        self.show_summary_on_complete = False

        exclude_patterns = [
            p.strip()
            for p in self.exclude_patterns_input.text().split(",")
            if p.strip()
        ]

        self.search_directories = directories
        self.exclude_patterns = exclude_patterns

        if not self._check_database_access():
            QMessageBox.critical(
                self,
                "資料庫錯誤",
                f"無法存取資料庫檔案：{self.db_path}\n"
                f"請確保該位置有寫入權限，或選擇其他位置。",
            )
            return

        self.save_settings_to_disk()

        self._prepare_indexing_ui("正在索引（全量）...")

        batch_size = self.index_batch_size.value()
        index_content_flag = self.index_content.isChecked()
        content_types_input = self.content_types.text()
        max_workers = self.index_threads.value()

        self.indexer = FileIndexer(
            self.db_path,
            directories,
            exclude_patterns,
            index_content=index_content_flag,
            content_types=content_types_input,
            max_workers=max_workers,
            batch_size=batch_size,
            mode="full",
            advanced_exclude_rules=self.advanced_exclude_rules,
        )
        self.indexer.progress_update.connect(self.update_index_progress)
        self.indexer.indexing_complete.connect(self.indexing_complete)
        self.indexer.start()

    def start_indexing_update(self):
        self.show_summary_on_complete = False
        update_targets = self._get_selected_update_targets()
        if not update_targets:
            directories = [
                d.strip()
                for d in self.directory_list.toPlainText().split("\n")
                if d.strip()
            ]
            if not directories:
                QMessageBox.warning(
                    self, "無法更新", "請先在索引目錄列表中添加目錄，或選擇已索引目標。"
                )
                return
            db_targets = self._get_indexed_targets_from_db()
            update_targets = [
                d for d in directories if any(p.startswith(d) for p in db_targets)
            ]
            if not update_targets:
                QMessageBox.information(
                    self,
                    "無已索引目標",
                    "資料庫中找不到與目前目錄清單相符的已索引目標。",
                )
                return

        if not self._check_database_access():
            QMessageBox.critical(
                self,
                "資料庫錯誤",
                f"無法存取資料庫檔案：{self.db_path}\n"
                f"請確保該位置有寫入權限，或選擇其他位置。",
            )
            return

        exclude_patterns = [
            p.strip()
            for p in self.exclude_patterns_input.text().split(",")
            if p.strip()
        ]

        self._prepare_indexing_ui("正在更新（增量）...")

        batch_size = self.index_batch_size.value()
        index_content_flag = self.index_content.isChecked()
        content_types_input = self.content_types.text()
        max_workers = self.index_threads.value()

        self.indexer = FileIndexer(
            self.db_path,
            update_targets,
            exclude_patterns,
            index_content=index_content_flag,
            content_types=content_types_input,
            max_workers=max_workers,
            batch_size=batch_size,
            mode="update",
            update_targets=update_targets,
            advanced_exclude_rules=self.advanced_exclude_rules,
        )
        self.indexer.progress_update.connect(self.update_index_progress)
        self.indexer.indexing_complete.connect(self.indexing_complete)
        self.indexer.start()

    def _prepare_indexing_ui(self, status_text):
        self.start_index_button.setEnabled(False)
        self.update_index_button.setEnabled(False)
        self.stop_index_button.setEnabled(True)
        self.tabs.setTabEnabled(0, False)

        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
        self.statusBar.showMessage(status_text)
        self.index_progress_info.setText("準備中...")

    def _collect_directories_or_warn(self):
        directories = self.directory_list.toPlainText().strip().split("\n")
        directories = [d for d in directories if d.strip()]
        if not directories:
            QMessageBox.warning(self, "無法索引", "請先添加至少一個目錄進行索引")
            return None
        return directories

    def _get_selected_update_targets(self):
        items = self.update_targets_list.selectedItems()
        return [i.text() for i in items]

    def _get_indexed_targets_from_db(self):
        targets = set()
        if not os.path.exists(self.db_path):
            return []
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute("SELECT path FROM files LIMIT 10000")
            rows = cursor.fetchall()
            conn.close()
        except Exception:
            rows = []

        ui_dirs = [
            d.strip()
            for d in self.directory_list.toPlainText().split("\n")
            if d.strip()
        ]
        for (p,) in rows:
            for d in ui_dirs:
                if p and os.path.abspath(p).startswith(os.path.abspath(d)):
                    targets.add(d)
        return sorted(targets)

    def refresh_indexed_targets_list(self):
        self.update_targets_list.clear()
        targets = self._get_indexed_targets_from_db()
        for t in targets:
            self.update_targets_list.addItem(QListWidgetItem(t))

    def _check_database_access(self):
        try:
            db_dir = os.path.dirname(self.db_path)
            os.makedirs(db_dir, exist_ok=True)

            conn = sqlite3.connect(self.db_path)
            conn.execute("CREATE TABLE IF NOT EXISTS test_access (id INTEGER)")
            conn.execute("INSERT INTO test_access (id) VALUES (1)")
            conn.execute("DELETE FROM test_access")
            conn.execute("DROP TABLE test_access")
            conn.commit()
            conn.close()
            return True
        except Exception as e:
            print(f"資料庫存取測試失敗: {e}")
            return False

    def stop_indexing(self):
        if hasattr(self, "indexer") and self.indexer.isRunning():
            self.indexer.stop()
            self.indexer.wait()
            stats = {
                "success": False,
                "mode": "unknown",
                "files_added": 0,
                "files_updated": 0,
                "files_deleted": 0,
                "total_files": 0,
                "duration": 0,
                "error_count": 0,
                "error_messages": [],
            }
            self.indexing_complete(stats)

    def update_index_progress(self, current, total):
        progress = int(current / max(total, 1) * 100)
        self.progress_bar.setValue(progress)
        self.statusBar.showMessage(f"索引中... ({current}/{total})")
        self.index_progress_info.setText(
            f"已處理: {current} | 估計總數: {total} | 進度: {progress}%"
        )

    def indexing_complete(self, stats):
        self.progress_bar.setVisible(False)
        self.start_index_button.setEnabled(True)
        self.update_index_button.setEnabled(True)
        self.stop_index_button.setEnabled(False)
        self.tabs.setTabEnabled(0, True)

        success = stats.get("success", False)

        if success:
            self.refresh_index_stats()
            self.refresh_index_tree()
            if hasattr(self, "statistics_widget"):
                self.statistics_widget.refresh_statistics()
            self.statusBar.showMessage("操作完成")
            self.index_progress_info.setText("完成")

            # 儲存更新歷史記錄
            status = "success" if stats.get("error_count", 0) == 0 else "partial"
            self.update_history_manager.add_record(
                mode=stats.get("mode", "full"),
                files_added=stats.get("files_added", 0),
                files_updated=stats.get("files_updated", 0),
                files_deleted=stats.get("files_deleted", 0),
                total_files=stats.get("total_files", 0),
                duration_seconds=stats.get("duration", 0),
                status=status,
                target_paths=self.search_directories
                if stats.get("mode") == "full"
                else list(
                    self.indexer.update_targets
                    if hasattr(self, "indexer")
                    and hasattr(self.indexer, "update_targets")
                    else []
                ),
                error_count=stats.get("error_count", 0),
                error_messages=stats.get("error_messages", []),
            )

            # 顯示更新摘要對話框
            if (
                hasattr(self, "show_summary_on_complete")
                and self.show_summary_on_complete
            ):
                self._show_update_summary(stats)

            last_ts = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            try:
                s = {}
                if os.path.exists(self.settings_path):
                    with open(self.settings_path, "r", encoding="utf-8") as f:
                        s = json.load(f)
                s["last_index_time"] = last_ts
                with open(self.settings_path, "w", encoding="utf-8") as f:
                    json.dump(s, f, ensure_ascii=False, indent=2)
            except Exception as e:
                print(f"更新最後索引時間失敗: {e}")
            self._update_last_index_label(last_ts)

            self.refresh_indexed_targets_list()
        else:
            self.statusBar.showMessage("操作被中斷")
            self.index_progress_info.setText("已中斷")

        self.save_settings_to_disk()

    def _show_update_summary(self, stats):
        """顯示更新摘要對話框"""
        dialog = UpdateSummaryDialog(stats, self)
        dialog.exec_()

    def show_update_history(self):
        """顯示更新歷史記錄對話框"""
        dialog = UpdateHistoryDialog(self.update_history_manager, self)
        dialog.exec_()

    def _show_advanced_exclude_dialog(self):
        """顯示進階排除規則設定對話框"""
        dialog = AdvancedExcludeDialog(self.advanced_exclude_rules, self)
        if dialog.exec_() == QDialog.Accepted:
            self.advanced_exclude_rules = dialog.get_rules()
            self.save_settings_to_disk()
            self.statusBar.showMessage("進階排除規則已更新", 3000)

    def _show_targets_context_menu(self, position):
        """顯示已索引目標右鍵選單"""
        item = self.update_targets_list.itemAt(position)
        if not item:
            return

        selected_items = self.update_targets_list.selectedItems()
        selected_paths = [item.text() for item in selected_items]

        menu = QMenu()

        if len(selected_paths) == 1:
            quick_update_action = QAction("立即更新此目錄", self)
            quick_update_action.triggered.connect(
                lambda: self._quick_update_target(selected_paths[0])
            )
            menu.addAction(quick_update_action)

            view_stats_action = QAction("查看目錄統計", self)
            view_stats_action.triggered.connect(
                lambda: self._show_target_stats(selected_paths[0])
            )
            menu.addAction(view_stats_action)
        else:
            quick_update_all_action = QAction(
                f"立即更新 {len(selected_paths)} 個目錄", self
            )
            quick_update_all_action.triggered.connect(
                lambda: self._quick_update_targets(selected_paths)
            )
            menu.addAction(quick_update_all_action)

        menu.addSeparator()

        remove_action = QAction("從列表移除", self)
        remove_action.triggered.connect(
            lambda: self._remove_targets_from_list(selected_paths)
        )
        menu.addAction(remove_action)

        menu.exec_(self.update_targets_list.viewport().mapToGlobal(position))

    def _quick_update_target(self, target_path):
        """快速更新單一目錄"""
        if hasattr(self, "indexer") and self.indexer.isRunning():
            QMessageBox.warning(self, "無法執行", "索引操作正在進行中，請稍後再試。")
            return

        reply = QMessageBox.question(
            self,
            "確認更新",
            f"確定要立即更新目錄：\n{target_path}\n嗎？",
            QMessageBox.Yes | QMessageBox.No,
        )

        if reply != QMessageBox.Yes:
            return

        self.show_summary_on_complete = True
        self._start_indexing_for_targets([target_path])

    def _quick_update_targets(self, target_paths):
        """快速更新多個目錄"""
        if hasattr(self, "indexer") and self.indexer.isRunning():
            QMessageBox.warning(self, "無法執行", "索引操作正在進行中，請稍後再試。")
            return

        reply = QMessageBox.question(
            self,
            "確認更新",
            f"確定要立即更新 {len(target_paths)} 個目錄嗎？",
            QMessageBox.Yes | QMessageBox.No,
        )

        if reply != QMessageBox.Yes:
            return

        self.show_summary_on_complete = True
        self._start_indexing_for_targets(target_paths)

    def _start_indexing_for_targets(self, targets):
        """為指定目標開始索引"""
        if not self._check_database_access():
            QMessageBox.critical(
                self,
                "資料庫錯誤",
                f"無法存取資料庫檔案：{self.db_path}\n"
                f"請確保該位置有寫入權限，或選擇其他位置。",
            )
            return

        exclude_patterns = [
            p.strip()
            for p in self.exclude_patterns_input.text().split(",")
            if p.strip()
        ]

        self._prepare_indexing_ui("正在更新（增量）...")

        batch_size = self.index_batch_size.value()
        index_content_flag = self.index_content.isChecked()
        content_types_input = self.content_types.text()
        max_workers = self.index_threads.value()

        self.indexer = FileIndexer(
            self.db_path,
            targets,
            exclude_patterns,
            index_content=index_content_flag,
            content_types=content_types_input,
            max_workers=max_workers,
            batch_size=batch_size,
            mode="update",
            update_targets=targets,
            advanced_exclude_rules=self.advanced_exclude_rules,
        )
        self.indexer.progress_update.connect(self.update_index_progress)
        self.indexer.indexing_complete.connect(self.indexing_complete)
        self.indexer.start()

    def _show_target_stats(self, target_path):
        """顯示特定目錄的統計資訊"""
        if not os.path.exists(self.db_path):
            QMessageBox.information(self, "提示", "資料庫尚未建立。")
            return

        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute(
                "SELECT COUNT(*) FROM files WHERE path LIKE ?", (f"{target_path}%",)
            )
            file_count = cursor.fetchone()[0]

            cursor.execute(
                "SELECT SUM(size) FROM files WHERE path LIKE ?", (f"{target_path}%",)
            )
            total_size = cursor.fetchone()[0] or 0

            cursor.execute(
                """
                SELECT ext, COUNT(*) as cnt
                FROM files
                WHERE path LIKE ?
                GROUP BY ext
                ORDER BY cnt DESC
                LIMIT 10
            """,
                (f"{target_path}%",),
            )
            type_stats = cursor.fetchall()

            conn.close()

            size_mb = total_size / (1024 * 1024)
            size_str = (
                f"{size_mb:.2f} MB" if size_mb < 1024 else f"{size_mb / 1024:.2f} GB"
            )

            message = f"目錄: {target_path}\n\n"
            message += f"檔案總數: {file_count}\n"
            message += f"總大小: {size_str}\n\n"

            if type_stats:
                message += "前10種檔案類型:\n"
                for ext, count in type_stats:
                    ext_display = ext if ext else "(無副檔名)"
                    message += f"  {ext_display}: {count}\n"

            QMessageBox.information(self, "目錄統計", message)

        except Exception as e:
            QMessageBox.warning(self, "錯誤", f"無法讀取統計資訊: {str(e)}")

    def _remove_targets_from_list(self, target_paths):
        """從列表移除目錄"""
        reply = QMessageBox.question(
            self,
            "確認移除",
            f"確定要從列表移除 {len(target_paths)} 個目錄嗎？\n"
            f"(這不會刪除資料庫中的索引)",
            QMessageBox.Yes | QMessageBox.No,
        )

        if reply != QMessageBox.Yes:
            return

        current_text = self.directory_list.toPlainText()
        current_dirs = [d.strip() for d in current_text.split("\n") if d.strip()]

        new_dirs = [d for d in current_dirs if d not in target_paths]

        self.directory_list.setPlainText("\n".join(new_dirs))
        self.save_settings_to_disk()

        self.refresh_indexed_targets_list()
        self.statusBar.showMessage(f"已移除 {len(target_paths)} 個目錄", 3000)

    def refresh_index_stats(self):
        try:
            if os.path.exists(self.settings_path):
                with open(self.settings_path, "r", encoding="utf-8") as f:
                    s = json.load(f)
                    self._update_last_index_label(s.get("last_index_time"))
        except Exception:
            pass

        self.refresh_indexed_targets_list()

        if not os.path.exists(self.db_path):
            self.index_stats.setText("尚未建立索引")
            return

        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute("SELECT COUNT(*) FROM files")
            file_count = cursor.fetchone()[0]

            cursor.execute("SELECT COUNT(*) FROM files WHERE content_indexed = 1")
            content_indexed = cursor.fetchone()[0]

            cursor.execute(
                "SELECT ext, COUNT(*) FROM files GROUP BY ext ORDER BY COUNT(*) DESC LIMIT 10"
            )
            type_stats = cursor.fetchall()

            cursor.execute("PRAGMA page_count")
            page_count = cursor.fetchone()[0]
            cursor.execute("PRAGMA page_size")
            page_size = cursor.fetchone()[0]
            db_size = page_count * page_size / (1024 * 1024)

            conn.close()

            stats_text = (
                f"共索引 {file_count} 個檔案，其中 {content_indexed} 個檔案索引了內容\n"
            )
            stats_text += f"索引資料庫大小: {db_size:.2f} MB\n\n"

            if type_stats:
                stats_text += "前10種檔案類型統計:\n"
                for ext, count in type_stats:
                    ext_display = ext if ext else "(無副檔名)"
                    stats_text += f"  {ext_display}: {count} 個檔案\n"

            self.index_stats.setText(stats_text)

        except Exception as e:
            self.index_stats.setText(f"讀取索引統計時出錯: {str(e)}")

    def _init_search_history(self):
        """初始化搜尋歷史功能"""
        self._refresh_search_history()
        self.search_history_combo.currentIndexChanged.connect(self._load_search_history)
        self.clear_history_button.clicked.connect(self._clear_search_history)

    def _refresh_search_history(self):
        """重新載入搜尋歷史下拉選單"""
        self.search_history_combo.clear()
        histories = self.search_history_manager.get_unique_searches(limit=20)
        if not histories:
            self.search_history_combo.addItem("(無歷史記錄)")
            return
        self.search_history_combo.addItem("(選擇歷史記錄)")
        for record in histories:
            filename = record[0] or ""
            content = record[1] or ""
            boolean_expr = record[4] or ""
            display_text = ""
            if boolean_expr:
                display_text = f"[布林] {boolean_expr}"
            elif filename and content:
                display_text = f"{filename} / {content}"
            elif filename:
                display_text = filename
            elif content:
                display_text = f"[內容] {content}"
            if len(display_text) > 50:
                display_text = display_text[:50] + "..."
            self.search_history_combo.addItem(display_text, record)

    def _load_search_history(self, index):
        """從歷史記錄載入搜尋條件"""
        if index <= 0:
            return
        record = self.search_history_combo.itemData(index)
        if not record:
            return
        filename = record[0] or ""
        content = record[1] or ""
        file_type = record[2] or ""
        boolean_expr = record[4] or ""
        path_filter = record[5] or ""
        use_regex_filename = record[6]
        use_regex_content = record[7]
        use_boolean = record[8]
        self.filename_input.setText(filename)
        self.content_input.setText(content)
        self.file_type_input.setText(file_type)
        self.boolean_input.setText(boolean_expr)
        self.path_filter.setText(path_filter)
        self.use_regex_filename.setChecked(bool(use_regex_filename))
        self.use_regex_content.setChecked(bool(use_regex_content))
        self.use_boolean.setChecked(bool(use_boolean))

    def _clear_search_history(self):
        """清除搜尋歷史記錄"""
        reply = QMessageBox.question(
            self,
            "確認清除",
            "確定要清除所有搜尋歷史記錄？",
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No
        )
        if reply == QMessageBox.Yes:
            count = self.search_history_manager.clear_all_records()
            self._refresh_search_history()
            self.statusBar.showMessage(f"已清除 {count} 搜尋歷史記錄")

    def _init_realtime_search(self):
        """初始化即時搜尋功能"""
        self.filename_input.textChanged.connect(self._on_search_input_changed)
        self.content_input.textChanged.connect(self._on_search_input_changed)
        self.boolean_input.textChanged.connect(self._on_search_input_changed)
        self.path_filter.textChanged.connect(self._on_search_input_changed)
        self.use_realtime_search.stateChanged.connect(self._on_realtime_search_toggled)

    def _init_file_watcher(self):
        """初始化檔案系統即時監控"""
        self.file_watcher = IndexWatchManager(
            self.db_path,
            set(self.content_types.text().replace(" ", "").split(",")),
            self.advanced_exclude_rules
        )

    def start_file_watching(self):
        """開始檔案監控"""
        dirs = self.directory_list.toPlainText().strip().split("\n")
        dirs = [d.strip() for d in dirs if d.strip() and os.path.isdir(d.strip())]
        if dirs:
            self.file_watcher.start_watching(dirs)
            self.enable_watcher_btn.setEnabled(False)
            self.disable_watcher_btn.setEnabled(True)
            self.statusBar.showMessage(f"已開始監控 {len(dirs)} 個目錄")
        else:
            QMessageBox.warning(self, "無法啟動監控", "沒有有效的索引目錄可供監控")

    def stop_file_watching(self):
        """停止檔案監控"""
        self.file_watcher.stop_watching()
        self.enable_watcher_btn.setEnabled(True)
        self.disable_watcher_btn.setEnabled(False)
        self.statusBar.showMessage("已停止檔案監控")

    def _on_realtime_search_toggled(self, state):
        """即時搜尋開關切換"""
        self.realtime_search_enabled = bool(state)

    def _on_filename_text_changed(self, text):
        """檔案名稱輸入變化，顯示建議"""
        text = text.strip()
        if len(text) < 2:
            if hasattr(self, "filename_suggestion"):
                self.filename_suggestion.hide()
            return
        suggestions = self._get_search_suggestions(text)
        if suggestions and hasattr(self, "filename_suggestion"):
            cursor_pos = self.filename_input.cursorRect().bottomRight()
            global_pos = self.filename_input.mapToGlobal(cursor_pos)
            self.filename_suggestion.show_suggestions(suggestions, global_pos)

    def _on_content_text_changed(self, text):
        """內容輸入變化，顯示建議"""
        text = text.strip()
        if len(text) < 2:
            if hasattr(self, "content_suggestion"):
                self.content_suggestion.hide()
            return
        suggestions = self._get_search_suggestions(text)
        if suggestions and hasattr(self, "content_suggestion"):
            cursor_pos = self.content_input.cursorRect().bottomRight()
            global_pos = self.content_input.mapToGlobal(cursor_pos)
            self.content_suggestion.show_suggestions(suggestions, global_pos)

    def _get_search_suggestions(self, prefix):
        """取得搜尋建議（歷史 + 模板）"""
        suggestions = []
        try:
            history_suggestions = self.search_history_manager.get_suggestions(prefix, limit=5)
            suggestions.extend(history_suggestions)
        except Exception:
            pass
        try:
            template_suggestions = self.search_template_manager.get_template_suggestions(prefix, limit=5)
            for name in template_suggestions:
                if name not in suggestions:
                    suggestions.append(name)
        except Exception:
            pass
        BOOLEAN_KEYWORDS = ["AND", "OR", "NOT", "NEAR", "REGEX", "LIKE"]
        for kw in BOOLEAN_KEYWORDS:
            if kw.lower().startswith(prefix.lower()) and kw not in suggestions:
                suggestions.append(kw)
        return suggestions[:10]

    def _on_filename_suggestion_selected(self, text):
        """選擇檔案名稱建議"""
        self.filename_input.setText(text)
        self.filename_input.setFocus()
        if hasattr(self, "filename_suggestion"):
            self.filename_suggestion.hide()

    def _on_content_suggestion_selected(self, text):
        """選擇內容建議"""
        self.content_input.setText(text)
        self.content_input.setFocus()
        if hasattr(self, "content_suggestion"):
            self.content_suggestion.hide()

    def _on_search_input_changed(self):
        """搜尋輸入變化時觸發 debounce"""
        if not self.realtime_search_enabled:
            return
        if hasattr(self, "searcher") and self.searcher.isRunning():
            return
        self.realtime_search_timer.start(300)

    def _do_realtime_search(self):
        """執行即時搜尋"""
        if not self.realtime_search_enabled:
            return
        if hasattr(self, "searcher") and self.searcher.isRunning():
            return
        filename = self.filename_input.text().strip()
        content = self.content_input.text().strip()
        boolean_expr = self.boolean_input.text().strip()
        path = self.path_filter.text().strip()
        if not (filename or content or boolean_expr or path):
            self.result_table.setRowCount(0)
            self.result_info.setText("即時搜尋：輸入條件以搜尋")
            return
        self._perform_search(save_history=False)

    # ---------------- 搜尋 ----------------
    def start_search(self):
        self.pagination_widget.reset()
        self._perform_search(save_history=True)

    def _perform_search(self, save_history=True):
        if not os.path.exists(self.db_path):
            QMessageBox.warning(self, "無法搜尋", "請先建立檔案索引")
            self.tabs.setCurrentIndex(1)
            return

        search_config = {
            "filename": self.filename_input.text(),
            "use_regex_filename": self.use_regex_filename.isChecked(),
            "use_fuzzy_filename": self.use_fuzzy_filename.isChecked(),
            "content": self.content_input.text(),
            "use_regex_content": self.use_regex_content.isChecked(),
            "use_fuzzy_content": self.use_fuzzy_content.isChecked(),
            "file_type": self.file_type_input.text(),
            "boolean_mode": self.use_boolean.isChecked(),
            "boolean_expr": self.boolean_input.text()
            if self.use_boolean.isChecked()
            else "",
            "path_filter": self.path_filter.text(),
            "sort_by": ["filename", "size", "modified", "path"][
                self.sort_by.currentIndex()
            ],
            "sort_desc": self.sort_desc.isChecked(),
            "save_history": save_history,
        }

        if self.min_size_input.value() > 0:
            search_config["min_size"] = self.min_size_input.value() * 1024

        if self.max_size_input.value() > 0:
            search_config["max_size"] = self.max_size_input.value() * 1024

        if self.use_date_filter.isChecked():
            search_config["date_after"] = self.date_after.date().toString(Qt.ISODate)
            search_config["date_before"] = self.date_before.date().toString(Qt.ISODate)

        self.search_button.setEnabled(False)
        self.stop_search_button.setEnabled(True)
        self.result_table.setRowCount(0)
        self.result_info.setText("正在搜尋...")
        self.content_preview.clear()

        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
        self.statusBar.showMessage("正在搜尋...")

        search_threads = self.search_threads.value()
        max_memory = self.max_memory.value()
        page = self.pagination_widget.current_page()
        page_size = self.pagination_widget.page_size()

        self.searcher = FileSearcher(
            self.db_path, search_config, search_threads, max_memory, page, page_size
        )
        self.searcher.search_complete.connect(self.search_complete)
        self.searcher.search_count.connect(self._on_search_count)
        self.searcher.search_progress.connect(self.update_search_progress)
        self.searcher.start()

    def stop_search(self):
        if hasattr(self, "searcher") and self.searcher.isRunning():
            self.searcher.stop()
            self.searcher.wait()

    def _on_page_changed(self, page):
        self._perform_search(save_history=False)

    def _on_search_count(self, total_count):
        self.pagination_widget.set_total_count(total_count)
        self._current_search_total = total_count

    def update_search_progress(self, current, total):
        progress = int(current / max(total, 1) * 100)
        self.progress_bar.setValue(progress)
        self.statusBar.showMessage(f"正在搜尋... ({current}/{total} 檔案)")

    def search_complete(self, results, search_time):
        self.progress_bar.setVisible(False)
        self.search_button.setEnabled(True)
        self.stop_search_button.setEnabled(False)

        keywords = []
        if hasattr(self, "searcher") and hasattr(self.searcher, "search_config"):
            config = self.searcher.search_config
            if config.get("filename"):
                keywords.append(config["filename"])
            if config.get("content"):
                keywords.append(config["content"])
            if config.get("boolean_expr"):
                terms = re.findall(r'\b(?!AND|OR|NOT|NEAR)\w+\b', config["boolean_expr"])
                keywords.extend(terms)

        if not hasattr(self, "highlight_delegate"):
            self.highlight_delegate = HighlightDelegate(keywords)
            self.result_table.setItemDelegateForColumn(0, self.highlight_delegate)
        else:
            self.highlight_delegate.set_keywords(keywords)

        self.result_table.setRowCount(len(results))

        for i, result in enumerate(results):
            self.result_table.setItem(i, 0, QTableWidgetItem(result["filename"]))
            normalized_path = os.path.normpath(result["path"])
            self.result_table.setItem(i, 1, QTableWidgetItem(normalized_path))

            size_kb = result["size"] / 1024
            if size_kb < 1024:
                size_str = f"{size_kb:.2f} KB"
            else:
                size_mb = size_kb / 1024
                if size_mb < 1024:
                    size_str = f"{size_mb:.2f} MB"
                else:
                    size_gb = size_mb / 1024
                    size_str = f"{size_gb:.2f} GB"

            self.result_table.setItem(i, 2, QTableWidgetItem(size_str))

            if isinstance(result["modified"], str):
                modified_date = result["modified"]
            else:
                try:
                    modified_date = datetime.datetime.fromtimestamp(
                        result["modified"]
                    ).strftime("%Y-%m-%d %H:%M:%S")
                except (ValueError, OSError, TypeError):
                    modified_date = str(result["modified"])

            self.result_table.setItem(i, 3, QTableWidgetItem(modified_date))
            self.result_table.setItem(i, 4, QTableWidgetItem(result["ext"]))

        total_count = getattr(self, "_current_search_total", len(results))
        current_page = self.pagination_widget.current_page()
        total_pages = max(1, (total_count + self.pagination_widget.page_size() - 1) // self.pagination_widget.page_size())

        self.result_info.setText(
            f" 第 {current_page}/{total_pages} 頁，本頁 {len(results)} 項，共 {total_count} 項 (耗時: {search_time:.3f}s)"
        )
        self.statusBar.showMessage(f"搜尋完成: 第 {current_page} 頁，{len(results)} 個結果")

        self.result_table.resizeColumnsToContents()

        if hasattr(self, "searcher") and hasattr(self.searcher, "search_config"):
            if self.searcher.search_config.get("save_history", False):
                self.search_history_manager.add_record(
                    self.searcher.search_config,
                    result_count=len(results),
                    search_duration=search_time
                )
                self._refresh_search_history()

    # ---------------- 右鍵選單與檔案操作 ----------------
    def open_file(self, index):
        row = index.row()
        file_path = self.result_table.item(row, 1).text()
        try:
            if sys.platform == "win32":
                os.startfile(file_path)
            elif sys.platform == "darwin":
                import subprocess

                subprocess.call(["open", file_path])
            else:
                import subprocess

                subprocess.call(["xdg-open", file_path])
            self.statusBar.showMessage(f"已開啟: {file_path}")
        except Exception as e:
            QMessageBox.warning(self, "無法開啟檔案", f"無法開啟檔案: {str(e)}")

    def show_context_menu(self, position):
        if not self.result_table.rowCount():
            return

        row = self.result_table.indexAt(position).row()
        if row < 0:
            return

        file_path = self.result_table.item(row, 1).text()
        file_name = self.result_table.item(row, 0).text()

        # 檢查是否有多個選中項
        selected_rows = sorted(
            {idx.row() for idx in self.result_table.selectedIndexes()}
        )
        is_multiple_selected = len(selected_rows) > 1

        menu = QMenu()

        # 檔案操作組
        open_action = QAction("🗂️ 開啟檔案", self)
        open_folder_action = QAction("📁 開啟所在資料夾", self)
        preview_action = QAction("👁️ 預覽內容", self)

        menu.addAction(open_action)
        menu.addAction(open_folder_action)
        menu.addAction(preview_action)
        menu.addSeparator()

        # 複製到剪貼簿組
        copy_menu = menu.addMenu("📋 複製到剪貼簿")

        if is_multiple_selected:
            copy_files_clipboard_action = QAction(
                f"📄 複製 {len(selected_rows)} 個檔案", self
            )
            copy_paths_action = QAction(f"📍 複製 {len(selected_rows)} 個路徑", self)
            copy_names_action = QAction(f"🏷️ 複製 {len(selected_rows)} 個檔名", self)
        else:
            copy_files_clipboard_action = QAction("📄 複製檔案", self)
            copy_paths_action = QAction("📍 複製路徑", self)
            copy_names_action = QAction("🏷️ 複製檔名", self)

        copy_menu.addAction(copy_files_clipboard_action)
        copy_menu.addAction(copy_paths_action)
        copy_menu.addAction(copy_names_action)

        # 批量操作
        if is_multiple_selected:
            menu.addSeparator()
            batch_menu = menu.addMenu(f"⚡ 批量操作 ({len(selected_rows)} 個檔案)")
            
            batch_copy_action = QAction("複製到資料夾...", self)
            batch_copy_action.triggered.connect(self.batch_copy_files)
            batch_menu.addAction(batch_copy_action)
            
            batch_move_action = QAction("移動到資料夾...", self)
            batch_move_action.triggered.connect(self.batch_move_files)
            batch_menu.addAction(batch_move_action)
            
            batch_delete_action = QAction("刪除...", self)
            batch_delete_action.triggered.connect(self.batch_delete_files)
            batch_menu.addAction(batch_delete_action)
            
            batch_rename_action = QAction("批量重新命名...", self)
            batch_rename_action.triggered.connect(self.batch_rename_files)
            batch_menu.addAction(batch_rename_action)

        # 檔案管理組
        if is_multiple_selected:
            menu.addSeparator()
            copy_to_folder_action = QAction(
                f"📂 複製 {len(selected_rows)} 個檔案到資料夾...", self
            )
            menu.addAction(copy_to_folder_action)
            copy_to_folder_action.triggered.connect(self.copy_selected_files)

        # 連接信號
        open_action.triggered.connect(
            lambda: self.open_file(self.result_table.indexAt(position))
        )
        open_folder_action.triggered.connect(lambda: self.open_folder(file_path))
        preview_action.triggered.connect(
            lambda: self.preview_file(file_path, file_name)
        )

        copy_files_clipboard_action.triggered.connect(
            self.copy_selected_files_to_clipboard
        )
        copy_paths_action.triggered.connect(self.copy_selected_paths)
        copy_names_action.triggered.connect(self.copy_selected_filenames)

        menu.exec_(self.result_table.viewport().mapToGlobal(position))

    def copy_selected_paths(self):
        rows = sorted({idx.row() for idx in self.result_table.selectedIndexes()})
        if not rows:
            QMessageBox.information(self, "提示", "請先選取至少一個檔案列。")
            return
        paths = []
        for r in rows:
            item = self.result_table.item(r, 1)
            if item:
                paths.append(item.text())
        if paths:
            QApplication.clipboard().setText("\n".join(paths))
            self.statusBar.showMessage(f"已複製 {len(paths)} 個路徑到剪貼簿")

    def copy_selected_filenames(self):
        rows = sorted({idx.row() for idx in self.result_table.selectedIndexes()})
        if not rows:
            QMessageBox.information(self, "提示", "請先選取至少一個檔案列。")
            return
        names = []
        for r in rows:
            item = self.result_table.item(r, 0)
            if item:
                names.append(item.text())
        if names:
            QApplication.clipboard().setText("\n".join(names))
            self.statusBar.showMessage(f"已複製 {len(names)} 個檔名到剪貼簿")

    def batch_copy_files(self):
        """批量複製檔案"""
        rows = sorted({idx.row() for idx in self.result_table.selectedIndexes()})
        if not rows:
            return
        
        target_folder = QFileDialog.getExistingDirectory(self, "選擇目標資料夾")
        if not target_folder:
            return
        
        success_count = 0
        for r in rows:
            path_item = self.result_table.item(r, 1)
            if path_item:
                src_path = path_item.text()
                if os.path.exists(src_path):
                    try:
                        import shutil
                        shutil.copy2(src_path, target_folder)
                        success_count += 1
                    except Exception as e:
                        print(f"複製 {src_path} 失敗: {e}")
        
        self.statusBar.showMessage(f"已複製 {success_count} 個檔案到 {target_folder}")

    def batch_move_files(self):
        """批量移動檔案"""
        rows = sorted({idx.row() for idx in self.result_table.selectedIndexes()})
        if not rows:
            return
        
        target_folder = QFileDialog.getExistingDirectory(self, "選擇目標資料夾")
        if not target_folder:
            return
        
        reply = QMessageBox.question(self, "確認移動", 
            f"確定要移動 {len(rows)} 個檔案？此操作無法復原。",
            QMessageBox.Yes | QMessageBox.No, QMessageBox.No)
        if reply != QMessageBox.Yes:
            return
        
        success_count = 0
        for r in rows:
            path_item = self.result_table.item(r, 1)
            if path_item:
                src_path = path_item.text()
                if os.path.exists(src_path):
                    try:
                        import shutil
                        shutil.move(src_path, target_folder)
                        success_count += 1
                    except Exception as e:
                        print(f"移動 {src_path} 失敗: {e}")
        
        self.statusBar.showMessage(f"已移動 {success_count} 個檔案到 {target_folder}")

    def batch_delete_files(self):
        """批量刪除檔案"""
        rows = sorted({idx.row() for idx in self.result_table.selectedIndexes()})
        if not rows:
            return
        
        reply = QMessageBox.warning(self, "確認刪除",
            f"確定要刪除 {len(rows)} 個檔案？此操作無法復原！",
            QMessageBox.Yes | QMessageBox.No, QMessageBox.No)
        if reply != QMessageBox.Yes:
            return
        
        success_count = 0
        for r in rows:
            path_item = self.result_table.item(r, 1)
            if path_item:
                file_path = path_item.text()
                if os.path.exists(file_path):
                    try:
                        os.remove(file_path)
                        success_count += 1
                    except Exception as e:
                        print(f"刪除 {file_path} 失敗: {e}")
        
        self.statusBar.showMessage(f"已刪除 {success_count} 個檔案")

    def batch_rename_files(self):
        """批量重新命名檔案"""
        rows = sorted({idx.row() for idx in self.result_table.selectedIndexes()})
        if not rows:
            return
        
        prefix = QInputDialog.getText(self, "批量重新命名", "輸入前綴 (可為空):")[0]
        suffix = QInputDialog.getText(self, "批量重新命名", "輸入後綴 (可為空):")[0]
        
        success_count = 0
        for idx, r in enumerate(rows, 1):
            path_item = self.result_table.item(r, 1)
            name_item = self.result_table.item(r, 0)
            if path_item and name_item:
                old_path = path_item.text()
                old_name = name_item.text()
                if os.path.exists(old_path):
                    name_part, ext = os.path.splitext(old_name)
                    new_name = f"{prefix}{name_part}{suffix}{ext}"
                    new_path = os.path.join(os.path.dirname(old_path), new_name)
                    try:
                        os.rename(old_path, new_path)
                        success_count += 1
                    except Exception as e:
                        print(f"重新命名 {old_path} 失敗: {e}")
        
        self.statusBar.showMessage(f"已重新命名 {success_count} 個檔案")

    def copy_selected_files_to_clipboard(self):
        """複製選中的檔案到剪貼簿"""
        rows = sorted({idx.row() for idx in self.result_table.selectedIndexes()})
        if not rows:
            QMessageBox.information(self, "提示", "請先選取至少一個檔案列。")
            return

        paths = []
        for r in rows:
            item = self.result_table.item(r, 1)
            if item and os.path.isfile(item.text()):
                paths.append(item.text())

        if not paths:
            QMessageBox.information(self, "提示", "無可複製的檔案。")
            return

        import platform

        if platform.system() == "Windows":
            try:
                # Windows 系統：嘗試使用 win32clipboard
                import win32clipboard
                import win32con
                import struct

                # 準備 DROPFILES 結構
                file_list = "\0".join(paths) + "\0\0"
                file_list_bytes = file_list.encode("utf-16le")

                # DROPFILES 結構頭部 (20 bytes)
                dropfiles_header = struct.pack("LLLLL", 20, 0, 0, 0, 1)

                # 組合完整數據
                clipboard_data = dropfiles_header + file_list_bytes

                win32clipboard.OpenClipboard()
                win32clipboard.EmptyClipboard()
                win32clipboard.SetClipboardData(win32con.CF_HDROP, clipboard_data)
                win32clipboard.CloseClipboard()

                self.statusBar.showMessage(f"已複製 {len(paths)} 個檔案到剪貼簿")

            except ImportError:
                # 如果沒有安裝 pywin32，使用文字模式
                QApplication.clipboard().setText("\n".join(paths))
                self.statusBar.showMessage(f"已複製 {len(paths)} 個檔案路徑到剪貼簿")
            except Exception as e:
                QMessageBox.warning(self, "複製失敗", f"複製檔案到剪貼簿失敗：{str(e)}")

        elif platform.system() == "Darwin":  # macOS
            try:
                # macOS：使用 pbcopy 或 PyQt 的剪貼簿
                # 在 macOS 上，檔案複製到剪貼簿較複雜，這裡提供路徑列表
                from PyQt5.QtCore import QMimeData, QUrl

                mime_data = QMimeData()

                # 設定檔案 URL 列表
                urls = [QUrl.fromLocalFile(path) for path in paths]
                mime_data.setUrls(urls)

                # 同時設定純文字路徑
                mime_data.setText("\n".join(paths))

                QApplication.clipboard().setMimeData(mime_data)
                self.statusBar.showMessage(f"已複製 {len(paths)} 個檔案到剪貼簿")

            except Exception as e:
                # 備用方案：只複製路徑文字
                QApplication.clipboard().setText("\n".join(paths))
                self.statusBar.showMessage(f"已複製 {len(paths)} 個檔案路徑到剪貼簿")

        else:  # Linux 和其他系統
            try:
                from PyQt5.QtCore import QMimeData, QUrl

                mime_data = QMimeData()
                urls = [QUrl.fromLocalFile(path) for path in paths]
                mime_data.setUrls(urls)
                mime_data.setText("\n".join(paths))

                QApplication.clipboard().setMimeData(mime_data)
                self.statusBar.showMessage(f"已複製 {len(paths)} 個檔案到剪貼簿")

            except Exception:
                # 備用方案
                QApplication.clipboard().setText("\n".join(paths))
                self.statusBar.showMessage(f"已複製 {len(paths)} 個檔案路徑到剪貼簿")

    def copy_selected_files(self):
        """複製選中的檔案到指定資料夾"""
        rows = sorted({idx.row() for idx in self.result_table.selectedIndexes()})
        if not rows:
            QMessageBox.information(self, "提示", "請先選取至少一個檔案列。")
            return
        paths = []
        for r in rows:
            item = self.result_table.item(r, 1)
            if item:
                paths.append(item.text())
        if not paths:
            QMessageBox.information(self, "提示", "無可複製的檔案路徑。")
            return

        dst_dir = QFileDialog.getExistingDirectory(self, "選擇複製目標資料夾")
        if not dst_dir:
            return

        apply_for_all = {"decision": None}  # 'overwrite', 'skip'
        success_count = 0
        fail_count = 0

        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
        total = len(paths)
        self.statusBar.showMessage("正在複製已選取檔案...")

        for i, src in enumerate(paths, start=1):
            try:
                if not os.path.isfile(src):
                    raise FileNotFoundError("來源非檔案或不存在")

                filename = os.path.basename(src)
                dst_path = os.path.join(dst_dir, filename)

                if os.path.exists(dst_path):
                    decision = apply_for_all["decision"]
                    if decision is None:
                        btn = QMessageBox.question(
                            self,
                            "檔名衝突",
                            f"目標已存在檔案：\n{dst_path}\n\n要覆蓋嗎？",
                            QMessageBox.Yes
                            | QMessageBox.No
                            | QMessageBox.YesAll
                            | QMessageBox.NoAll,
                            QMessageBox.No,
                        )
                        if btn == QMessageBox.Yes:
                            decision = "overwrite"
                        elif btn == QMessageBox.No:
                            decision = "skip"
                        elif btn == QMessageBox.YesAll:
                            decision = "overwrite"
                            apply_for_all["decision"] = "overwrite"
                        elif btn == QMessageBox.NoAll:
                            decision = "skip"
                            apply_for_all["decision"] = "skip"
                    if decision == "skip":
                        pass
                    else:
                        shutil.copy2(src, dst_path)
                        success_count += 1
                else:
                    shutil.copy2(src, dst_path)
                    success_count += 1
            except Exception as e:
                fail_count += 1
                self.statusBar.showMessage(f"複製失敗：{os.path.basename(src)} - {e}")

            self.progress_bar.setValue(int(i / total * 100))

        self.statusBar.showMessage(f"複製完成：成功 {success_count}，失敗 {fail_count}")
        QMessageBox.information(
            self,
            "複製完成",
            f"成功：{success_count}\n失敗：{fail_count}\n目標資料夾：\n{dst_dir}",
        )
        self.progress_bar.setVisible(False)

    def open_folder(self, file_path):
        folder_path = os.path.dirname(file_path)
        try:
            if sys.platform == "win32":
                os.startfile(folder_path)
            elif sys.platform == "darwin":
                import subprocess

                subprocess.call(["open", folder_path])
            else:
                import subprocess

                subprocess.call(["xdg-open", folder_path])
            self.statusBar.showMessage(f"已開啟資料夾: {folder_path}")
        except Exception as e:
            QMessageBox.warning(self, "無法開啟資料夾", f"無法開啟資料夾: {str(e)}")

    def _highlight_text_html(self, text, keywords):
        """將文字中的關鍵字高亮顯示為 HTML"""
        if not keywords:
            return text
        escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        for kw in keywords:
            if kw and len(kw) > 0:
                escaped_kw = kw.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                escaped = escaped.replace(
                    escaped_kw,
                    f'<span style="background-color: #e5f2ff; color: #007aff; font-weight: 600;">{escaped_kw}</span>'
                )
        return f'<pre style="font-family: Consolas, Menlo, Monaco, monospace; font-size: 12px; white-space: pre-wrap;">{escaped}</pre>'

    def preview_file(self, file_path, file_name):
        try:
            _, ext = os.path.splitext(file_name)
            ext = ext.lower()
            previewable_exts = {
                ".txt",
                ".py",
                ".java",
                ".c",
                ".cpp",
                ".html",
                ".xml",
                ".json",
                ".csv",
                ".md",
                ".log",
                ".ini",
                ".conf",
                ".js",
                ".css",
                ".yml",
                ".yaml",
            }
            rich_preview_exts = {".pdf", ".docx", ".xlsx", ".pptx"}
            keywords = []
            if hasattr(self, "searcher") and hasattr(self.searcher, "search_config"):
                config = self.searcher.search_config
                if config.get("filename"):
                    keywords.append(config["filename"])
                if config.get("content"):
                    keywords.append(config["content"])

            if ext in previewable_exts:
                try:
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read(102400)
                        if len(content) == 102400:
                            content += "\n\n[檔案過大，只顯示前100KB內容...]"
                except Exception:
                    content = "無法讀取檔案內容，可能是編碼問題或權限問題。"
                encoding_info = self._detect_encoding(file_path, ext)
                if keywords:
                    content = self._highlight_text_html(content, keywords)
                    self.content_preview.setHtml(content)
                else:
                    if encoding_info:
                        content = f"[偵測編碼: {encoding_info}]\n\n{content}"
                    self.content_preview.setText(content)
            elif ext in rich_preview_exts:
                content = ContentExtractor.extract_text(file_path, ext)
                if not content:
                    content = "此檔案沒有可提取的文字內容。"
                else:
                    page_info = {
                        ".pdf": "前 20 頁",
                        ".docx": "全部段落",
                        ".xlsx": "前 5 個工作表",
                        ".pptx": "前 20 張投影片",
                    }.get(ext, "")
                    content = f"[{ext.upper()} 檔案 — 提取{page_info}]\n\n{content}"
                if keywords:
                    content = self._highlight_text_html(content, keywords)
                    self.content_preview.setHtml(content)
                else:
                    self.content_preview.setText(content)
            else:
                content = f"無法預覽此類型檔案: {ext}"
                self.content_preview.setText(content)

            # 自動顯示預覽視窗
            if not self.preview_visible:
                self.preview_widget.show()
                self.preview_visible = True

            self.statusBar.showMessage(f"已預覽: {file_path}")
        except Exception as e:
            self.content_preview.setText(f"預覽檔案時發生錯誤: {str(e)}")

    def _detect_encoding(self, file_path, ext):
        """簡單偵測文字檔案的編碼"""
        text_exts = {
            ".txt", ".py", ".java", ".c", ".cpp", ".html", ".xml",
            ".json", ".csv", ".md", ".log", ".ini", ".conf",
            ".js", ".css", ".yml", ".yaml",
        }
        if ext not in text_exts:
            return None
        try:
            with open(file_path, "rb") as f:
                raw = f.read(8192)
            if not raw:
                return None
            for enc, bom in [("UTF-8 BOM", b"\xef\xbb\xbf"), ("UTF-16 BE", b"\xfe\xff"), ("UTF-16 LE", b"\xff\xfe")]:
                if raw.startswith(bom):
                    return enc
            try:
                raw.decode("utf-8")
                return "UTF-8"
            except UnicodeDecodeError:
                pass
            try:
                raw.decode("big5")
                return "Big5"
            except UnicodeDecodeError:
                pass
            try:
                raw.decode("gbk")
                return "GBK"
            except UnicodeDecodeError:
                pass
            return "Latin-1 (fallback)"
        except Exception:
            return None

    def copy_path(self, file_path):
        QApplication.clipboard().setText(file_path)
        self.statusBar.showMessage(f"已複製路徑: {file_path}")

    # ---------------- 鍵盤快捷鍵 ----------------
    def keyPressEvent(self, event):
        """處理鍵盤按鍵事件"""
        # F3 鍵切換預覽視窗
        if event.key() == Qt.Key_F3:
            self.toggle_preview_panel()
            event.accept()
        else:
            super().keyPressEvent(event)

    def toggle_preview_panel(self):
        """切換檔案預覽視窗的顯示/隱藏"""
        if self.preview_visible:
            self.preview_widget.hide()
            self.preview_visible = False
            self.statusBar.showMessage("檔案預覽視窗已隱藏 (按 F3 重新顯示)")
        else:
            self.preview_widget.show()
            self.preview_visible = True
            self.statusBar.showMessage("檔案預覽視窗已顯示 (按 F3 隱藏)")

    # ---------------- 拖放支援 ----------------
    def dragEnterEvent(self, event):
        """接受拖放事件"""
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
            self._drag_overlay_active = True
            self._apply_drag_style()
        else:
            event.ignore()

    def dragMoveEvent(self, event):
        """拖放移動時的視覺回饋"""
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
        else:
            event.ignore()

    def dragLeaveEvent(self, event):
        """拖放離開時恢復背景"""
        self._drag_overlay_active = False
        self._apply_theme()

    def _apply_drag_style(self):
        """套用拖放視覺回饋（保留主題樣式）"""
        base_style = DARK_STYLE if self.dark_mode else LIGHT_STYLE
        overlay = "QMainWindow { background-color: #e8f0fe; }"
        self.setStyleSheet(base_style + "\n" + overlay)

    def dropEvent(self, event):
        """處理拖放事件"""
        self._drag_overlay_active = False
        self._apply_theme()
        paths = []
        for url in event.mimeData().urls():
            local_path = url.toLocalFile()
            if local_path and os.path.exists(local_path):
                paths.append(local_path)

        if not paths:
            return

        self._handle_dropped_files(paths)
        event.acceptProposedAction()

    def _handle_dropped_files(self, paths):
        """根據當前頁籤處理拖放的檔案/資料夾"""
        current_tab = self.tab_widget.currentIndex()

        if current_tab == 0:  # 搜尋頁籤
            self._handle_drop_on_search(paths)
        elif current_tab == 1:  # 索引頁籤
            self._handle_drop_on_index(paths)
        else:  # 其他頁籤
            self._handle_drop_default(paths)

    def _handle_drop_on_search(self, paths):
        """拖放到搜尋頁籤: 單一檔案直接預覽，多個檔案/資料夾設定路徑過濾"""
        if len(paths) == 1 and os.path.isfile(paths[0]):
            file_path = paths[0]
            file_name = os.path.basename(file_path)
            self.preview_file(file_path, file_name)
            self.statusBar.showMessage(f"已預覽拖放檔案: {file_name}")
        else:
            dir_paths = [p for p in paths if os.path.isdir(p)]
            file_paths = [p for p in paths if os.path.isfile(p)]
            filter_parts = []
            if dir_paths:
                filter_parts.extend(dir_paths)
            if file_paths:
                filter_parts.extend(file_paths)
            current_filter = self.path_filter.text().strip()
            if current_filter:
                new_filter = current_filter + " " + " ".join(f'"{p}"' for p in filter_parts)
            else:
                new_filter = " ".join(f'"{p}"' for p in filter_parts)
            self.path_filter.setText(new_filter)
            self._perform_search()
            self.statusBar.showMessage(f"已加入 {len(paths)} 個拖放項目到搜尋路徑過濾")

    def _handle_drop_on_index(self, paths):
        """拖放到索引頁籤: 將資料夾加入索引目錄列表"""
        existing_dirs = set()
        if hasattr(self, "index_directories"):
            existing_text = self.index_directories.toPlainText().strip()
            if existing_text:
                existing_dirs = set(l.strip() for l in existing_text.splitlines() if l.strip())

        new_dirs = []
        for p in paths:
            if os.path.isdir(p) and p not in existing_dirs:
                new_dirs.append(p)

        if new_dirs:
            current_text = self.index_directories.toPlainText().strip()
            added = "\n".join(new_dirs)
            if current_text:
                self.index_directories.setPlainText(current_text + "\n" + added)
            else:
                self.index_directories.setPlainText(added)
            self.statusBar.showMessage(f"已加入 {len(new_dirs)} 個資料夾到索引目錄")
        else:
            self.statusBar.showMessage("拖放的資料夾已在索引目錄中")

    def _handle_drop_default(self, paths):
        """預設行為: 單一檔案預覽，多個檔案顯示路徑"""
        if len(paths) == 1 and os.path.isfile(paths[0]):
            self.preview_file(paths[0], os.path.basename(paths[0]))
            self.tab_widget.setCurrentIndex(0)
            self.statusBar.showMessage(f"已預覽拖放檔案: {os.path.basename(paths[0])}")
        else:
            self.tab_widget.setCurrentIndex(0)
            filter_str = " ".join(f'"{p}"' for p in paths)
            self.path_filter.setText(filter_str)
            self._perform_search()
            self.statusBar.showMessage(f"已搜尋 {len(paths)} 個拖放項目")


def check_single_instance():
    """使用套接字檢查單實例（跨平台方案）"""
    import socket
    import platform

    # 定義一個唯一的端口號
    PORT = 62847  # 選擇一個不太可能被其他程式使用的端口

    try:
        # 嘗試綁定到本地端口
        global single_instance_socket
        single_instance_socket = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
        if platform.system() != "Windows":
            # 在 Unix-like 系統上設置 SO_REUSEADDR
            single_instance_socket.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
        single_instance_socket.bind(("127.0.0.1", PORT))

        # 成功綁定，表示沒有其他實例在運行
        return True

    except socket.error:
        # 端口已被佔用，表示已有實例在運行
        print(f"{APP_NAME} 已經在運行中")

        # 在 macOS 上嘗試激活現有視窗
        if platform.system() == "Darwin":
            try:
                import subprocess

                app_name = "SearchingPro" if getattr(sys, "frozen", False) else "Python"
                subprocess.run(
                    [
                        "osascript",
                        "-e",
                        f'tell application "System Events" to set frontmost of (first process whose name contains "{app_name}") to true',
                    ],
                    capture_output=True,
                )
            except Exception:
                pass

        return False


def main():
    # 檢查單實例
    if not check_single_instance():
        # 已有實例在運行，直接退出
        sys.exit(0)

    app = QApplication(sys.argv)
    window = AdvancedFileSearcher()
    window.show()
    sys.exit(app.exec_())


if __name__ == "__main__":
    main()
